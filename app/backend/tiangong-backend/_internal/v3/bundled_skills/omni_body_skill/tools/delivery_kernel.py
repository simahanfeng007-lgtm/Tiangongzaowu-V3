"""
Tiangong Omni Body v3.2 Delivery Kernel
======================================

This module keeps the system a tool, not an agent. It provides deterministic
quality gates, template operations, preview extraction, packaging, and basic
repair-plan generation. The model must still choose actions and iterate.
"""
from __future__ import annotations

import ast
import csv
import json
import os
import re
import shutil
import subprocess
import sys
import textwrap
import time
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Tuple

try:
    from .delivery_v33 import V33_DELIVERY_ACTIONS, handle_v33_action  # type: ignore
except ImportError:
    V33_DELIVERY_ACTIONS = {}
    handle_v33_action = None  # type: ignore

try:
    from .skill_router import SKILL_ROUTER_ACTIONS, handle_skill_router_action  # type: ignore
except ImportError:
    SKILL_ROUTER_ACTIONS = {}
    handle_skill_router_action = None  # type: ignore

try:
    from .pro_apps_v34 import PRO_APP_ACTIONS, handle_pro_app_action  # type: ignore
except ImportError:
    PRO_APP_ACTIONS = {}
    handle_pro_app_action = None  # type: ignore

try:
    from ..model_adapters.core import MODEL_ADAPTER_ACTIONS, handle_model_adapter_action  # type: ignore
except ImportError:
    # The skill can also be mounted with its root directly on sys.path, where
    # ``tools`` and ``model_adapters`` are sibling top-level packages.
    try:
        from model_adapters.core import MODEL_ADAPTER_ACTIONS, handle_model_adapter_action  # type: ignore
    except ImportError:
        MODEL_ADAPTER_ACTIONS = {}
        handle_model_adapter_action = None  # type: ignore

try:
    from .novel_system import NOVEL_SYSTEM_ACTIONS, handle_novel_system_action  # type: ignore
except ImportError:
    NOVEL_SYSTEM_ACTIONS = {}
    handle_novel_system_action = None  # type: ignore

DELIVERY_ACTIONS: Dict[str, Dict[str, Any]] = {
    "delivery.kernel.info": {"risk": "A0", "implemented": True, "summary": "Inspect v3.2 delivery kernel standards, rubrics, and available quality gates."},
    "template.list": {"risk": "A0", "implemented": True, "summary": "List delivery templates and rubrics shipped with the package."},
    "template.apply": {"risk": "A2", "implemented": True, "summary": "Apply a template skeleton and create a structured draft markdown/json file."},
    "preview.generate": {"risk": "A0", "implemented": True, "summary": "Generate lightweight preview/summary evidence for docx/pptx/xlsx/image/video/text deliverables."},
    "rubric.evaluate": {"risk": "A0", "implemented": True, "summary": "Evaluate supplied content or target file against a named delivery rubric."},

    "qc.docx.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check business-grade Word/document deliverables for structure, clarity, evidence, actionability, and openability."},
    "qc.ppt.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check executive-grade PPT deliverables for storyline, slide density, title quality, structure, evidence, and CTA."},
    "qc.sheet.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check spreadsheets for headers, empty cells, duplicate rows, numeric consistency, formulas, and delivery readiness."},
    "qc.code.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check code deliverables for syntax, tests, README, structure, security smells, and maintainability evidence."},
    "qc.research.evidence_check": {"risk": "A0", "implemented": True, "summary": "Check research/literature deliverables for search strategy, inclusion/exclusion, citations, evidence table, limitations, and uncertainty."},
    "qc.video.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check short-video deliverables for playability, duration, aspect ratio, audio/subtitle evidence, hook/CTA metadata, and package readiness."},
    "qc.image.delivery_check": {"risk": "A0", "implemented": True, "summary": "Check image/poster deliverables for dimensions, readability, text overflow risk, contrast proxy, and export readiness."},
    "qc.writing.ai_tone_check": {"risk": "A0", "implemented": True, "summary": "Heuristically check writing for generic AI tone, repetition, vague claims, and weak specificity."},

    "writing.outline.create": {"risk": "A2", "implemented": True, "summary": "Create a structured outline markdown for proposal, deck, research, novel, or video script workflows."},
    "research.evidence_table.create": {"risk": "A2", "implemented": True, "summary": "Create a structured research evidence table CSV/Markdown from supplied sources."},
    "repair.plan": {"risk": "A2", "implemented": True, "summary": "Write a repair plan file from QC issues; does not autonomously modify deliverables."},
    "deliverable.package": {"risk": "A2", "implemented": True, "summary": "Package final deliverables, QC reports, source notes, and manifests into a zip archive."},
}
DELIVERY_ACTIONS.update(V33_DELIVERY_ACTIONS)
DELIVERY_ACTIONS.update(SKILL_ROUTER_ACTIONS)
DELIVERY_ACTIONS.update(PRO_APP_ACTIONS)
DELIVERY_ACTIONS.update(MODEL_ADAPTER_ACTIONS)
DELIVERY_ACTIONS.update(NOVEL_SYSTEM_ACTIONS)

RUBRIC_WEIGHTS = {
    "business_proposal": {
        "customer_focus": 15,
        "executive_summary": 15,
        "problem_solution_fit": 15,
        "evidence_and_proof": 15,
        "implementation_plan": 12,
        "risk_and_assumptions": 10,
        "commercial_actionability": 10,
        "clarity": 8,
    },
    "executive_ppt": {
        "single_big_idea": 16,
        "storyline": 16,
        "slide_titles": 14,
        "evidence": 14,
        "visual_density": 12,
        "audience_transformation": 10,
        "cta": 10,
        "consistency": 8,
    },
    "code_project": {
        "correctness": 18,
        "tests": 16,
        "readability": 15,
        "maintainability": 15,
        "security": 12,
        "documentation": 10,
        "packaging": 8,
        "rollback": 6,
    },
    "research_review": {
        "question": 12,
        "search_strategy": 14,
        "screening": 12,
        "evidence_table": 14,
        "citation_traceability": 14,
        "synthesis": 14,
        "limitations": 10,
        "uncertainty": 10,
    },
    "short_video": {
        "hook": 16,
        "narrative": 14,
        "vertical_mobile_fit": 14,
        "caption_sound": 14,
        "pace": 12,
        "brand_message": 10,
        "cta": 10,
        "technical_export": 10,
    },
}

GENERIC_AI_PHRASES = [
    "在当今快速发展的", "赋能", "闭环", "抓手", "生态", "降本增效", "全方位", "多维度",
    "显著提升", "深度融合", "未来可期", "以用户为中心", "打造", "助力", "全面提升",
]


def handle_delivery_action(runtime: Any, op_id: str, action: str, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    if action in globals().get("NOVEL_SYSTEM_ACTIONS", {}):
        if globals().get("handle_novel_system_action") is None:
            return {"success": False, "op_id": op_id, "action": action, "message": "novel system unavailable"}
        return globals()["handle_novel_system_action"](runtime, op_id, action, target, args)
    if action in globals().get("MODEL_ADAPTER_ACTIONS", {}):
        if globals().get("handle_model_adapter_action") is None:
            return {"success": False, "op_id": op_id, "action": action, "message": "v3.5 model adapter layer unavailable"}
        return globals()["handle_model_adapter_action"](runtime, op_id, action, target, args)
    if action in globals().get("PRO_APP_ACTIONS", {}):
        if globals().get("handle_pro_app_action") is None:
            return {"success": False, "op_id": op_id, "action": action, "message": "v3.4 professional app layer unavailable"}
        return globals()["handle_pro_app_action"](runtime, op_id, action, target, args)
    if action in globals().get("SKILL_ROUTER_ACTIONS", {}):
        if globals().get("handle_skill_router_action") is None:
            return {"success": False, "op_id": op_id, "action": action, "message": "v3.3.1 skill router unavailable"}
        return globals()["handle_skill_router_action"](runtime, op_id, action, target, args)
    if action in globals().get("V33_DELIVERY_ACTIONS", {}):
        if globals().get("handle_v33_action") is None:
            return {"success": False, "op_id": op_id, "action": action, "message": "v3.3 delivery expansion unavailable"}
        return globals()["handle_v33_action"](runtime, op_id, action, target, args)
    if action == "delivery.kernel.info":
        return _delivery_kernel_info(runtime, target, args)
    if action == "template.list":
        return _template_list(runtime, target, args)
    if action == "template.apply":
        return _template_apply(runtime, target, args)
    if action == "preview.generate":
        return _preview_generate(runtime, target, args)
    if action == "rubric.evaluate":
        return _rubric_evaluate(runtime, target, args)
    if action == "qc.docx.delivery_check":
        return _qc_docx(runtime, target, args)
    if action == "qc.ppt.delivery_check":
        return _qc_ppt(runtime, target, args)
    if action == "qc.sheet.delivery_check":
        return _qc_sheet(runtime, target, args)
    if action == "qc.code.delivery_check":
        return _qc_code(runtime, target, args)
    if action == "qc.research.evidence_check":
        return _qc_research(runtime, target, args)
    if action == "qc.video.delivery_check":
        return _qc_video(runtime, target, args)
    if action == "qc.image.delivery_check":
        return _qc_image(runtime, target, args)
    if action == "qc.writing.ai_tone_check":
        return _qc_writing(runtime, target, args)
    if action == "writing.outline.create":
        return _writing_outline_create(runtime, target, args)
    if action == "research.evidence_table.create":
        return _research_evidence_table_create(runtime, target, args)
    if action == "repair.plan":
        return _repair_plan(runtime, target, args)
    if action == "deliverable.package":
        return _deliverable_package(runtime, target, args)
    return {"success": False, "op_id": op_id, "action": action, "message": f"Delivery action not implemented: {action}"}


def _resolve(runtime: Any, target: str | None, must_exist: bool = False) -> Path:
    return runtime._resolve(target, must_exist=must_exist)


def _rel(runtime: Any, path: Path) -> str:
    return runtime._rel(path)


def _write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _read_text_any(path: Path, max_chars: int = 300_000) -> str:
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt", ".json", ".csv", ".py", ".js", ".ts", ".html", ".xml", ".opml"}:
        return path.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    if suffix == ".docx":
        return _extract_docx_text(path)[:max_chars]
    if suffix == ".pptx":
        return _extract_pptx_text(path)[:max_chars]
    if suffix == ".xlsx":
        return _extract_xlsx_text(path)[:max_chars]
    if suffix == ".pdf":
        try:
            import pypdf  # type: ignore
            reader = pypdf.PdfReader(str(path))
            return "\n".join((page.extract_text() or "") for page in reader.pages)[:max_chars]
        except Exception:
            return ""
    return ""


def _zip_xml_text(path: Path, patterns: Tuple[str, ...]) -> str:
    out: List[str] = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if any(name.startswith(p) for p in patterns) and name.endswith(".xml"):
                    raw = zf.read(name).decode("utf-8", errors="ignore")
                    text = re.sub(r"<[^>]+>", " ", raw)
                    text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        out.append(text)
    except Exception:
        pass
    return "\n".join(out)


def _extract_docx_text(path: Path) -> str:
    try:
        import docx  # type: ignore
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except Exception:
        return _zip_xml_text(path, ("word/",))


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        lines: List[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"[slide {idx}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    lines.append(str(shape.text).strip())
        return "\n".join(lines)
    except Exception:
        return _zip_xml_text(path, ("ppt/slides/",))


def _extract_xlsx_text(path: Path) -> str:
    rows: List[str] = []
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=False)
        try:
            for ws in wb.worksheets:
                rows.append(f"[sheet {ws.title}]")
                for r in ws.iter_rows(max_row=50, values_only=True):
                    rows.append(" | ".join("" if c is None else str(c) for c in r))
        finally:
            wb.close()
        return "\n".join(rows)
    except Exception:
        return _zip_xml_text(path, ("xl/worksheets/",))


def _sentence_stats(text: str) -> Dict[str, Any]:
    sentences = [s.strip() for s in re.split(r"[。！？.!?]\s*", text) if s.strip()]
    lengths = [len(s) for s in sentences]
    return {
        "sentences": len(sentences),
        "avg_sentence_chars": round(sum(lengths) / max(1, len(lengths)), 1),
        "long_sentence_count": sum(1 for n in lengths if n > 90),
    }


def _score_from_issues(max_score: int, issues: List[Dict[str, Any]], warnings: List[Dict[str, Any]] | None = None) -> int:
    score = max_score
    has_critical = False
    for issue in issues:
        sev = issue.get("severity", "medium")
        has_critical = has_critical or sev == "critical"
        score -= {"critical": 20, "high": 12, "medium": 7, "low": 3}.get(sev, 5)
    for warning in warnings or []:
        score -= 2 if warning.get("severity", "low") == "low" else 4
    # A corrupt/unreadable artifact or a critical correctness defect must never
    # cross the delivery threshold merely because the rubric starts at 100.
    if has_critical:
        score = min(score, 59)
    return max(0, min(max_score, score))


def _grade(score: int) -> str:
    if score >= 90:
        return "world_class_ready"
    if score >= 80:
        return "delivery_ready"
    if score >= 70:
        return "acceptable_with_minor_repair"
    if score >= 60:
        return "needs_repair"
    return "not_ready"


def _issue(code: str, message: str, severity: str = "medium", repair: str = "") -> Dict[str, Any]:
    return {"code": code, "severity": severity, "message": message, "repair": repair or message}


def _delivery_kernel_info(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    root = Path(__file__).resolve().parents[1]
    return {
        "success": True,
        "result": {
            "schema": "tiangong.v3.delivery_kernel.v1",
            "version": "3.3.0",
            "principle": "tool-only: deterministic actions, evidence, quality gates, repair plans; no autonomous planning.",
            "rubrics": sorted(RUBRIC_WEIGHTS.keys()),
            "quality_gates": sorted(k for k in DELIVERY_ACTIONS if k.startswith("qc.")),
            "root": str(root),
        },
        "evidence": {"path": "delivery_kernel", "exists": True, "bytes": 0},
    }


def _template_list(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    root = Path(__file__).resolve().parents[1]
    manifest = root / "templates" / "manifest.json"
    data = json.loads(manifest.read_text(encoding="utf-8")) if manifest.exists() else {"templates": []}
    return {"success": True, "result": data, "evidence": {"path": _rel(runtime, manifest) if manifest.exists() else "templates", "exists": manifest.exists(), "bytes": manifest.stat().st_size if manifest.exists() else 0}}


def _template_apply(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    template_id = str(args.get("template_id") or args.get("id") or "business_proposal").strip()
    output = _resolve(runtime, target or args.get("output") or f"{template_id}_draft.md")
    variables = args.get("variables") if isinstance(args.get("variables"), dict) else {}
    skeleton = _template_skeleton(template_id, variables)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(skeleton, encoding="utf-8")
    result: Dict[str, Any] = {
        "success": True,
        "output": {
            "path": _rel(runtime, output),
            "exists": output.exists(),
            "bytes": output.stat().st_size,
            "template_id": template_id,
        },
    }
    try:
        from .ppt_design import template_asset_root  # type: ignore
        design_root = template_asset_root()
    except Exception:
        design_root = Path(__file__).resolve().parents[1]
    design_source = design_root / "templates" / f"{template_id}.design.json"
    if design_source.is_file():
        design = json.loads(design_source.read_text(encoding="utf-8", errors="strict"))
        if not isinstance(design, dict) or design.get("schema") != "tiangong.v3.ppt_design.v1":
            raise ValueError(f"invalid machine-readable design contract for template {template_id}")
        design_output = _resolve(runtime, args.get("design_output") or output.with_suffix(".design.json"))
        design_output.parent.mkdir(parents=True, exist_ok=True)
        _write_json(design_output, design)
        result["design_spec"] = {
            "path": _rel(runtime, design_output),
            "exists": design_output.exists(),
            "bytes": design_output.stat().st_size,
            "schema": design.get("schema"),
        }
        result["next_action_args"] = {
            "template_id": template_id,
            "design_spec": _rel(runtime, design_output),
        }
    return result


def _template_skeleton(template_id: str, v: Dict[str, Any]) -> str:
    title = v.get("title") or {
        "business_proposal": "商业方案初稿",
        "executive_ppt": "商业汇报故事线",
        "code_project": "代码工程交付说明",
        "research_review": "资料/论文综述初稿",
        "short_video": "短视频脚本与交付说明",
    }.get(template_id, f"{template_id} 模板")
    audience = v.get("audience", "待明确受众")
    if template_id == "business_proposal":
        sections = ["执行摘要", "受众与决策目标", "现状问题", "解决方案", "实施路径", "收益与证据", "风险与假设", "报价/资源", "行动建议"]
    elif template_id == "executive_ppt":
        sections = ["Big Idea", "受众现状", "核心结论", "三条支撑证据", "反对意见与回应", "实施路径", "决策请求"]
    elif template_id == "code_project":
        sections = ["需求边界", "架构设计", "运行方式", "测试证据", "安全与回滚", "交付清单"]
    elif template_id == "research_review":
        sections = ["研究问题", "搜索策略", "纳入/排除标准", "证据表", "综合结论", "局限性", "不确定性与下一步"]
    elif template_id == "short_video":
        sections = ["目标受众", "前3秒钩子", "脚本", "镜头节奏", "字幕/配乐", "封面", "CTA", "导出规格"]
    else:
        sections = ["目标", "输入", "流程", "质检", "交付"]
    body = [f"# {title}", "", f"- 受众：{audience}", f"- 交付目标：{v.get('objective', '待明确')}", f"- 版本：v0.1", ""]
    for sec in sections:
        body.append(f"## {sec}")
        body.append(v.get(sec, "待补充。"))
        body.append("")
    return "\n".join(body)


def _preview_generate(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    text = _read_text_any(path, max_chars=int(args.get("max_chars", 12000)))
    preview = {
        "path": _rel(runtime, path),
        "suffix": path.suffix.lower(),
        "bytes": path.stat().st_size,
        "text_chars": len(text),
        "text_preview": text[:1500],
        "line_count": text.count("\n") + 1 if text else 0,
    }
    if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}:
        try:
            from PIL import Image  # type: ignore
            with Image.open(path) as im:
                preview.update({"width": im.width, "height": im.height, "mode": im.mode, "format": im.format})
        except Exception as exc:
            preview["image_error"] = str(exc)
    return {"success": True, "result": preview, "evidence": {"path": _rel(runtime, path), "exists": True, "bytes": path.stat().st_size}}


def _rubric_evaluate(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    rubric = str(args.get("rubric") or args.get("rubric_id") or "business_proposal")
    content = str(args.get("content") or "")
    if target:
        p = _resolve(runtime, target, must_exist=True)
        content += "\n" + _read_text_any(p)
    weights = RUBRIC_WEIGHTS.get(rubric, RUBRIC_WEIGHTS["business_proposal"])
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    lower = content.lower()
    for key, weight in weights.items():
        cn_key = key.replace("_", " ")
        # Minimal deterministic proxy: require evidence markers or equivalent section words.
        has_signal = key in lower or cn_key in lower or _cn_signal(key, content)
        if not has_signal and weight >= 14:
            issues.append(_issue(f"missing_{key}", f"缺少高权重评分项：{key}", "high", f"补充 {key} 对应内容并给出证据。"))
        elif not has_signal:
            warnings.append(_issue(f"weak_{key}", f"评分项不明显：{key}", "low", f"强化 {key}。"))
    ai = _ai_tone_issues(content)
    warnings.extend(ai[:5])
    score = _score_from_issues(100, issues, warnings)
    return {"success": True, "result": {"rubric": rubric, "score": score, "grade": _grade(score), "issues": issues, "warnings": warnings, "weights": weights}, "evidence": {"path": target or "content", "exists": bool(target), "bytes": len(content.encode('utf-8'))}}


def _cn_signal(key: str, text: str) -> bool:
    signals = {
        "customer_focus": ["客户", "受众", "决策", "痛点"],
        "executive_summary": ["执行摘要", "核心结论", "摘要"],
        "problem_solution_fit": ["问题", "解决方案", "匹配"],
        "evidence_and_proof": ["证据", "案例", "数据", "来源"],
        "implementation_plan": ["实施", "里程碑", "计划"],
        "risk_and_assumptions": ["风险", "假设", "边界"],
        "commercial_actionability": ["行动", "报价", "预算", "ROI", "收益"],
        "single_big_idea": ["Big Idea", "大观点", "核心主张"],
        "storyline": ["故事线", "SCQA", "金字塔", "逻辑"],
        "slide_titles": ["标题", "结论句"],
        "visual_density": ["留白", "视觉", "版式"],
        "audience_transformation": ["受众转变", "当前", "未来"],
        "correctness": ["正确", "运行", "验证"],
        "tests": ["测试", "pytest", "unittest"],
        "readability": ["可读", "命名", "注释"],
        "maintainability": ["维护", "模块", "架构"],
        "security": ["安全", "注入", "权限"],
        "documentation": ["README", "文档"],
        "search_strategy": ["搜索策略", "关键词", "数据库"],
        "screening": ["纳入", "排除", "筛选"],
        "evidence_table": ["证据表", "研究", "样本"],
        "citation_traceability": ["引用", "来源", "doi", "url"],
        "synthesis": ["综合", "共识", "分歧"],
        "limitations": ["局限", "限制"],
        "uncertainty": ["不确定", "置信", "可能"],
        "hook": ["钩子", "前3秒", "开头"],
        "narrative": ["脚本", "镜头", "叙事"],
        "vertical_mobile_fit": ["9:16", "竖屏", "移动端"],
        "caption_sound": ["字幕", "配乐", "声音"],
        "pace": ["节奏", "剪辑点"],
        "brand_message": ["品牌", "卖点"],
        "cta": ["CTA", "行动", "转化"],
        "technical_export": ["导出", "mp4", "分辨率"],
    }
    return any(s.lower() in text.lower() for s in signals.get(key, []))


def _qc_docx(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    text = _read_text_any(path)
    document_type = str(args.get("document_type") or args.get("mode") or "").strip().lower()
    if document_type in {"long_document", "managed_long_document", "longform"}:
        return _qc_managed_long_document(runtime, path, text, args)
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    if path.suffix.lower() not in {".docx", ".md", ".txt", ".pdf"}:
        issues.append(_issue("wrong_format", "交付文件不是常见文档格式。", "high", "导出为 docx/pdf/md。"))
    required = ["执行摘要", "问题", "解决方案", "实施", "风险", "行动"]
    for word in required:
        if word not in text:
            issues.append(_issue(f"missing_{word}", f"缺少关键章节或内容：{word}", "medium", f"补充“{word}”相关章节。"))
    stats = _sentence_stats(text)
    if len(text) < 1200:
        issues.append(_issue("too_short", "方案正文过短，难以达到可交付级。", "high", "补充背景、分析、路径、风险、收益与证据。"))
    if stats["long_sentence_count"] > 6:
        warnings.append(_issue("dense_sentences", "长句过多，阅读负担偏高。", "low", "拆分长句，优先主谓宾短句。"))
    generic = _ai_tone_issues(text)
    warnings.extend(generic[:8])
    score = _score_from_issues(100, issues, warnings)
    report = {"type": "docx_business_delivery", "score": score, "grade": _grade(score), "stats": stats, "issues": issues, "warnings": warnings, "acceptance": score >= 80}
    return {"success": True, "result": report, "evidence": {"path": _rel(runtime, path), "exists": path.exists(), "bytes": path.stat().st_size, "score": score, "grade": report["grade"]}}


def _qc_managed_long_document(runtime: Any, path: Path, text: str, args: Dict[str, Any]) -> Dict[str, Any]:
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    manifest_path: Path | None = None
    manifest: Dict[str, Any] = {}
    raw_manifest = args.get("project_manifest")
    if not isinstance(raw_manifest, str) or not raw_manifest.strip():
        issues.append(_issue("project_manifest_required", "受管超长文档 QC 必须提供项目 manifest。", "critical", "传入工作区内 project_manifest.json。"))
    else:
        try:
            manifest_path = _resolve(runtime, raw_manifest, must_exist=True)
            loaded = json.loads(manifest_path.read_text(encoding="utf-8", errors="strict"))
            if not isinstance(loaded, dict):
                raise ValueError("manifest root must be an object")
            manifest = loaded
        except (OSError, UnicodeError, ValueError, json.JSONDecodeError) as exc:
            issues.append(_issue("project_manifest_invalid", f"项目 manifest 不可解析：{exc}", "critical", "修复 manifest 后重新质检。"))

    target_words = manifest.get("target_words")
    if not isinstance(target_words, int) or isinstance(target_words, bool) or target_words < 5_000:
        issues.append(_issue("target_words_invalid", "manifest.target_words 必须是至少 5000 的整数。", "critical", "写入真实目标字数。"))
        target_words = 5_000
    chapter_files = manifest.get("chapter_files")
    if not isinstance(chapter_files, list) or not chapter_files or not all(isinstance(item, str) and item.strip() for item in chapter_files):
        issues.append(_issue("chapter_files_invalid", "manifest.chapter_files 必须是非空相对路径数组。", "critical", "列出全部章节源文件。"))
        chapter_files = []

    missing: List[str] = []
    unsafe: List[str] = []
    if manifest_path is not None:
        project_root = manifest_path.parent.resolve()
        seen: set[str] = set()
        for item in chapter_files:
            folded = item.replace("\\", "/").casefold()
            if folded in seen:
                issues.append(_issue("duplicate_chapter_path", f"章节路径重复：{item}", "critical", "修复章节清单并保持唯一顺序。"))
                continue
            seen.add(folded)
            candidate = (project_root / item).resolve()
            try:
                candidate.relative_to(project_root)
            except ValueError:
                unsafe.append(item)
                continue
            if candidate.is_symlink() or not candidate.is_file():
                missing.append(item)
        if unsafe:
            issues.append(_issue("unsafe_chapter_paths", f"章节路径逃逸项目目录：{unsafe[:5]}", "critical", "只使用项目内相对路径。"))
        if missing:
            issues.append(_issue("missing_chapter_files", f"缺少 {len(missing)} 个章节源文件。", "critical", "先补齐章节再汇编。"))

    compact_chars = len(re.sub(r"\s+", "", text))
    minimum_chars = max(5_000, int(target_words * 0.75))
    if compact_chars < minimum_chars:
        issues.append(_issue("document_incomplete", f"正文有效字符 {compact_chars}，低于目标完成门 {minimum_chars}。", "critical", "补齐全部计划章节并重新生成 DOCX。"))
    heading_count = len(re.findall(r"(?m)^(?:第\s*[0-9一二三四五六七八九十百千万]+\s*[章节篇部]|[^\n]{1,80}\n[-=]{3,})", text))
    if len(chapter_files) >= 2 and heading_count < 2:
        warnings.append(_issue("weak_section_structure", "正文中可识别的章节结构偏少。", "medium", "保留清晰章节标题后重新生成。"))
    placeholders = len(re.findall(r"待补充|占位|TODO|TBD|lorem ipsum", text, flags=re.I))
    if placeholders:
        issues.append(_issue("placeholder_content", f"正文仍含 {placeholders} 处占位内容。", "high", "删除占位符并写入最终内容。"))
    paragraphs = [re.sub(r"\s+", "", item) for item in re.split(r"\n+", text) if len(re.sub(r"\s+", "", item)) >= 40]
    duplicate_ratio = (len(paragraphs) - len(set(paragraphs))) / max(1, len(paragraphs))
    if duplicate_ratio > 0.08:
        issues.append(_issue("duplicate_paragraphs", f"长段落重复率 {duplicate_ratio:.1%}。", "high", "删除重复章节或段落后重新汇编。"))

    score = _score_from_issues(100, issues, warnings)
    hard_gate = not any(item.get("severity") == "critical" for item in issues)
    acceptance = hard_gate and score >= 80
    report = {
        "type": "managed_long_document_delivery",
        "score": score,
        "grade": _grade(score),
        "target_words": target_words,
        "effective_chars": compact_chars,
        "minimum_chars": minimum_chars,
        "chapter_file_count": len(chapter_files),
        "missing_chapter_files": missing,
        "hard_gate_passed": hard_gate,
        "issues": issues,
        "warnings": warnings,
        "acceptance": acceptance,
    }
    return {
        "success": True,
        "result": report,
        "evidence": {
            "path": _rel(runtime, path),
            "manifest": _rel(runtime, manifest_path) if manifest_path else "",
            "exists": path.exists(),
            "bytes": path.stat().st_size,
            "score": score,
            "acceptance": acceptance,
        },
    }


def _qc_ppt(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    inspection = _ppt_inspection(path)
    slides = list(inspection.get("slides") or [])
    if not slides:
        issues.append(_issue("ppt_unreadable", "无法读取 PPT 或没有幻灯片。", "critical", "重新生成 pptx 并检查文件可打开。"))
    if len(slides) < int(args.get("min_slides", 5)):
        issues.append(_issue("too_few_slides", "幻灯片数量不足，难以形成完整商业汇报。", "medium", "补充背景、核心结论、证据、路径、决策请求。"))

    conclusion_titles = dense = weak_titles = low_content = placeholder_slides = visual_slides = 0
    placeholder_layout_slides = 0
    signatures: List[str] = []
    placeholder_pattern = re.compile(r"(待补充|占位|placeholder|lorem ipsum|(?:目标|证据|验收|下一步)\s*(?:编号|项)?\s*[：:]\s*\d+)", re.I)
    for slide in slides:
        title = str(slide.get("title", "")).strip()
        body = str(slide.get("content_text") or slide.get("text", ""))
        content_only = body[len(title):].strip() if title and body.startswith(title) else body.strip()
        weak_titles += int(len(title) < 6)
        conclusion_titles += int(any(word in title for word in ["结论", "建议", "必须", "预计", "应", "可", "将", "需要"]))
        dense += int(len(body) > 650 or body.count("\n") > 12)
        low_content += int(len(re.sub(r"\s+", "", content_only)) < 24 and int(slide.get("visual_count") or 0) == 0)
        placeholder_slides += int(bool(placeholder_pattern.search(body)))
        visual_slides += int(int(slide.get("visual_count") or 0) > 0)
        placeholder_layout_slides += int(int(slide.get("placeholder_count") or 0) > 0)
        signature = re.sub(r"\d+", "#", re.sub(r"\s+", "", body.lower()))
        signature = re.sub(r"[^a-z\u3400-\u9fff#]", "", signature)
        if signature:
            signatures.append(signature)

    duplicate_ratio = (len(signatures) - len(set(signatures))) / max(1, len(signatures))
    if slides and conclusion_titles < max(1, len(slides) // 3):
        issues.append(_issue("weak_conclusion_titles", "结论句标题比例低，偏资料堆砌。", "high", "将页面标题改成可独立阅读的结论句。"))
    if weak_titles:
        warnings.append(_issue("weak_titles", f"有 {weak_titles} 页标题过短或不明确。", "low", "补充标题中的判断/结论。"))
    if dense:
        issues.append(_issue("dense_slides", f"有 {dense} 页信息密度过高。", "medium", "拆页、压缩文字、用图表替代段落。"))
    if slides and low_content > max(1, len(slides) // 2):
        issues.append(_issue("mostly_empty_slides", f"有 {low_content}/{len(slides)} 页有效内容过少。", "critical", "补充真实结论、证据、数据或删除空泛页面。"))
    if placeholder_slides:
        issues.append(_issue("placeholder_content", f"有 {placeholder_slides} 页仍包含占位或机械编号内容。", "critical", "清除占位词并替换为真实内容。"))
    if duplicate_ratio >= 0.60 and len(slides) >= 6:
        issues.append(_issue("highly_repetitive_slides", f"页面结构化文本重复率为 {duplicate_ratio:.0%}。", "critical", "重写重复页面并建立不同证据与叙事角色。"))
    elif duplicate_ratio >= 0.30:
        issues.append(_issue("repetitive_slides", f"页面结构化文本重复率为 {duplicate_ratio:.0%}。", "high", "合并重复页面或增加差异化证据。"))

    visual_coverage = visual_slides / max(1, len(slides))
    native_visual_count = int(inspection.get("native_visual_count") or 0)
    if len(slides) >= 3 and visual_slides == 0:
        issues.append(_issue("no_meaningful_visuals", "整套演示稿没有图表、图片、表格或语义化视觉组件。", "critical", "应用设计模板，并把关键观点转成图表、证据表、路径图或信息卡片。"))
    elif len(slides) >= 5 and visual_coverage < float(args.get("min_visual_coverage", 0.4)):
        issues.append(_issue("weak_visual_coverage", f"有效视觉覆盖率仅 {visual_coverage:.0%}。", "high", "至少让40%的页面具有服务于结论的视觉表达。"))
    if len(slides) >= 8 and native_visual_count == 0:
        issues.append(_issue("no_native_evidence_visuals", "长演示稿没有图片、图表或表格证据。", "high", "在有真实数据或素材的页面加入至少一类原生证据视觉；不要编造数据。"))
    if slides and placeholder_layout_slides / len(slides) >= 0.6:
        severity = "critical" if visual_slides == 0 else "high"
        issues.append(_issue("default_placeholder_layout", f"{placeholder_layout_slides}/{len(slides)} 页仍依赖默认占位符版式。", severity, "改用无默认占位符的设计系统版式。"))
    if slides and not bool(inspection.get("is_widescreen")):
        issues.append(_issue("legacy_aspect_ratio", f"页面比例为 {inspection.get('aspect_ratio') or 'unknown'}，不是16:9宽屏。", "critical", "改为16:9宽屏后重新排版，不能只拉伸页面。"))
    if slides and not bool(inspection.get("has_explicit_fonts")):
        issues.append(_issue("missing_explicit_fonts", "没有检测到明确字体设置，结果会依赖Office默认主题。", "high", "显式设置标题和正文字体并验证中文字体回退。"))
    if not any(word in "\n".join(str(slide.get("text", "")) for slide in slides) for word in ["行动", "决策", "下一步", "建议", "CTA"]):
        issues.append(_issue("missing_cta", "缺少行动建议或决策请求。", "high", "末页补充明确决策请求/下一步。"))

    hard_gate_passed = not any(item.get("severity") == "critical" for item in issues)
    score = _score_from_issues(100, issues, warnings)
    if len(slides) >= 8 and native_visual_count == 0 and hard_gate_passed:
        score = min(score, 89)
    acceptance = hard_gate_passed and score >= 80
    report = {
        "type": "executive_ppt_delivery", "score": score, "grade": _grade(score), "slides": len(slides),
        "conclusion_title_count": conclusion_titles, "dense_slide_count": dense, "low_content_slide_count": low_content,
        "placeholder_slide_count": placeholder_slides, "layout_placeholder_slide_count": placeholder_layout_slides,
        "duplicate_ratio": round(duplicate_ratio, 4), "visual_slide_count": visual_slides,
        "visual_coverage": round(visual_coverage, 4), "native_visual_count": native_visual_count,
        "designed_visual_count": int(inspection.get("designed_visual_count") or 0), "aspect_ratio": inspection.get("aspect_ratio"),
        "is_widescreen": bool(inspection.get("is_widescreen")), "has_explicit_fonts": bool(inspection.get("has_explicit_fonts")),
        "font_names": list(inspection.get("font_names") or []), "hard_gate_passed": hard_gate_passed,
        "issues": issues, "warnings": warnings, "acceptance": acceptance,
    }
    return {"success": True, "result": report, "evidence": {"path": _rel(runtime, path), "exists": path.exists(), "bytes": path.stat().st_size, "score": score, "grade": report["grade"], "acceptance": acceptance}}


def _ppt_inspection(path: Path) -> Dict[str, Any]:
    try:
        from .ppt_design import inspect_presentation  # type: ignore

        return inspect_presentation(path)
    except Exception:
        slides: List[Dict[str, Any]] = []
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path)); fonts: set[str] = set(); placeholders = native_visuals = 0
        for slide in prs.slides:
            texts: List[str] = []; visual_count = slide_placeholders = 0
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text: texts.append(text)
                if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False) or int(getattr(shape, "shape_type", 0) or 0) == 13: visual_count += 1
                if bool(getattr(shape, "is_placeholder", False)): placeholders += 1; slide_placeholders += 1
                if bool(getattr(shape, "has_text_frame", False)):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name: fonts.add(str(run.font.name))
            native_visuals += visual_count; text = "\n".join(texts)
            slides.append({"title": texts[0] if texts else "", "text": text, "content_text": text, "visual_count": visual_count, "native_visual_count": visual_count, "designed_visual_count": 0, "placeholder_count": slide_placeholders})
        ratio = round(int(prs.slide_width) / int(prs.slide_height), 4) if int(prs.slide_height) else 0.0
        return {"slides": slides, "slide_count": len(slides), "aspect_ratio": ratio, "is_widescreen": 1.70 <= ratio <= 1.82, "placeholder_count": placeholders, "native_visual_count": native_visuals, "designed_visual_count": 0, "font_names": sorted(fonts), "has_explicit_fonts": bool(fonts)}
    except Exception:
        text = _extract_pptx_text(path)
        for chunk in re.split(r"\[slide \d+\]", text):
            lines = [item.strip() for item in chunk.splitlines() if item.strip()]
            if lines:
                joined = "\n".join(lines); slides.append({"title": lines[0], "text": joined, "content_text": joined, "visual_count": 0, "native_visual_count": 0, "designed_visual_count": 0, "placeholder_count": 0})
        return {"slides": slides, "slide_count": len(slides), "aspect_ratio": 0.0, "is_widescreen": False, "placeholder_count": 0, "native_visual_count": 0, "designed_visual_count": 0, "font_names": [], "has_explicit_fonts": False}


def _ppt_slides(path: Path) -> List[Dict[str, Any]]:
    return list(_ppt_inspection(path).get("slides") or [])


def _qc_sheet(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    rows = _read_sheet_rows(path)
    if not rows:
        issues.append(_issue("empty_sheet", "表格为空或不可读。", "critical", "重新生成表格并确保至少有表头和数据。"))
    else:
        header = rows[0]
        if any(str(h).strip() == "" for h in header):
            issues.append(_issue("blank_header", "表头存在空列。", "high", "补齐字段名。"))
        blank_cells = sum(1 for r in rows[1:] for c in r if str(c).strip() == "")
        if blank_cells > max(5, len(rows) * len(header) * 0.2):
            warnings.append(_issue("many_blank_cells", "空值比例偏高。", "low", "标注缺失原因或补齐数据。"))
        seen = set(); dup = 0
        for r in rows[1:]:
            key = tuple(str(c) for c in r)
            if key in seen: dup += 1
            seen.add(key)
        if dup:
            warnings.append(_issue("duplicate_rows", f"发现 {dup} 条重复行。", "low", "去重或说明重复原因。"))
    score = _score_from_issues(100, issues, warnings)
    return {"success": True, "result": {"type": "sheet_delivery", "score": score, "grade": _grade(score), "rows": len(rows), "cols": len(rows[0]) if rows else 0, "issues": issues, "warnings": warnings, "acceptance": score >= 80}, "evidence": {"path": _rel(runtime, path), "exists": True, "bytes": path.stat().st_size, "score": score}}


def _read_sheet_rows(path: Path) -> List[List[Any]]:
    if path.suffix.lower() == ".csv":
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            return list(csv.reader(f))[:1000]
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=False)
        try:
            ws = wb.active
            return [[c for c in r] for r in ws.iter_rows(max_row=1000, values_only=True)]
        finally:
            wb.close()
    except Exception:
        return []


def _qc_code(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    root=_resolve(runtime,target or ".",must_exist=True)
    code_suffixes={".py",".js",".mjs",".cjs",".ts",".tsx",".jsx",".java",".go",".rs",".cpp",".cc",".c",".h",".cs",".php",".rb",".swift",".kt"}
    project_type=str(args.get("project_type") or args.get("mode") or "").strip().lower()
    miniapp_mode=project_type in {"wechat_miniapp","wechat_miniprogram","miniapp","miniprogram"}
    if miniapp_mode: code_suffixes.update({".json",".wxml",".wxss"})
    candidates=[root] if root.is_file() else [p for p in root.rglob("*") if p.is_file() and ".omni_" not in p.parts and "__pycache__" not in p.parts]
    files=[p for p in candidates if p.suffix.lower() in code_suffixes]
    issues:List[Dict[str,Any]]=[]; warnings:List[Dict[str,Any]]=[]; syntax_errors=[]
    if not files: issues.append(_issue("no_code_files","目标中没有真实代码文件。","critical","提供源代码目录或生成真实代码文件。"))
    total_lines=0
    for p in files[:500]:
        txt=p.read_text(encoding="utf-8",errors="ignore"); total_lines += len([line for line in txt.splitlines() if line.strip() and not line.lstrip().startswith(("#","//"))])
        if p.suffix.lower()==".py":
            try: ast.parse(txt)
            except SyntaxError as exc: syntax_errors.append({"file":_rel(runtime,p),"line":exc.lineno,"message":exc.msg})
    if syntax_errors: issues.append(_issue("syntax_errors",f"发现 {len(syntax_errors)} 个 Python 语法错误。","critical","先修复语法错误。"))
    if miniapp_mode and root.is_dir(): issues.extend(_miniapp_project_issues(runtime,root,candidates))
    if files and total_lines < int(args.get("min_effective_lines",5)): issues.append(_issue("insufficient_code",f"有效代码行仅 {total_lines}。","critical","补充可运行实现，而不是占位文件。"))
    test_files=[p for p in candidates if "test" in p.name.lower() or "tests" in p.parts]
    if not test_files: issues.append(_issue("missing_tests","缺少测试文件或测试目录。","critical","补充测试并真实执行。"))
    test_exec={"executed":False,"returncode":None,"stdout":"","stderr":""}
    if test_files and root.is_dir():
        command=args.get("test_command")
        if miniapp_mode and not command:
            issues.append(_issue("miniapp_test_command_required","小程序工程必须提供可执行的离线测试命令。","critical","提供 node 测试脚本或项目自带测试命令。"))
            command=[]
        elif not command:
            # Delayed import: omni_body_tool imports this module at load time.
            # In frozen builds sys.executable is the backend exe, never reuse it.
            try:
                from .omni_body_tool import _resolve_python_interpreter
                command=[_resolve_python_interpreter(),"-m","pytest","-q"]
            except Exception as exc:
                issues.append(_issue("tests_not_executable",f"测试无法执行：{exc}","critical","修复测试环境和命令。"))
                command=[]
        if isinstance(command,str): command=command.split()
        if command:
            try:
                cp=subprocess.run(list(command),cwd=str(root),capture_output=True,text=True,timeout=int(args.get("timeout",180)))
                test_exec={"executed":True,"returncode":cp.returncode,"stdout":cp.stdout[-8000:],"stderr":cp.stderr[-8000:]}
                if cp.returncode != 0: issues.append(_issue("tests_failed",f"测试返回码为 {cp.returncode}。","critical","修复失败测试后重新运行。"))
            except subprocess.TimeoutExpired: issues.append(_issue("tests_timeout","测试执行超时。","critical","定位卡死测试或调整合理超时。")); test_exec["executed"]=True; test_exec["returncode"]=-1
            except Exception as exc: issues.append(_issue("tests_not_executable",f"测试无法执行：{exc}","critical","修复测试环境和命令。"))
    readme=(root/"README.md") if root.is_dir() else (root.parent/"README.md")
    if not readme.exists(): warnings.append(_issue("missing_readme","缺少 README.md。","low","补充安装、运行、测试和边界说明。"))
    smells=[]
    for p in files[:200]:
        txt=p.read_text(encoding="utf-8",errors="ignore")
        for pat in ["eval(","exec(","shell=True","pickle.loads"]:
            if pat in txt: smells.append({"file":_rel(runtime,p),"pattern":pat})
    if smells: issues.append(_issue("security_smells",f"发现 {len(smells)} 个敏感模式。","medium","逐项确认并限制。"))
    hard=not any(x.get("severity")=="critical" for x in issues); score=_score_from_issues(100,issues,warnings); acceptance=hard and score>=80
    report={"type":"code_project_delivery","score":score,"grade":_grade(score),"files_checked":len(files),"total_lines":total_lines,"syntax_errors":syntax_errors,"test_files":[_rel(runtime,p) for p in test_files[:50]],"test_execution":test_exec,"security_smells":smells[:50],"hard_gate_passed":hard,"issues":issues,"warnings":warnings,"acceptance":acceptance}
    return {"success":True,"result":report,"evidence":{"path":_rel(runtime,root),"exists":True,"bytes":0,"score":score,"acceptance":acceptance}}


def _miniapp_project_issues(runtime: Any, root: Path, candidates: List[Path]) -> List[Dict[str, Any]]:
    issues: List[Dict[str, Any]] = []
    by_relative = {path.relative_to(root).as_posix(): path for path in candidates}
    for required in ("app.js", "app.json", "project.config.json"):
        if required not in by_relative:
            issues.append(_issue("miniapp_required_file", f"小程序缺少 {required}。", "critical", f"创建有效的 {required}。"))
    app_config: Dict[str, Any] = {}
    for relative, path in by_relative.items():
        if path.suffix.lower() != ".json":
            continue
        try:
            value = json.loads(path.read_text(encoding="utf-8", errors="strict"))
            if not isinstance(value, dict):
                raise ValueError("JSON root must be an object")
            if relative == "app.json":
                app_config = value
        except (OSError, UnicodeError, ValueError, json.JSONDecodeError) as exc:
            issues.append(_issue("miniapp_invalid_json", f"{relative} 不是有效 JSON 对象：{exc}", "critical", "修复 JSON 语法和根结构。"))
    pages = app_config.get("pages") if isinstance(app_config, dict) else None
    if not isinstance(pages, list) or not pages or not all(isinstance(page, str) and page.strip() for page in pages):
        issues.append(_issue("miniapp_pages_invalid", "app.json.pages 必须是非空页面路径数组。", "critical", "声明至少一个真实页面。"))
        pages = []
    for page in pages:
        normalized = page.replace("\\", "/").strip("/")
        if not normalized or ".." in Path(normalized).parts:
            issues.append(_issue("miniapp_page_path_unsafe", f"页面路径不安全：{page}", "critical", "只使用项目内规范相对路径。"))
            continue
        for suffix in (".js", ".wxml", ".wxss"):
            expected = normalized + suffix
            if expected not in by_relative:
                issues.append(_issue("miniapp_page_file_missing", f"页面缺少 {expected}。", "critical", "补齐页面 JS/WXML/WXSS 文件。"))
    return issues


def _qc_research(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True) if target else None
    text = (str(args.get("content") or "") + "\n" + (_read_text_any(path) if path else "")).strip()
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    required = {
        "research_question": ["研究问题", "问题", "PICO", "RQ"],
        "search_strategy": ["搜索策略", "关键词", "数据库", "检索"],
        "screening": ["纳入", "排除", "筛选"],
        "evidence_table": ["证据表", "样本", "研究", "作者", "年份"],
        "limitations": ["局限", "限制", "偏倚", "不确定"],
    }
    for code, keys in required.items():
        if not any(k.lower() in text.lower() for k in keys):
            issues.append(_issue(f"missing_{code}", f"缺少研究交付关键模块：{code}", "high", f"补充 {code}。"))
    citation_like = len(re.findall(r"(doi\.org|https?://|\[\d+\]|（\d{4}）|\(\d{4}\))", text, flags=re.I))
    if citation_like < 3:
        issues.append(_issue("weak_citations", "引用/来源线索不足。", "high", "补充可追溯来源链接、DOI、年份或引用编号。"))
    if "PRISMA" not in text and "系统综述" in text:
        warnings.append(_issue("prisma_not_named", "系统综述类任务未显式使用 PRISMA 检查项。", "low", "按 PRISMA 2020 核查标题、摘要、方法、结果和流程图。"))
    score = _score_from_issues(100, issues, warnings)
    return {"success": True, "result": {"type": "research_review_delivery", "score": score, "grade": _grade(score), "citation_like_count": citation_like, "issues": issues, "warnings": warnings, "acceptance": score >= 80}, "evidence": {"path": _rel(runtime, path) if path else "content", "exists": bool(path), "bytes": len(text.encode('utf-8')), "score": score}}


def _qc_video(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    info = _ffprobe(runtime, path)
    if not info:
        issues.append(_issue("video_unreadable", "视频不可读或 ffprobe 不可用。", "critical", "重新导出 MP4 并验证可播放。"))
    duration = float(info.get("duration", 0) or 0)
    width = int(info.get("width", 0) or 0); height = int(info.get("height", 0) or 0)
    if duration and duration > float(args.get("max_duration", 180)):
        issues.append(_issue("too_long", f"短视频时长 {duration:.1f}s 超过目标。", "medium", "压缩节奏或裁剪冗余段落。"))
    if width and height and height < width:
        warnings.append(_issue("not_vertical", "视频不是竖屏，移动端沉浸感不足。", "low", "按 9:16 重新构图或导出。"))
    meta_text = str(args.get("script") or args.get("brief") or "")
    if not any(x in meta_text for x in ["钩子", "前3秒", "hook", "开头"]):
        issues.append(_issue("missing_hook_spec", "缺少前3秒钩子说明。", "high", "补充开头钩子、视觉冲击点和第一句字幕。"))
    if not any(x in meta_text for x in ["CTA", "行动", "联系", "点击", "私信"]):
        warnings.append(_issue("missing_cta_spec", "缺少 CTA 设计。", "low", "在结尾2-4秒加入明确行动提示。"))
    score = _score_from_issues(100, issues, warnings)
    return {"success": True, "result": {"type": "short_video_delivery", "score": score, "grade": _grade(score), "video_info": info, "issues": issues, "warnings": warnings, "acceptance": score >= 80}, "evidence": {"path": _rel(runtime, path), "exists": True, "bytes": path.stat().st_size, "score": score}}


def _ffprobe(runtime: Any, path: Path) -> Dict[str, Any]:
    ffprobe = getattr(runtime, "ffprobe", None) or shutil.which("ffprobe")
    if not ffprobe:
        return {}
    try:
        cmd = [ffprobe, "-v", "error", "-select_streams", "v:0", "-show_entries", "stream=width,height:format=duration", "-of", "json", str(path)]
        out = subprocess.run(cmd, capture_output=True, text=True, timeout=20)
        if out.returncode != 0:
            return {}
        data = json.loads(out.stdout or "{}")
        streams = data.get("streams") or []
        if not streams:
            return {}
        fmt = data.get("format") or {}
        width = streams[0].get("width")
        height = streams[0].get("height")
        if not width or not height:
            return {}
        return {"duration": float(fmt.get("duration") or 0), "width": width, "height": height}
    except Exception:
        return {}


def _qc_image(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    path = _resolve(runtime, target, must_exist=True)
    issues: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []
    info: Dict[str, Any] = {}
    try:
        from PIL import Image, ImageStat  # type: ignore
        with Image.open(path) as im:
            info = {"width": im.width, "height": im.height, "mode": im.mode, "format": im.format}
            if im.width < int(args.get("min_width", 1080)) or im.height < int(args.get("min_height", 1080)):
                warnings.append(_issue("low_resolution", "图片分辨率偏低。", "low", "导出更高分辨率版本。"))
            gray = im.convert("L")
            stat = ImageStat.Stat(gray)
            if stat.stddev and stat.stddev[0] < 25:
                warnings.append(_issue("low_contrast_proxy", "整体对比度偏低，文字可读性可能受影响。", "low", "增强标题区对比度或增加遮罩。"))
    except Exception as exc:
        issues.append(_issue("image_unreadable", f"图片不可读：{exc}", "critical", "重新导出 PNG/JPG。"))
    score = _score_from_issues(100, issues, warnings)
    return {"success": True, "result": {"type": "image_delivery", "score": score, "grade": _grade(score), "image_info": info, "issues": issues, "warnings": warnings, "acceptance": score >= 80}, "evidence": {"path": _rel(runtime, path), "exists": True, "bytes": path.stat().st_size, "score": score}}


def _qc_writing(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    text = str(args.get("content") or "")
    path = None
    if target:
        path = _resolve(runtime, target, must_exist=True)
        text += "\n" + _read_text_any(path)
    issues = _ai_tone_issues(text)
    stats = _sentence_stats(text)
    if stats["avg_sentence_chars"] > 80:
        issues.append(_issue("avg_sentence_too_long", "平均句长偏长，像AI堆叠。", "medium", "拆分长句，减少复合从句。"))
    score = _score_from_issues(100, issues, [])
    return {"success": True, "result": {"type": "writing_ai_tone", "score": score, "grade": _grade(score), "stats": stats, "issues": issues, "acceptance": score >= 80}, "evidence": {"path": _rel(runtime, path) if path else "content", "exists": bool(path), "bytes": len(text.encode('utf-8')), "score": score}}


def _ai_tone_issues(text: str) -> List[Dict[str, Any]]:
    issues: List[Dict[str, Any]] = []
    for phrase in GENERIC_AI_PHRASES:
        count = text.count(phrase)
        if count >= 2:
            issues.append(_issue("generic_phrase", f"泛化表达重复：{phrase} ×{count}", "low", f"替换“{phrase}”为具体事实、数据或动作。"))
    vague_claims = len(re.findall(r"(显著|大幅|全面|有效|深度|极大|明显).{0,6}(提升|优化|改善|增强)", text))
    if vague_claims > 5:
        issues.append(_issue("vague_claims", "模糊效果词过多。", "medium", "为每个效果词补充量化口径或删掉。"))
    return issues[:20]


def _writing_outline_create(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    outline_type = str(args.get("type") or "business_proposal")
    output = _resolve(runtime, target or f"{outline_type}_outline.md")
    variables = args.get("variables") if isinstance(args.get("variables"), dict) else {}
    content = _template_skeleton(outline_type, variables)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(content, encoding="utf-8")
    return {"success": True, "output": {"path": _rel(runtime, output), "exists": True, "bytes": output.stat().st_size, "outline_type": outline_type}}


def _research_evidence_table_create(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    output = _resolve(runtime, target or "research_evidence_table.csv")
    sources = args.get("sources") if isinstance(args.get("sources"), list) else []
    headers = ["id", "title", "year", "source", "method", "sample", "key_finding", "limitations", "relevance", "url_or_doi"]
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=headers)
        writer.writeheader()
        for idx, s in enumerate(sources, 1):
            row = {h: "" for h in headers}
            row.update(s if isinstance(s, dict) else {"title": str(s)})
            row["id"] = row.get("id") or idx
            writer.writerow({h: row.get(h, "") for h in headers})
    return {"success": True, "output": {"path": _rel(runtime, output), "exists": True, "bytes": output.stat().st_size, "rows": len(sources)}}


def _repair_plan(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    output = _resolve(runtime, target or args.get("output") or "repair_plan.md")
    issues = args.get("issues") if isinstance(args.get("issues"), list) else []
    source_action = args.get("source_action", "qc")
    lines = ["# 返工计划", "", f"- 来源动作：{source_action}", f"- 生成时间：{time.strftime('%Y-%m-%d %H:%M:%S')}", "", "## 待处理问题"]
    if not issues:
        lines.append("暂无问题。")
    for i, issue in enumerate(issues, 1):
        if isinstance(issue, dict):
            lines.append(f"{i}. **{issue.get('severity','')} / {issue.get('code','')}**：{issue.get('message','')}")
            lines.append(f"   - 修复：{issue.get('repair','')}")
        else:
            lines.append(f"{i}. {issue}")
    lines.append("\n## 返工原则\n先修 critical/high，再修 medium，最后处理 low；每次返工后重新运行对应 qc.* delivery_check。")
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text("\n".join(lines), encoding="utf-8")
    return {"success": True, "output": {"path": _rel(runtime, output), "exists": True, "bytes": output.stat().st_size, "issue_count": len(issues)}}


def _deliverable_package(runtime: Any, target: str | None, args: Dict[str, Any]) -> Dict[str, Any]:
    output = _resolve(runtime, target or args.get("output") or "delivery_package.zip")
    items = args.get("items") if isinstance(args.get("items"), list) else []
    if output.suffix.lower() != ".zip":
        raise ValueError("deliverable.package output must be a new .zip file")
    if output.exists():
        raise FileExistsError("deliverable.package refuses to overwrite an existing output")
    if not items or any(not isinstance(item, str) or not item.strip() for item in items):
        raise ValueError("deliverable.package requires a non-empty items list")

    output_resolved = output.resolve(strict=False)
    entries: List[Tuple[Path, str]] = []
    manifest_items: List[Dict[str, Any]] = []
    seen_archives: set[str] = set()
    for item in items:
        source = _resolve(runtime, item, must_exist=True)
        if source.is_symlink() or (not source.is_file() and not source.is_dir()):
            raise ValueError("deliverable.package items must be regular files or directories")
        source_resolved = source.resolve(strict=True)
        if source_resolved == output_resolved:
            raise ValueError("deliverable.package input and output must be different paths")
        if source.is_dir() and output_resolved.is_relative_to(source_resolved):
            raise ValueError("deliverable.package output cannot be inside an input directory")
        source_files = [source] if source.is_file() else sorted(source.rglob("*"))
        file_count = 0
        total_bytes = 0
        for child in source_files:
            if child.is_symlink():
                raise ValueError("deliverable.package refuses symbolic links")
            if not child.is_file():
                continue
            child_resolved = child.resolve(strict=True)
            if child_resolved == output_resolved:
                raise ValueError("deliverable.package input and output must be different paths")
            # Items may live in user-authorized roots outside the workspace;
            # those are archived relative to their own source directory.
            archive_name = (
                child.name
                if source.is_file()
                else (
                    child.relative_to(runtime.workspace).as_posix()
                    if child.is_relative_to(runtime.workspace)
                    else f"{source.name}/{child.relative_to(source).as_posix()}"
                )
            )
            folded = archive_name.casefold()
            if folded in seen_archives or folded == "delivery_manifest.json":
                raise ValueError("deliverable.package contains a duplicate archive path")
            seen_archives.add(folded)
            entries.append((child, archive_name))
            file_count += 1
            total_bytes += child.stat().st_size
        manifest_items.append(
            {
                "path": _rel(runtime, source),
                "type": "file" if source.is_file() else "dir",
                "file_count": file_count,
                "bytes": total_bytes,
            }
        )
    if not entries:
        raise ValueError("deliverable.package cannot create a manifest-only archive")

    manifest = {
        "schema": "tiangong.v3.delivery_package.v1",
        "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "items": manifest_items,
        "notes": args.get("notes", ""),
    }
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(
        f".{output.name}.{os.getpid()}.{time.time_ns()}.tmp"
    )
    try:
        with temporary.open("xb") as stream:
            with zipfile.ZipFile(stream, "w", zipfile.ZIP_DEFLATED, allowZip64=True) as zf:
                for source, archive_name in entries:
                    zf.write(source, archive_name)
                zf.writestr(
                    "DELIVERY_MANIFEST.json",
                    json.dumps(manifest, ensure_ascii=False, indent=2),
                )
            stream.flush()
            os.fsync(stream.fileno())
        with zipfile.ZipFile(temporary, "r") as verify:
            if verify.testzip() is not None:
                raise ValueError("deliverable.package temporary archive failed CRC verification")
            expected_names = {archive_name for _, archive_name in entries}
            expected_names.add("DELIVERY_MANIFEST.json")
            if set(verify.namelist()) != expected_names:
                raise ValueError("deliverable.package temporary archive membership mismatch")
            stored_manifest = json.loads(verify.read("DELIVERY_MANIFEST.json").decode("utf-8"))
            if stored_manifest != manifest:
                raise ValueError("deliverable.package temporary manifest readback mismatch")
        if output.exists():
            raise FileExistsError("deliverable.package output appeared before atomic commit")
        os.rename(temporary, output)
    finally:
        if temporary.exists():
            temporary.unlink()
    return {
        "success": True,
        "output": {
            "path": _rel(runtime, output),
            "exists": output.is_file(),
            "bytes": output.stat().st_size,
            "item_count": len(manifest["items"]),
            "file_count": len(entries),
        },
    }
