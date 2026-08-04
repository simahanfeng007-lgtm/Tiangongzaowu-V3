"""Deterministic, design-aware PowerPoint generation for ``omni_body``.

The model still chooses the workflow through a Skill.  This module owns the
mechanical design work so the model only needs to provide slide semantics and a
short design-spec reference instead of reproducing layout instructions in the
conversation context.
"""
from __future__ import annotations

import json
import math
import re
from pathlib import Path
from typing import Any, Callable, Dict, Iterable, List, Tuple


_FALLBACK_DESIGN: Dict[str, Any] = {
    "schema": "tiangong.v3.ppt_design.v1",
    "template_id": "executive_ppt",
    "style_name": "executive_blue",
    "aspect_ratio": "16:9",
    "brand": "TIANGONG",
    "fonts": {"heading": "Microsoft YaHei", "body": "Microsoft YaHei", "latin": "Aptos"},
    "colors": {
        "background": "#F4F7FB",
        "surface": "#FFFFFF",
        "ink": "#10233F",
        "muted": "#5E6F86",
        "primary": "#1261FF",
        "accent": "#00A9B8",
        "dark": "#0A1733",
        "line": "#DCE4EF",
        "success": "#16A36A",
        "warning": "#F59E0B",
    },
    "layout": {"margin_x": 0.72, "title_y": 0.5, "content_y": 1.62, "footer_y": 7.12, "max_items_per_slide": 6},
    "quality_contract": {
        "requires_widescreen": True,
        "forbid_default_placeholders": True,
        "minimum_meaningful_visual_coverage": 0.4,
        "require_explicit_fonts": True,
    },
}

_STYLE_PRESETS: Dict[str, Dict[str, Any]] = {
    "executive_blue": {},
    "professional": {},
    "business": {},
    "tech_cyan": {
        "colors": {"primary": "#006DFF", "accent": "#00C2C7", "dark": "#06182E", "background": "#F2F8FA"}
    },
    "warm_amber": {
        "colors": {"primary": "#C65A11", "accent": "#E9A23B", "dark": "#2B1A12", "background": "#FBF7F1"}
    },
}

_HEX_COLOR = re.compile(r"^#[0-9A-Fa-f]{6}$")
_HEADING = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
_LIST_ITEM = re.compile(r"^\s*(?:[-*+]\s+|\d+[.)、]\s*)(.+?)\s*$")
_TABLE_DIVIDER = re.compile(r"^\s*\|?(?:\s*:?-{3,}:?\s*\|)+\s*$")


def template_asset_root() -> Path:
    """Locate templates in source, frozen-module, and legacy mirror layouts."""
    here = Path(__file__).resolve()
    candidates = [here.parents[1]]
    for parent in here.parents:
        if parent.name in {"frozen_modules", "legacy_pyz_modules"}:
            candidates.append(parent.parent / "omni_body_skill")
    for candidate in candidates:
        if (candidate / "templates" / "manifest.json").is_file():
            return candidate
    return candidates[0]


def _deep_merge(base: Dict[str, Any], patch: Dict[str, Any]) -> Dict[str, Any]:
    result: Dict[str, Any] = {}
    for key, value in base.items():
        result[key] = _deep_merge(value, {}) if isinstance(value, dict) else value
    for key, value in patch.items():
        if isinstance(value, dict) and isinstance(result.get(key), dict):
            result[key] = _deep_merge(result[key], value)
        else:
            result[key] = value
    return result


def _load_json_object(path: Path) -> Dict[str, Any]:
    if not path.is_file() or path.is_symlink():
        raise ValueError(f"PPT design spec is not a regular file: {path}")
    if path.stat().st_size > 256 * 1024:
        raise ValueError("PPT design spec exceeds 256 KiB")
    value = json.loads(path.read_text(encoding="utf-8", errors="strict"))
    if not isinstance(value, dict):
        raise ValueError("PPT design spec root must be a JSON object")
    return value


def _validate_design(design: Dict[str, Any]) -> Dict[str, Any]:
    if str(design.get("aspect_ratio") or "16:9") != "16:9":
        raise ValueError("pptx.create currently requires a 16:9 design spec")
    colors = design.get("colors")
    fonts = design.get("fonts")
    layout = design.get("layout")
    if not isinstance(colors, dict) or not isinstance(fonts, dict) or not isinstance(layout, dict):
        raise ValueError("PPT design spec requires colors, fonts, and layout objects")
    for key in ("background", "surface", "ink", "muted", "primary", "accent", "dark", "line"):
        value = str(colors.get(key) or "")
        if not _HEX_COLOR.fullmatch(value):
            raise ValueError(f"PPT design color {key} must be #RRGGBB")
        colors[key] = value.upper()
    for key in ("heading", "body"):
        value = str(fonts.get(key) or "").strip()
        if not value or len(value) > 80:
            raise ValueError(f"PPT design font {key} must be a non-empty font family")
        fonts[key] = value
    maximum = layout.get("max_items_per_slide", 6)
    if not isinstance(maximum, int) or isinstance(maximum, bool) or not 2 <= maximum <= 8:
        raise ValueError("PPT design max_items_per_slide must be an integer from 2 to 8")
    return design


def load_design_spec(
    resolve: Callable[..., Path],
    args: Dict[str, Any],
    source_path: Path | None = None,
) -> Tuple[Dict[str, Any], str]:
    """Load shipped defaults, then a sidecar or explicit design override."""
    template_id = str(args.get("template_id") or "executive_ppt").strip().lower()
    if not re.fullmatch(r"[a-z0-9_-]{1,64}", template_id):
        raise ValueError("pptx.create template_id contains unsupported characters")
    design = _deep_merge({}, _FALLBACK_DESIGN)
    design["template_id"] = template_id
    sources: List[str] = ["builtin:fallback"]

    shipped = template_asset_root() / "templates" / f"{template_id}.design.json"
    if shipped.is_file():
        design = _deep_merge(design, _load_json_object(shipped))
        sources.append(str(shipped))

    style = args.get("style")
    if isinstance(style, str) and style.strip():
        style_name = style.strip().lower()
        design = _deep_merge(design, _STYLE_PRESETS.get(style_name, {}))
        design["style_name"] = style_name
    elif isinstance(style, dict):
        design = _deep_merge(design, style)
        sources.append("args.style")

    explicit = args.get("design_spec")
    sidecar: Path | None = None
    if isinstance(explicit, str) and explicit.strip():
        sidecar = resolve(explicit, must_exist=True)
    elif isinstance(explicit, dict):
        design = _deep_merge(design, explicit)
        sources.append("args.design_spec")
    elif explicit is not None:
        raise ValueError("pptx.create design_spec must be a JSON object or workspace .json path")
    elif source_path is not None:
        candidate = source_path.with_suffix(".design.json")
        if candidate.is_file() and not candidate.is_symlink():
            sidecar = candidate
    if sidecar is not None:
        if sidecar.suffix.lower() != ".json":
            raise ValueError("pptx.create design_spec path must end with .json")
        design = _deep_merge(design, _load_json_object(sidecar))
        sources.append(str(sidecar))

    return _validate_design(design), " -> ".join(sources)


def _clean_inline(text: str) -> str:
    value = re.sub(r"[`*_~]", "", str(text or ""))
    value = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", value)
    return re.sub(r"\s+", " ", value).strip()


def _parse_table(lines: List[str]) -> Dict[str, Any] | None:
    rows: List[List[str]] = []
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|") or not stripped.endswith("|") or _TABLE_DIVIDER.fullmatch(stripped):
            continue
        cells = [_clean_inline(cell) for cell in stripped.strip("|").split("|")]
        if any(cells):
            rows.append(cells)
    if len(rows) < 2:
        return None
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]
    return {"headers": padded[0], "rows": padded[1:]}


def _section_to_spec(section: str, index: int, total: int) -> Dict[str, Any] | None:
    lines = [line.rstrip() for line in section.strip().splitlines()]
    while lines and not lines[0].strip():
        lines.pop(0)
    if not lines:
        return None
    heading = _HEADING.match(lines[0].strip())
    title = _clean_inline(heading.group(2) if heading else lines[0])
    body_lines = lines[1:]
    items: List[str] = []
    notes: List[str] = []
    in_notes = False
    for raw in body_lines:
        stripped = raw.strip()
        if not stripped:
            continue
        if re.match(r"^(?:备注|讲稿|speaker\s*notes?)\s*[:：]", stripped, re.I):
            in_notes = True
            note = re.split(r"[:：]", stripped, maxsplit=1)[-1].strip()
            if note:
                notes.append(_clean_inline(note))
            continue
        if in_notes:
            notes.append(_clean_inline(stripped))
            continue
        if stripped.startswith("|"):
            continue
        match = _LIST_ITEM.match(stripped)
        items.append(_clean_inline(match.group(1) if match else stripped))
    table = _parse_table(body_lines)
    is_cover = bool(
        index == 0
        and (
            (heading is not None and len(heading.group(1)) == 1)
            or "封面" in title
            or (total > 1 and len(items) <= 2)
        )
    )
    spec: Dict[str, Any] = {"title": title, "items": [item for item in items if item], "is_cover": is_cover}
    if notes:
        spec["notes"] = "\n".join(notes)
    if table:
        spec["table"] = table
    return spec


def parse_markdown_slides(content: str) -> List[Dict[str, Any]]:
    """Parse compact slide Markdown without dropping any non-empty section."""
    normalized = str(content or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not normalized:
        return []
    if re.search(r"(?m)^\s*---+\s*$", normalized):
        sections = re.split(r"(?m)^\s*---+\s*$", normalized)
    else:
        sections: List[str] = []
        current: List[str] = []
        for line in normalized.splitlines():
            if re.match(r"^##\s+", line) and current:
                sections.append("\n".join(current))
                current = [line]
            else:
                current.append(line)
        if current:
            sections.append("\n".join(current))
    specs = [_section_to_spec(section, index, len(sections)) for index, section in enumerate(sections)]
    return [spec for spec in specs if spec and spec.get("title")]


def structured_slides(args: Dict[str, Any]) -> List[Dict[str, Any]]:
    specs: List[Dict[str, Any]] = []
    title = str(args.get("title") or "").strip()
    if title:
        subtitle = str(args.get("subtitle") or "").strip()
        specs.append({"title": title, "items": [subtitle] if subtitle else [], "is_cover": True})
    raw_slides = args.get("slides") if isinstance(args.get("slides"), list) else []
    for raw in raw_slides:
        if not isinstance(raw, dict):
            continue
        spec = dict(raw)
        spec["title"] = _clean_inline(str(raw.get("title") or ""))
        items = raw.get("bullets") if raw.get("bullets") is not None else raw.get("body")
        if isinstance(items, str):
            items = [line for line in items.splitlines() if line.strip()]
        if not isinstance(items, list):
            items = []
        spec["items"] = [_clean_inline(str(item)) for item in items if str(item).strip()]
        spec["is_cover"] = bool(raw.get("is_cover", False))
        if spec["title"]:
            specs.append(spec)
    return specs


def _expand_dense_specs(specs: List[Dict[str, Any]], maximum: int) -> List[Dict[str, Any]]:
    expanded: List[Dict[str, Any]] = []
    for spec in specs:
        items = list(spec.get("items") or [])
        if spec.get("is_cover") or spec.get("table") or spec.get("chart") or spec.get("image") or len(items) <= maximum:
            expanded.append(spec)
            continue
        for chunk_index in range(0, len(items), maximum):
            clone = dict(spec)
            clone["items"] = items[chunk_index : chunk_index + maximum]
            if chunk_index:
                clone["title"] = f"{spec.get('title', '')}（续）"
                clone.pop("notes", None)
            expanded.append(clone)
    return expanded


def _rgb(value: str):
    from pptx.dml.color import RGBColor

    raw = str(value).lstrip("#")
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _fill(shape: Any, color: str, line_color: str | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
    else:
        shape.line.fill.background()


def _add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    font: str,
    size: float,
    color: str,
    bold: bool = False,
    align: str = "left",
    valign: str = "top",
    name: str = "TGText",
) -> Any:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = {
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return shape


def _add_rect(slide: Any, x: float, y: float, width: float, height: float, color: str, *, radius: bool = True, line: str | None = None, name: str = "TGDecor") -> Any:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.name = name
    _fill(shape, color, line)
    return shape


def _set_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _font_size_for_item(text: str, count: int) -> float:
    length = len(str(text or ""))
    if length > 120:
        return 14
    if length > 75 or count >= 6:
        return 16
    if length > 42 or count >= 4:
        return 18
    return 20


def _render_cover(slide: Any, spec: Dict[str, Any], design: Dict[str, Any]) -> None:
    colors = design["colors"]
    fonts = design["fonts"]
    _set_background(slide, colors["dark"])
    _add_rect(slide, 0, 0, 13.333, 0.10, colors["accent"], radius=False, name="TGDecor:accent-bar")
    _add_text(slide, str(design.get("brand") or "TIANGONG"), 0.78, 0.55, 3.0, 0.35, font=fonts["body"], size=11, color=colors["accent"], bold=True, name="TGText:brand")
    title = str(spec.get("title") or "")
    title_size = 34 if len(title) > 34 else 42
    _add_text(slide, title, 0.78, 1.45, 7.6, 2.45, font=fonts["heading"], size=title_size, color="#FFFFFF", bold=True, valign="middle", name="TGTitle")
    items = [str(item) for item in spec.get("items") or []]
    if items:
        _add_text(slide, " · ".join(items[:2]), 0.82, 4.25, 7.0, 0.9, font=fonts["body"], size=17, color="#C8D4E8", name="TGText:subtitle")
    _add_rect(slide, 9.15, 1.05, 3.25, 5.25, "#102A52", name="TGVisual:cover-panel")
    for index, (label, color) in enumerate((("洞察", colors["primary"]), ("证据", colors["accent"]), ("决策", colors.get("success", "#16A36A")))):
        y = 1.62 + index * 1.35
        _add_rect(slide, 9.65, y, 2.2, 0.82, color, name=f"TGVisual:cover-step:{index + 1}")
        _add_text(slide, f"0{index + 1}  {label}", 9.82, y + 0.17, 1.85, 0.4, font=fonts["heading"], size=15, color="#FFFFFF", bold=True, name="TGText:cover-step")
    _add_text(slide, "先结论 · 后证明 · 明确行动", 9.42, 5.65, 2.7, 0.35, font=fonts["body"], size=10, color="#9FB1CA", align="center", name="TGText:cover-caption")


def _render_header(slide: Any, title: str, page_number: int, design: Dict[str, Any]) -> None:
    colors = design["colors"]
    fonts = design["fonts"]
    _set_background(slide, colors["background"])
    _add_text(slide, f"{page_number:02d}", 0.72, 0.42, 0.52, 0.34, font=fonts["body"], size=11, color=colors["primary"], bold=True, name="TGText:page-kicker")
    _add_text(slide, title, 1.35, 0.38, 10.9, 0.86, font=fonts["heading"], size=25 if len(title) > 34 else 29, color=colors["ink"], bold=True, valign="middle", name="TGTitle")
    _add_rect(slide, 0.72, 1.31, 11.9, 0.025, colors["line"], radius=False, name="TGDecor:header-line")


def _render_footer(slide: Any, page_number: int, design: Dict[str, Any]) -> None:
    colors = design["colors"]
    fonts = design["fonts"]
    _add_text(slide, str(design.get("brand") or "TIANGONG"), 0.74, 7.08, 2.0, 0.20, font=fonts["body"], size=8, color=colors["muted"], name="TGText:footer-brand")
    _add_text(slide, str(page_number), 12.05, 7.02, 0.55, 0.25, font=fonts["body"], size=9, color=colors["muted"], align="right", name="TGText:footer-page")


def _render_cards(slide: Any, items: List[str], design: Dict[str, Any]) -> None:
    colors = design["colors"]
    fonts = design["fonts"]
    if not items:
        raise ValueError(
            "slide has no content items; refusing to render placeholder card content. "
            "Provide non-empty items for every slide before calling pptx.create."
        )
    count = len(items)
    columns = 1 if count == 1 else 2
    rows = int(math.ceil(count / columns))
    gap_x, gap_y = 0.28, 0.25
    area_x, area_y, area_w, area_h = 0.76, 1.62, 11.82, 4.98
    card_w = (area_w - gap_x * (columns - 1)) / columns
    card_h = (area_h - gap_y * (rows - 1)) / rows
    accents = [colors["primary"], colors["accent"], colors.get("success", colors["primary"]), colors.get("warning", colors["accent"])]
    for index, item in enumerate(items):
        row, column = divmod(index, columns)
        x = area_x + column * (card_w + gap_x)
        y = area_y + row * (card_h + gap_y)
        _add_rect(slide, x, y, card_w, card_h, colors["surface"], line=colors["line"], name=f"TGVisual:card:{index + 1}")
        _add_rect(slide, x, y, 0.08, card_h, accents[index % len(accents)], radius=False, name="TGDecor:card-accent")
        _add_text(slide, f"{index + 1:02d}", x + 0.25, y + 0.20, 0.55, 0.28, font=fonts["body"], size=9, color=accents[index % len(accents)], bold=True, name="TGText:card-index")
        size = _font_size_for_item(item, count)
        _add_text(slide, item, x + 0.25, y + 0.58, card_w - 0.55, max(0.55, card_h - 0.78), font=fonts["body"], size=size, color=colors["ink"], bold=count <= 2, valign="middle" if count <= 2 else "top", name="TGText:card-body")


def _render_timeline(slide: Any, items: List[str], design: Dict[str, Any]) -> None:
    colors = design["colors"]
    fonts = design["fonts"]
    count = max(1, len(items))
    area_x, area_w = 0.80, 11.72
    gap = 0.18
    card_w = (area_w - gap * (count - 1)) / count
    _add_rect(slide, 1.0, 3.05, 11.2, 0.04, colors["line"], radius=False, name="TGDecor:timeline-line")
    for index, item in enumerate(items):
        x = area_x + index * (card_w + gap)
        accent = colors["primary"] if index % 2 == 0 else colors["accent"]
        _add_rect(slide, x, 1.82, card_w, 3.85, colors["surface"], line=colors["line"], name=f"TGVisual:timeline:{index + 1}")
        _add_rect(slide, x + card_w / 2 - 0.22, 2.84, 0.44, 0.44, accent, name="TGDecor:timeline-node")
        _add_text(slide, f"STEP {index + 1}", x + 0.16, 2.03, card_w - 0.32, 0.35, font=fonts["body"], size=9, color=accent, bold=True, align="center", name="TGText:timeline-step")
        _add_text(slide, item, x + 0.17, 3.52, card_w - 0.34, 1.55, font=fonts["body"], size=14 if len(item) > 55 else 16, color=colors["ink"], bold=True, align="center", valign="middle", name="TGText:timeline-body")


def _render_table(slide: Any, table_spec: Dict[str, Any], design: Dict[str, Any]) -> bool:
    from pptx.util import Inches, Pt

    headers = table_spec.get("headers") if isinstance(table_spec.get("headers"), list) else []
    rows = table_spec.get("rows") if isinstance(table_spec.get("rows"), list) else []
    if not headers or not rows:
        return False
    width = len(headers)
    clean_rows = [list(row)[:width] if isinstance(row, list) else [row] for row in rows[:7]]
    shape = slide.shapes.add_table(len(clean_rows) + 1, width, Inches(0.78), Inches(1.72), Inches(11.78), Inches(4.76))
    shape.name = "TGVisual:table"
    table = shape.table
    colors = design["colors"]
    fonts = design["fonts"]
    for row_index, row in enumerate([headers, *clean_rows]):
        for column_index in range(width):
            cell = table.cell(row_index, column_index)
            cell.text = str(row[column_index] if column_index < len(row) else "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(colors["dark"] if row_index == 0 else (colors["surface"] if row_index % 2 else colors["background"]))
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = fonts["body"]
                paragraph.font.size = Pt(12 if len(clean_rows) > 5 else 14)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = _rgb("#FFFFFF" if row_index == 0 else colors["ink"])
    return True


def _render_chart(slide: Any, chart_spec: Dict[str, Any], design: Dict[str, Any]) -> bool:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    categories = chart_spec.get("categories")
    series = chart_spec.get("series")
    if not isinstance(categories, list) or not categories or not isinstance(series, list) or not series:
        return False
    data = ChartData()
    data.categories = [str(item) for item in categories[:8]]
    valid = 0
    for item in series[:4]:
        if not isinstance(item, dict) or not isinstance(item.get("values"), list):
            continue
        values = item["values"][: len(data.categories)]
        try:
            numeric = [float(value) for value in values]
        except (TypeError, ValueError):
            continue
        data.add_series(str(item.get("name") or f"Series {valid + 1}"), numeric)
        valid += 1
    if not valid:
        return False
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if str(chart_spec.get("type") or "column").lower() != "bar" else XL_CHART_TYPE.BAR_CLUSTERED
    chart_shape = slide.shapes.add_chart(chart_type, Inches(0.92), Inches(1.72), Inches(11.45), Inches(4.78), data)
    chart_shape.name = "TGVisual:chart"
    chart = chart_shape.chart
    chart.has_legend = valid > 1
    chart.has_title = False
    chart.value_axis.has_major_gridlines = True
    colors = design["colors"]
    palette = [colors["primary"], colors["accent"], colors.get("success", "#16A36A"), colors.get("warning", "#F59E0B")]
    for index, chart_series in enumerate(chart.series):
        chart_series.format.fill.solid()
        chart_series.format.fill.fore_color.rgb = _rgb(palette[index % len(palette)])
    return True


def _render_picture(slide: Any, image_path: Path, items: List[str], design: Dict[str, Any]) -> bool:
    from pptx.util import Inches

    if not image_path.is_file() or image_path.is_symlink():
        return False
    picture = slide.shapes.add_picture(str(image_path), Inches(7.35), Inches(1.72), width=Inches(5.18), height=Inches(4.78))
    picture.name = "TGVisual:picture"
    colors = design["colors"]
    fonts = design["fonts"]
    _add_rect(slide, 0.80, 1.72, 6.22, 4.78, colors["surface"], line=colors["line"], name="TGVisual:picture-caption-card")
    text = "\n\n".join(items[:4]) or "图像证据"
    _add_text(slide, text, 1.12, 2.05, 5.58, 4.05, font=fonts["body"], size=17, color=colors["ink"], valign="middle", name="TGText:picture-caption")
    return True


def _add_notes(slide: Any, notes: str) -> None:
    if not str(notes or "").strip():
        return
    try:
        slide.notes_slide.notes_text_frame.text = str(notes)
    except Exception:
        return


def build_presentation(
    resolve: Callable[..., Path],
    args: Dict[str, Any],
    *,
    content: str = "",
    source_path: Path | None = None,
) -> Tuple[Any, Dict[str, Any]]:
    """Create a widescreen, placeholder-free presentation and compact evidence."""
    from pptx import Presentation
    from pptx.util import Inches

    design, design_source = load_design_spec(resolve, args, source_path)
    specs = parse_markdown_slides(content) if str(content or "").strip() else structured_slides(args)
    if not specs:
        raise ValueError("pptx.create received no valid slide content")
    source_slide_count = len(specs)
    maximum = int(design["layout"].get("max_items_per_slide", 6))
    specs = _expand_dense_specs(specs, maximum)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = str(args.get("title") or specs[0].get("title") or "")
    prs.core_properties.subject = "Tiangong design-aware executive presentation"
    prs.core_properties.author = "Tiangong Omni Body"
    prs.core_properties.keywords = "tiangong, executive, design-system, widescreen"
    blank_layout = prs.slide_layouts[6]
    for index, spec in enumerate(specs):
        slide = prs.slides.add_slide(blank_layout)
        if spec.get("is_cover"):
            _render_cover(slide, spec, design)
        else:
            _render_header(slide, str(spec.get("title") or ""), index + 1, design)
            rendered = False
            table = spec.get("table")
            chart = spec.get("chart")
            image = spec.get("image")
            if isinstance(chart, dict):
                rendered = _render_chart(slide, chart, design)
            if not rendered and isinstance(table, dict):
                rendered = _render_table(slide, table, design)
            if not rendered and isinstance(image, str) and image.strip():
                rendered = _render_picture(slide, resolve(image, must_exist=True), list(spec.get("items") or []), design)
            if not rendered:
                items = [str(item) for item in spec.get("items") or []]
                title = str(spec.get("title") or "")
                if 2 <= len(items) <= 5 and any(word in title for word in ("路径", "阶段", "流程", "步骤", "计划", "实施", "路线")):
                    _render_timeline(slide, items, design)
                else:
                    _render_cards(slide, items, design)
            _render_footer(slide, index + 1, design)
        _add_notes(slide, str(spec.get("notes") or ""))

    metadata = {
        "design_schema": design.get("schema"),
        "template_id": design.get("template_id"),
        "style_name": design.get("style_name"),
        "design_source": design_source,
        "aspect_ratio": "16:9",
        "source_slide_count": source_slide_count,
        "generated_slide_count": len(prs.slides),
        "placeholder_free": True,
        "font_heading": design["fonts"]["heading"],
        "font_body": design["fonts"]["body"],
    }
    return prs, metadata


def inspect_presentation(path: Path) -> Dict[str, Any]:
    """Return semantic and visual evidence used by pptx.read and PPT QC."""
    from pptx import Presentation

    prs = Presentation(str(path))
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    ratio = round(width / height, 4) if height else 0.0
    slides: List[Dict[str, Any]] = []
    inspection_issues: List[Dict[str, Any]] = []
    font_names: set[str] = set()
    total_placeholders = total_shapes = native_visuals = designed_visuals = 0
    for index, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        semantic_texts: List[str] = []
        title = ""
        slide_placeholders = slide_native = slide_designed = slide_decorative = 0
        for shape in slide.shapes:
            total_shapes += 1
            name = str(getattr(shape, "name", "") or "")
            if bool(getattr(shape, "is_placeholder", False)):
                slide_placeholders += 1
                total_placeholders += 1
            has_chart = bool(getattr(shape, "has_chart", False))
            has_table = bool(getattr(shape, "has_table", False))
            is_picture = int(getattr(shape, "shape_type", 0) or 0) == 13
            if has_chart or has_table or is_picture:
                slide_native += 1
                native_visuals += 1
            if name.startswith("TGVisual:"):
                slide_designed += 1
                designed_visuals += 1
            elif name.startswith("TGDecor:"):
                slide_decorative += 1
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                texts.append(text)
                if not name.startswith(("TGText:footer", "TGText:page-kicker", "TGText:card-index", "TGText:timeline-step", "TGText:brand")):
                    semantic_texts.append(text)
                if name == "TGTitle":
                    title = text
            if bool(getattr(shape, "has_text_frame", False)):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            font_names.add(str(run.font.name))
        if not title and texts:
            title = texts[0]
        meaningful = slide_native + slide_designed
        slide_issues: List[Dict[str, Any]] = []
        content_text_joined = "\n".join(semantic_texts)
        # MM-P1-2 / MM-P2-GROUP-2: placeholders and pure decoration are not
        # valid expression.  A slide with no semantic content must never look
        # "visually rich" to the QC.
        if not content_text_joined.strip():
            slide_issues.append({
                "code": "slide_content_missing",
                "severity": "critical",
                "message": f"slide {index} has no semantic content",
                "repair": "add real heading/body content to the slide before declaring the deck complete",
            })
        elif "本页尚无可呈现内容" in content_text_joined:
            slide_issues.append({
                "code": "slide_placeholder_text",
                "severity": "critical",
                "message": f"slide {index} contains placeholder text",
                "repair": "replace the placeholder with real content and regenerate the deck",
            })
        if not content_text_joined.strip() and slide_designed > 0 and slide_native == 0:
            slide_issues.append({
                "code": "decor_only_slide",
                "severity": "high",
                "message": f"slide {index} has only decorative/designed shapes without semantic content",
                "repair": "add semantic text or a native visual to the slide",
            })
        inspection_issues.extend(slide_issues)
        slides.append(
            {
                "index": index,
                "title": title,
                "text": "\n".join(texts),
                "content_text": "\n".join(semantic_texts),
                "shape_count": len(slide.shapes),
                "placeholder_count": slide_placeholders,
                "native_visual_count": slide_native,
                "designed_visual_count": slide_designed,
                "decorative_shape_count": slide_decorative,
                "visual_count": meaningful,
                "issues": slide_issues,
            }
        )
    return {
        "schema": "tiangong.v3.ppt_inspection.v1",
        "slides": slides,
        "slide_count": len(slides),
        "slide_width": width,
        "slide_height": height,
        "aspect_ratio": ratio,
        "is_widescreen": 1.70 <= ratio <= 1.82,
        "total_shapes": total_shapes,
        "placeholder_count": total_placeholders,
        "native_visual_count": native_visuals,
        "designed_visual_count": designed_visuals,
        "font_names": sorted(font_names),
        "has_explicit_fonts": bool(font_names),
        "generator": str(prs.core_properties.author or ""),
        "design_keywords": str(prs.core_properties.keywords or ""),
        "issues": inspection_issues,
        "critical_issue_count": sum(1 for issue in inspection_issues if issue.get("severity") == "critical"),
    }
