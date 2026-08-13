"""
天工造物 v3：起源 — 肌肉层
执行L4适配器 + 结果原样回传
"""
from __future__ import annotations

import base64
import errno
import hashlib
import html
import http.client
import ipaddress
import json
import mimetypes
import os
import re
import shutil
import socket
import ssl
import subprocess
import sys
import threading
import time
import urllib.error
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta
from email.utils import parsedate_to_datetime
from io import BytesIO
from pathlib import Path
from typing import Any

from .guge_ceng import GongjuYingshe
from .omni_grant_client import issue_omni_grant
from ..json_guards import error_payload, loads_json_object
from ..run_context import current_run_context


_SEARCH_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
    "Accept": "application/rss+xml, text/xml, text/html, */*",
    "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
}

_DYNAMIC_TOOL_IMPORT_LOCK = threading.RLock()

_NEWS_QUERY_MARKERS = (
    "新闻",
    "最新",
    "今天",
    "今日",
    "昨天",
    "刚刚",
    "实时",
    "快讯",
    "近况",
    "发布",
    "宣布",
    "发生",
    "进展",
    "热点",
    "today",
    "latest",
    "news",
    "breaking",
    "current",
)

TERMINAL_DANGEROUS_COMMAND_PATTERNS = (
    r"\brm\s+-rf\b",
    r"\brmdir\s+/s\b",
    r"\brd\s+/s\b",
    r"\bdel\s+/[fsq]\b",
    r"\bformat\s+[a-z]:",
    r"\b(?:powershell|pwsh)(?:\.exe)?\s+-enc(?:odedcommand)?\b",
    r"\b(?:curl|wget|iwr|irm|invoke-webrequest|invoke-restmethod)\b.*\|\s*(?:sh|bash|powershell|pwsh|iex|invoke-expression)\b",
    r"\|\s*(?:iex|invoke-expression)\b",
    r"\bset-executionpolicy\b",
    r"\breg\s+(?:delete|add)\b",
)

_CHINA_NEWS_RSS_FEEDS = (
    ("chinanews_scroll", "https://www.chinanews.com.cn/rss/scroll-news.xml"),
    ("chinanews_china", "https://www.chinanews.com.cn/rss/china.xml"),
    ("xinhua_politics", "http://www.xinhuanet.com/politics/news_politics.xml"),
    ("xinhua_world", "http://www.xinhuanet.com/world/news_world.xml"),
    ("xinhua_finance", "http://www.xinhuanet.com/fortune/news_fortune.xml"),
    ("xinhua_tech", "http://www.xinhuanet.com/tech/news_tech.xml"),
    ("xinhua_local", "http://www.xinhuanet.com/local/news_local.xml"),
    ("xinhua_legal", "http://www.xinhuanet.com/legal/news_legal.xml"),
)

# 国际新闻 RSS 源
_GLOBAL_NEWS_RSS_FEEDS = (
    ("reuters_world", "https://feeds.reuters.com/reuters/worldNews"),
    ("reuters_top", "https://feeds.reuters.com/reuters/topNews"),
    ("bbc_world", "https://feeds.bbci.co.uk/news/world/rss.xml"),
    ("bbc_asia", "https://feeds.bbci.co.uk/news/world/asia/rss.xml"),
    ("ap_top", "https://rss.app/feeds/3AOgGkAV4oH8SMGP.xml"),  # AP News mirror
    ("npr_world", "https://feeds.npr.org/1004/rss.xml"),
    ("aljazeera", "https://www.aljazeera.com/xml/rss/all.xml"),
    ("cnn_world", "https://rss.app/feeds/BwlRC8u8TeFcu3l3.xml"),  # CNN mirror
)

_NEWS_STOPWORDS = (
    "今天", "今日", "昨天", "新闻", "最新", "实时", "快讯", "热点", "国内", "国际",
    "一下", "帮我", "查询", "搜索", "看看", "current", "today", "latest", "news",
    "breaking", "the", "and", "with",
)

_SEARCH_INTENT_STOPWORDS = {
    "about", "current", "find", "information", "latest", "news", "official",
    "please", "search", "today", "官网", "官方网站", "帮我", "今天", "今日",
    "信息", "搜索", "新闻", "最新", "查询", "看看", "相关",
}

_SEARCH_AD_MARKERS = (
    "sponsored", "新用户送", "注册送", "立即购买", "立即领取", "免费领取",
    "限时优惠", "广告", "推广", "赞助",
)

_WEB_MIN_MEANINGFUL_CHARS = 80


class _WebUrlSecurityError(ValueError):
    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code


def _looks_like_web_url(value: str) -> bool:
    raw = str(value or "").strip()
    if not raw or re.search(r"\s", raw):
        return False
    if raw.startswith("//"):
        raw = raw[2:]
    authority = re.split(r"[/?#]", raw, maxsplit=1)[0]
    if not authority:
        return False
    if "@" in authority:
        authority = authority.rsplit("@", 1)[-1]
    if authority.startswith("["):
        return bool(re.match(r"^\[[0-9a-fA-F:.%]+\](?::\d+)?$", authority))
    host = authority
    if authority.count(":") == 1:
        host, port = authority.rsplit(":", 1)
        if port and not port.isdigit():
            return False
    host = host.rstrip(".")
    if not host:
        return False
    try:
        ipaddress.ip_address(host)
        return True
    except ValueError:
        pass
    if host.lower() == "localhost":
        return True
    return bool(
        "." in host
        and all(re.fullmatch(r"[A-Za-z0-9-]{1,63}", label) and not label.startswith("-") and not label.endswith("-")
                for label in host.split("."))
    )


def _is_public_ip(value: str) -> bool:
    try:
        address = ipaddress.ip_address(str(value or "").split("%", 1)[0])
    except ValueError:
        return False
    return address.is_global


def _resolve_public_socket_addresses(host: str, port: int) -> list[tuple[int, int, int, tuple[Any, ...]]]:
    try:
        infos = socket.getaddrinfo(
            host,
            port,
            family=socket.AF_UNSPEC,
            type=socket.SOCK_STREAM,
            proto=socket.IPPROTO_TCP,
        )
    except socket.gaierror as exc:
        raise _WebUrlSecurityError("url_resolution_failed", f"网页主机无法解析：{host}") from exc

    blocked: list[str] = []
    resolved: list[tuple[int, int, int, tuple[Any, ...]]] = []
    seen: set[tuple[int, int, int, tuple[Any, ...]]] = set()
    for family, socktype, proto, _canonname, sockaddr in infos:
        if not sockaddr:
            continue
        address = str(sockaddr[0]).split("%", 1)[0]
        if not _is_public_ip(address):
            blocked.append(address)
            continue
        item = (family, socktype or socket.SOCK_STREAM, proto or socket.IPPROTO_TCP, tuple(sockaddr))
        if item not in seen:
            seen.add(item)
            resolved.append(item)
    if blocked:
        raise _WebUrlSecurityError("unsafe_url", f"网页主机解析到非公网地址：{', '.join(sorted(set(blocked))[:3])}")
    if not resolved:
        raise _WebUrlSecurityError("url_resolution_failed", f"网页主机没有可用公网地址：{host}")
    return resolved


def _connect_pinned_public_address(
    host: str,
    port: int,
    timeout: Any,
    source_address: tuple[str, int] | None,
):
    addresses = _resolve_public_socket_addresses(host, port)
    last_error: OSError | None = None
    for family, socktype, proto, sockaddr in addresses:
        sock = socket.socket(family, socktype, proto)
        try:
            if timeout is not socket._GLOBAL_DEFAULT_TIMEOUT:
                sock.settimeout(timeout)
            if source_address:
                sock.bind(source_address)
            sock.connect(sockaddr)
            return sock
        except OSError as exc:
            last_error = exc
            sock.close()
    if last_error is not None:
        raise last_error
    raise OSError(f"无法连接已验证的公网地址：{host}:{port}")


def _set_tcp_nodelay(sock: Any) -> None:
    try:
        sock.setsockopt(socket.IPPROTO_TCP, socket.TCP_NODELAY, 1)
    except OSError as exc:
        if exc.errno != errno.ENOPROTOOPT:
            raise


class _PinnedHTTPConnection(http.client.HTTPConnection):
    def connect(self):
        if self._tunnel_host:
            raise _WebUrlSecurityError("unsafe_proxy", "网页读取禁止代理隧道")
        sys.audit("http.client.connect", self, self.host, self.port)
        self.sock = _connect_pinned_public_address(self.host, self.port, self.timeout, self.source_address)
        _set_tcp_nodelay(self.sock)


class _PinnedHTTPSConnection(http.client.HTTPSConnection):
    def connect(self):
        if self._tunnel_host:
            raise _WebUrlSecurityError("unsafe_proxy", "网页读取禁止代理隧道")
        server_hostname = self.host
        sys.audit("http.client.connect", self, self.host, self.port)
        self.sock = _connect_pinned_public_address(self.host, self.port, self.timeout, self.source_address)
        _set_tcp_nodelay(self.sock)
        try:
            self.sock = self._context.wrap_socket(self.sock, server_hostname=server_hostname)
        except Exception:
            self.sock.close()
            self.sock = None
            raise


class _PinnedHTTPHandler(urllib.request.HTTPHandler):
    def http_open(self, req):
        return self.do_open(_PinnedHTTPConnection, req)


class _PinnedHTTPSHandler(urllib.request.HTTPSHandler):
    def https_open(self, req):
        return self.do_open(_PinnedHTTPSConnection, req, context=self._context)


def _normalise_public_web_url(value: str, *, resolve: bool = True) -> str:
    raw = str(value or "").strip()
    if not raw:
        raise _WebUrlSecurityError("invalid_url", "URL 为空")
    if raw.startswith("//"):
        raw = "https:" + raw
    elif "://" not in raw and _looks_like_web_url(raw):
        raw = "https://" + raw

    parsed = urllib.parse.urlsplit(raw)
    scheme = parsed.scheme.lower()
    if scheme not in {"http", "https"}:
        raise _WebUrlSecurityError("unsupported_url_scheme", "网页读取仅允许 http/https URL")
    if parsed.username is not None or parsed.password is not None:
        raise _WebUrlSecurityError("unsafe_url", "网页 URL 不允许包含用户凭据")

    host = str(parsed.hostname or "").rstrip(".").lower()
    if not host:
        raise _WebUrlSecurityError("invalid_url", "网页 URL 缺少主机名")
    if "%" in host or host == "localhost" or host.endswith((".localhost", ".local", ".lan", ".home", ".internal")):
        raise _WebUrlSecurityError("unsafe_url", "网页 URL 指向本机或内部主机")
    try:
        port = parsed.port
    except ValueError as exc:
        raise _WebUrlSecurityError("invalid_url", f"网页 URL 端口无效：{exc}") from exc

    try:
        ascii_host = host.encode("idna").decode("ascii")
    except UnicodeError as exc:
        raise _WebUrlSecurityError("invalid_url", "网页 URL 主机名无效") from exc

    try:
        literal_ip = ipaddress.ip_address(ascii_host)
    except ValueError:
        literal_ip = None
    if literal_ip is not None and not literal_ip.is_global:
        raise _WebUrlSecurityError("unsafe_url", f"网页 URL 指向非公网地址：{literal_ip}")

    if resolve:
        _resolve_public_socket_addresses(ascii_host, port or (443 if scheme == "https" else 80))

    netloc = f"[{ascii_host}]" if ":" in ascii_host else ascii_host
    if port is not None:
        netloc += f":{port}"
    return urllib.parse.urlunsplit((scheme, netloc, parsed.path or "/", parsed.query, parsed.fragment))


class _PublicWebRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        safe_url = _normalise_public_web_url(newurl, resolve=True)
        return super().redirect_request(req, fp, code, msg, headers, safe_url)


def _public_web_opener():
    context = ssl.create_default_context()
    try:
        context.set_alpn_protocols(["http/1.1"])
    except NotImplementedError:
        pass
    return urllib.request.build_opener(
        urllib.request.ProxyHandler({}),
        _PinnedHTTPHandler(),
        _PinnedHTTPSHandler(context=context),
        _PublicWebRedirectHandler(),
    )


def _is_tls_certificate_error(exc: BaseException) -> bool:
    pending: list[BaseException] = [exc]
    seen: set[int] = set()
    while pending:
        current = pending.pop()
        if id(current) in seen:
            continue
        seen.add(id(current))
        if isinstance(current, ssl.SSLCertVerificationError):
            return True
        if "certificate verify failed" in str(current).casefold():
            return True
        for nested in (getattr(current, "reason", None), current.__cause__, current.__context__):
            if isinstance(nested, BaseException):
                pending.append(nested)
    return False


def _web_fetch_error_payload(url: str, final_url: str, exc: BaseException) -> dict[str, Any]:
    certificate_error = _is_tls_certificate_error(exc)
    return {
        "ok": False,
        "zhuangtai": "cuowu",
        "url": str(url or ""),
        "final_url": str(final_url or url or ""),
        "recoverable": not certificate_error,
        "error_code": "tls_certificate_error" if certificate_error else "web_fetch_failed",
        "cuowu": str(exc)[:300],
    }


def _web_url_error_payload(url: str, exc: _WebUrlSecurityError) -> dict[str, Any]:
    return {
        "ok": False,
        "zhuangtai": "cuowu",
        "url": str(url or ""),
        "recoverable": False,
        "error_code": exc.code,
        "cuowu": str(exc),
    }


def _web_content_block_reason(source_text: str, plain: str) -> str:
    raw = str(source_text or "")
    text = str(plain or "")
    lowered = text.casefold()
    leading = lowered[:1600].strip()
    title_match = re.search(r"(?is)<title[^>]*>(.*?)</title>", raw)
    title = re.sub(r"\s+", " ", html.unescape(re.sub(r"(?is)<[^>]+>", " ", title_match.group(1)))).strip().casefold() if title_match else ""
    compact_title = re.sub(r"[^\w\u4e00-\u9fff]+", " ", title).strip()

    strong_markers = (
        "[stub]", "参数错误", "访问过于频繁", "操作频繁", "请在微信客户端打开",
        "环境异常", "请输入验证码", "verify you are human", "checking your browser",
        "enable javascript and cookies to continue",
    )
    for marker in strong_markers:
        if marker.casefold() in leading or marker.casefold() in title:
            return marker

    generic_titles = (
        "access denied", "attention required", "forbidden", "unauthorized", "just a moment",
        "captcha", "security check", "sign in", "log in", "login", "登录", "用户登录",
        "账号登录", "安全验证", "人机验证", "页面不存在", "404 not found",
    )
    for marker in generic_titles:
        if compact_title == marker or compact_title.startswith(marker + " "):
            return marker

    if leading.startswith(("access denied", "forbidden", "unauthorized", "just a moment", "checking your browser")):
        return leading.splitlines()[0][:80]

    raw_lower = raw.casefold()
    has_password_input = bool(re.search(r"(?is)<input[^>]+type\s*=\s*['\"]?password", raw_lower))
    has_login_form = bool(re.search(r"(?is)<form[^>]+(?:login|signin|sign-in|auth)", raw_lower))
    login_markers = ("sign in", "log in", "login", "请登录", "请先登录", "登录后查看", "账号登录", "扫码登录")
    if (has_password_input or has_login_form) and any(marker in leading for marker in login_markers):
        return "login_required"

    challenge_markers = (
        "access denied", "cloudflare", "ray id", "captcha", "security check", "安全验证",
        "人机验证", "forbidden", "unauthorized", "checking your browser",
    )
    hits = [marker for marker in challenge_markers if marker in leading or marker in raw_lower[:5000]]
    if len(hits) >= 2:
        return hits[0]

    meaningful_len = len(re.sub(r"\s+", "", text))
    if meaningful_len < _WEB_MIN_MEANINGFUL_CHARS:
        return "content_too_short"
    return ""


def _workspace_root() -> Path:
    forced = os.environ.get("TIANGONG_FORCE_WORKSPACE_ROOT") or os.environ.get("TIANGONG_TOOL_TEST_WORKSPACE_ROOT")
    if forced:
        root = Path(forced).expanduser().resolve(strict=False)
        root.mkdir(parents=True, exist_ok=True)
        return root
    try:
        from ..workspace_settings import duqu_workspace_root

        root = duqu_workspace_root()
    except Exception:
        raw = (
            os.environ.get("TIANGONG_DESKTOP_WORKSPACE_ROOT")
            or os.environ.get("TIANGONG_WORKSPACE_ROOT")
            or os.getcwd()
        )
        root = Path(raw).expanduser().resolve(strict=False)
    root.mkdir(parents=True, exist_ok=True)
    return root


def _workspace_path(value: Any = "", *, default: str = ".") -> Path:
    raw = str(value if value not in (None, "") else default).strip()
    if not raw:
        raw = default
    try:
        from ..path_resolver import resolve_path_text

        resolved = resolve_path_text(raw, base=str(_workspace_root()))
        if resolved.get("ok") and str(resolved.get("scope") or "") != "url":
            resolved_path = str(resolved.get("resolved_path") or "").strip()
            if resolved_path:
                return Path(resolved_path).expanduser().resolve(strict=False)
    except Exception:
        pass
    path = Path(raw).expanduser()
    if not path.is_absolute():
        path = _workspace_root() / path
    return path.resolve(strict=False)


def _office_output_path(value: Any, suffix: str, folder: str) -> Path:
    raw = str(value or "").strip()
    if raw:
        path = _workspace_path(raw)
    else:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = _workspace_root() / "generated" / folder / f"{folder}-{stamp}{suffix}"
    if path.suffix.lower() != suffix:
        path = path.with_suffix(suffix)
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _tool_success(path: Path, file_type: str, **extra: Any) -> dict:
    result = {
        "ok": True,
        "zhuangtai": "wancheng",
        "effect": "write",
        "leixing": file_type,
        "lujing": str(path),
        "wenjian_daxiao": path.stat().st_size,
        "sha256": _file_sha256(path),
        "updated_paths": [str(path)],
        "readback": {"ok": True, "bytes": path.stat().st_size},
    }
    result.update(extra)
    return result


def _backend_root() -> Path:
    return Path(__file__).resolve().parents[2]


def _activate_bundled_runtime() -> None:
    root = _backend_root()
    vendor = root / "_vendor"
    bin_dir = root / "bin"
    if vendor.exists() and str(vendor) not in sys.path:
        sys.path.insert(0, str(vendor))
    if bin_dir.exists():
        current_path = os.environ.get("PATH", "")
        bin_text = str(bin_dir)
        if bin_text.lower() not in [item.lower() for item in current_path.split(os.pathsep) if item]:
            os.environ["PATH"] = bin_text + os.pathsep + current_path


def _bundled_binary(name: str) -> str:
    clean = str(name or "").strip()
    candidates = []
    if clean:
        candidates.append(_backend_root() / "bin" / clean)
        if os.name == "nt" and not clean.lower().endswith(".exe"):
            candidates.append(_backend_root() / "bin" / f"{clean}.exe")
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    found = shutil.which(clean)
    return found or clean


def _novel_tool_script() -> Path:
    return Path(__file__).resolve().parents[1] / "bundled_skills" / "novel-creation" / "scripts" / "novel_tool.py"


def _load_novel_tool_module() -> Any:
    script = _novel_tool_script()
    if not script.exists():
        raise FileNotFoundError(f"novel_tool_not_found:{script}")
    import importlib.util

    module_name = "_tiangong_v3_bundled_novel_tool"
    with _DYNAMIC_TOOL_IMPORT_LOCK:
        loaded = sys.modules.get(module_name)
        if loaded is not None:
            return loaded
        spec = importlib.util.spec_from_file_location(module_name, str(script))
        if spec is None or spec.loader is None:
            raise RuntimeError(f"novel_tool_import_failed:{script}")
        module = importlib.util.module_from_spec(spec)
        sys.modules[module_name] = module
        previous_dont_write_bytecode = sys.dont_write_bytecode
        sys.dont_write_bytecode = True
        try:
            spec.loader.exec_module(module)
        except Exception:
            sys.modules.pop(module_name, None)
            raise
        finally:
            sys.dont_write_bytecode = previous_dont_write_bytecode
        return module


def _omni_body_skill_root() -> Path:
    forced = os.environ.get("TIANGONG_OMNI_BODY_ROOT")
    candidates = []
    if forced:
        candidates.append(Path(forced).expanduser())
    candidates.extend([
        _backend_root() / "omni_body_skill",
        Path(__file__).resolve().parents[1] / "bundled_skills" / "omni_body_skill",
    ])
    if str(os.environ.get("TIANGONG_ALLOW_USER_SKILL_OVERRIDE") or "").strip().lower() in {"1", "true", "yes", "on"}:
        candidates.insert(1 if forced else 0, Path.home() / ".tiangong" / "v3" / "omni_body_skill")
    for candidate in candidates:
        root = candidate.resolve(strict=False)
        if (root / "api" / "v1" / "v3" / "tools" / "omni_body.py").exists() and (root / "tools" / "omni_body_tool.py").exists():
            return root
    raise FileNotFoundError("omni_body_skill_not_found")


def _load_omni_body_module() -> Any:
    root = _omni_body_skill_root()
    script = root / "api" / "v1" / "v3" / "tools" / "omni_body.py"
    import importlib.util

    _activate_bundled_runtime()
    os.environ["TIANGONG_OMNI_BODY_ROOT"] = str(root)
    module_name = "_tiangong_v3_omni_body_tool"
    with _DYNAMIC_TOOL_IMPORT_LOCK:
        loaded = sys.modules.get(module_name)
        if loaded is not None and (
            callable(getattr(loaded, "run_omni_body", None))
            or callable(getattr(loaded, "run", None))
        ):
            return loaded
        sys.modules.pop(module_name, None)
        spec = importlib.util.spec_from_file_location(module_name, str(script))
        if spec is None or spec.loader is None:
            raise RuntimeError(f"omni_body_import_failed:{script}")
        module = importlib.util.module_from_spec(spec)
        sys.modules[module_name] = module
        previous_dont_write_bytecode = sys.dont_write_bytecode
        sys.dont_write_bytecode = True
        try:
            spec.loader.exec_module(module)
        except Exception:
            sys.modules.pop(module_name, None)
            raise
        finally:
            sys.dont_write_bytecode = previous_dont_write_bytecode
        return module


def _run_omni_body_tool(canshu: dict[str, Any]) -> dict[str, Any]:
    module = _load_omni_body_module()
    runner = getattr(module, "run_omni_body", None) or getattr(module, "run", None)
    if not callable(runner):
        raise RuntimeError("omni_body_runner_missing")
    payload = dict(canshu or {})
    if not isinstance(payload.get("__capability_grant"), dict) or not isinstance(payload.get("__runtime"), dict):
        raise PermissionError("CAPABILITY_REQUIRED")
    payload["workspace"] = str(_workspace_root())
    result = runner(payload)
    if isinstance(result, dict):
        return result
    return {"ok": True, "zhuangtai": "wancheng", "gongju": "omni_body", "result": result}


def _json_from_stdout(text: str) -> dict[str, Any]:
    raw = str(text or "").strip()
    if not raw:
        return {}
    try:
        data = json.loads(raw)
        return data if isinstance(data, dict) else {"value": data}
    except Exception:
        pass
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        try:
            data = json.loads(raw[start : end + 1])
            return data if isinstance(data, dict) else {"value": data}
        except Exception:
            return {}
    return {}


_activate_bundled_runtime()


def _result_readback(path_value: Any, *, deleted: bool = False) -> dict[str, Any]:
    if not path_value:
        return {"ok": False}
    try:
        path = Path(str(path_value))
        exists = path.exists()
        payload: dict[str, Any] = {"ok": (not exists if deleted else exists), "exists": exists}
        if exists and path.is_file():
            payload["bytes"] = path.stat().st_size
        return payload
    except Exception as exc:
        return {"ok": False, "error": str(exc)[:200]}


def _normalise_tool_result(tool_name: str, effect: str, result: Any) -> dict[str, Any]:
    if not isinstance(result, dict):
        result = {"value": result}
    out = dict(result)
    state = str(out.get("zhuangtai") or out.get("status") or "").strip()
    error_text = str(out.get("cuowu") or out.get("error") or out.get("detail") or "").strip()
    success_states = {"wancheng", "yixieru", "yixiazai", "yifuzhi", "yiyidong", "yishanchu", "yicunzai", "wu_zimu"}
    failure_states = {"cuowu", "jujue", "buzhichi", "weishixian", "plan_only"}
    if "ok" not in out:
        if out.get("fanhui_ma") == 0:
            out["ok"] = True
        elif state in success_states:
            out["ok"] = True
        elif state in failure_states or error_text:
            out["ok"] = False
        elif any(key in out for key in ("neirong", "jieguo", "pipei_shu", "wenjian", "biaozhun_shuchu", "sha256")):
            out["ok"] = True
        else:
            out["ok"] = False
    if not state:
        out["zhuangtai"] = "wancheng" if out.get("ok") else "cuowu"
    if "gongju" not in out:
        out["gongju"] = tool_name
    if "effect" not in out:
        out["effect"] = effect or "unknown"
    if out.get("lujing") and "path" not in out:
        out["path"] = out.get("lujing")
    if out.get("ok") is False and error_text and "error" not in out:
        out["error"] = error_text
    if "updated_paths" not in out and out.get("effect") in {"write", "generate"}:
        paths = [out.get(key) for key in ("target", "lujing", "shuchu", "zimu_wenjian", "zhen_mulu") if out.get(key)]
        if paths and out.get("zhuangtai") != "yishanchu":
            out["updated_paths"] = [str(item) for item in paths]
    if "readback" not in out:
        if out.get("zhuangtai") == "yishanchu":
            out["readback"] = _result_readback(out.get("lujing") or out.get("target"), deleted=True)
        elif out.get("lujing") and out.get("effect") in {"read", "write", "generate"}:
            out["readback"] = _result_readback(out.get("lujing"))
        elif out.get("target") and out.get("effect") in {"write", "generate"}:
            out["readback"] = _result_readback(out.get("target"))
    if out.get("ok") is True and out.get("effect") in {"write", "generate"}:
        readback = out.get("readback")
        if isinstance(readback, dict) and readback.get("ok") is False:
            out["ok"] = False
            out["zhuangtai"] = "cuowu"
            out.setdefault("error", "WRITE_READBACK_FAILED")
            out.setdefault("cuowu", out.get("error"))
    return out


def _authority_receipt(authority: Any) -> dict[str, Any]:
    """Return auditable authority identifiers without exposing signed grants."""

    value = authority if isinstance(authority, dict) else {}
    grant = value.get("grant") if isinstance(value.get("grant"), dict) else {}
    grant_payload = grant.get("payload") if isinstance(grant.get("payload"), dict) else {}
    runtime = value.get("runtime") if isinstance(value.get("runtime"), dict) else {}
    context = current_run_context()
    return {
        "schema": "tiangong.omni.execution-receipt.v1",
        "outer_execution_ticket_id": str(context.outer_execution_ticket_id or ""),
        "execution_ticket_id": str(runtime.get("execution_ticket_id") or ""),
        "omni_grant_id": str(grant_payload.get("grant_id") or ""),
        "decision_id": str(grant_payload.get("decision_id") or ""),
        "decision_sha256": str(runtime.get("decision_sha256") or ""),
        "action_permission_sha256": str(runtime.get("action_permission_sha256") or ""),
    }


def _as_list(value: Any) -> list[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, tuple):
        return list(value)
    if isinstance(value, str):
        lines = [line.strip() for line in value.splitlines() if line.strip()]
        return lines or ([value.strip()] if value.strip() else [])
    return [value]


def _plain_text(value: Any, limit: int = 2000) -> str:
    if value is None:
        return ""
    if isinstance(value, (dict, list, tuple)):
        text = json.dumps(value, ensure_ascii=False)
    else:
        text = str(value)
    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit]


def _normalise_table_rows(rows: Any, headers: Any = None) -> tuple[list[str], list[list[Any]]]:
    items = _as_list(rows)
    header_list = [str(item) for item in _as_list(headers)]
    if not items:
        return header_list, []
    if all(isinstance(item, dict) for item in items):
        wrapped_rows = [_dict_row_values(item, header_list) for item in items]
        if any(row is not None for row in wrapped_rows):
            table_rows = [
                row if row is not None else [item.get(key, "") for key in header_list]
                for item, row in zip(items, wrapped_rows)
            ]
            return header_list, table_rows
        if not header_list:
            seen: list[str] = []
            for item in items:
                for key in item.keys():
                    key_text = str(key)
                    if key_text not in seen:
                        seen.append(key_text)
            header_list = seen
        return header_list, [[item.get(key, "") for key in header_list] for item in items]
    table_rows = []
    for item in items:
        if isinstance(item, (list, tuple)):
            table_rows.append(list(item))
        else:
            table_rows.append([item])
    return header_list, table_rows


def _dict_row_values(item: dict, headers: list[str]) -> list[Any] | None:
    for key in ("item", "row", "values", "cells"):
        if key not in item:
            continue
        value = item.get(key)
        if isinstance(value, (list, tuple)):
            return list(value)
        if isinstance(value, dict):
            if headers:
                return [value.get(header, "") for header in headers]
            return list(value.values())
    if len(item) == 1:
        value = next(iter(item.values()))
        if isinstance(value, (list, tuple)):
            return list(value)
    return None


def _normalise_table_payload(table: Any) -> tuple[list[str], list[list[Any]]]:
    if isinstance(table, dict):
        return _normalise_table_rows(table.get("rows") or table.get("data") or [], table.get("headers") or table.get("columns"))
    return _normalise_table_rows(table)


def _safe_xlsx_sheet_name(value: Any, fallback: str, used: set[str]) -> str:
    name = str(value or fallback or "Sheet").strip()
    name = re.sub(r"[\[\]\:\*\?\/\\]", "_", name)[:31].strip() or fallback or "Sheet"
    base = name[:31]
    index = 2
    while name in used:
        suffix = f"_{index}"
        name = f"{base[:31 - len(suffix)]}{suffix}"
        index += 1
    used.add(name)
    return name


def _image_metadata(path: Path) -> dict:
    try:
        from PIL import Image, ImageStat

        with Image.open(str(path)) as img:
            stat = ImageStat.Stat(img.convert("RGB").resize((1, 1)))
            avg = [int(v) for v in stat.mean]
            info = {
                "leixing": "image",
                "geshi": img.format or path.suffix.lower().lstrip("."),
                "moshi": img.mode,
                "kuandu": img.width,
                "gaodu": img.height,
                "wenjian_daxiao": path.stat().st_size,
                "mime": mimetypes.guess_type(str(path))[0] or "",
                "avg_rgb": avg,
            }
            try:
                info["exif_count"] = len(img.getexif() or {})
            except Exception:
                info["exif_count"] = 0
            return info
    except Exception as exc:
        return {
            "leixing": "image",
            "wenjian_daxiao": path.stat().st_size if path.exists() else 0,
            "cuowu": str(exc)[:200],
        }


def _image_data_url(path: Path, max_bytes: int = 10 * 1024 * 1024) -> str:
    size = path.stat().st_size
    if size > max_bytes:
        raise ValueError(f"image_too_large:{size}>{max_bytes}")
    mime = mimetypes.guess_type(str(path))[0] or "image/png"
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{data}"


def _bytes_data_url(data: bytes, mime: str) -> str:
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


def _pil_data_url(image: Any, prefer_png: bool = True) -> str:
    buf = BytesIO()
    pixel_count = int(image.size[0]) * int(image.size[1])
    use_png = prefer_png and pixel_count <= 3_500_000
    if use_png:
        image.save(buf, format="PNG", optimize=True)
        return _bytes_data_url(buf.getvalue(), "image/png")
    image.convert("RGB").save(buf, format="JPEG", quality=90, optimize=True)
    return _bytes_data_url(buf.getvalue(), "image/jpeg")


def _resize_long_side(image: Any, long_side: int) -> Any:
    w, h = image.size
    current = max(w, h)
    if current <= 0 or current == long_side:
        return image
    scale = long_side / current
    new_size = (max(1, int(w * scale)), max(1, int(h * scale)))
    try:
        resample = image.Resampling.LANCZOS
    except AttributeError:
        from PIL import Image
        resample = Image.LANCZOS
    return image.resize(new_size, resample)


def _original_image_view(path: Path) -> dict[str, str]:
    try:
        return {"label": "original", "url": _image_data_url(path)}
    except ValueError:
        from PIL import Image, ImageOps

        with Image.open(str(path)) as src:
            img = ImageOps.exif_transpose(src).convert("RGB")
        resized = _resize_long_side(img, 2600)
        return {"label": "original_resized", "url": _pil_data_url(resized, prefer_png=False)}


def _enhanced_image_views(path: Path, purpose: str = "general") -> list[dict[str, str]]:
    views = [_original_image_view(path)]
    if purpose not in {"ocr", "table"}:
        return views
    try:
        from PIL import Image, ImageEnhance, ImageOps

        with Image.open(str(path)) as src:
            img = ImageOps.exif_transpose(src).convert("RGB")
        w, h = img.size
        long_side = max(w, h)
        target_long = 2200 if long_side < 1600 else min(2600, long_side)
        enhanced = _resize_long_side(img, target_long)
        gray = ImageOps.grayscale(enhanced)
        gray = ImageOps.autocontrast(gray, cutoff=1)
        tuned = gray.convert("RGB")
        tuned = ImageEnhance.Contrast(tuned).enhance(1.65)
        tuned = ImageEnhance.Sharpness(tuned).enhance(1.55)
        views.append({"label": "enhanced_full", "url": _pil_data_url(tuned, prefer_png=True)})

        if max(w, h) >= 1200:
            overlap = 0.08
            boxes = [
                (0, 0, int(w * (0.5 + overlap)), int(h * (0.5 + overlap))),
                (int(w * (0.5 - overlap)), 0, w, int(h * (0.5 + overlap))),
                (0, int(h * (0.5 - overlap)), int(w * (0.5 + overlap)), h),
                (int(w * (0.5 - overlap)), int(h * (0.5 - overlap)), w, h),
            ]
            for idx, box in enumerate(boxes, start=1):
                crop = img.crop(box)
                crop = _resize_long_side(crop, 1700)
                crop_gray = ImageOps.autocontrast(ImageOps.grayscale(crop), cutoff=1).convert("RGB")
                crop_gray = ImageEnhance.Contrast(crop_gray).enhance(1.5)
                crop_gray = ImageEnhance.Sharpness(crop_gray).enhance(1.4)
                views.append({"label": f"enhanced_tile_{idx}", "url": _pil_data_url(crop_gray, prefer_png=False)})
    except Exception as exc:
        views.append({"label": "enhance_error", "error": str(exc)[:200]})
    return [view for view in views if view.get("url")]


def _vision_settings() -> dict:
    from ..peizhi import MOREN_PROVIDER, duqu_api_miyao, duqu_moren_provider, duqu_provider_base_url

    provider = os.environ.get("TIANGONG_VISION_PROVIDER") or duqu_moren_provider(MOREN_PROVIDER)
    provider = str(provider or "").strip()
    base_url = os.environ.get("TIANGONG_VISION_BASE_URL") or duqu_provider_base_url(provider) or ""
    api_key = os.environ.get("TIANGONG_VISION_API_KEY") or duqu_api_miyao(provider) or ""
    model = os.environ.get("TIANGONG_VISION_MODEL") or _default_vision_model(provider)
    return {
        "provider": provider,
        "base_url": str(base_url or "").rstrip("/"),
        "api_key": api_key,
        "model": model,
    }


def _image_generation_settings() -> dict:
    from ..peizhi import (
        MOREN_PROVIDER,
        duqu_api_miyao,
        duqu_configured_model_ming,
        duqu_model_ming,
        duqu_moren_provider,
        duqu_provider_base_url,
    )

    image_override = any(os.environ.get(key) for key in (
        "TIANGONG_IMAGE_PROVIDER",
        "TIANGONG_IMAGE_BASE_URL",
        "TIANGONG_IMAGE_MODEL",
        "TIANGONG_IMAGE_API_KEY",
    ))
    provider = os.environ.get("TIANGONG_IMAGE_PROVIDER") or duqu_moren_provider(MOREN_PROVIDER)
    provider = str(provider or "").strip()
    base_url = os.environ.get("TIANGONG_IMAGE_BASE_URL") or duqu_provider_base_url(provider) or ""
    api_key = os.environ.get("TIANGONG_IMAGE_API_KEY") or duqu_api_miyao(provider) or ""
    configured_model = duqu_configured_model_ming(provider)
    default_image_model = _default_image_generation_model(provider)
    model = (
        os.environ.get("TIANGONG_IMAGE_MODEL")
        or default_image_model
        or configured_model
        or duqu_model_ming(provider)
    )
    return {
        "provider": provider,
        "base_url": str(base_url or "").rstrip("/"),
        "api_key": api_key,
        "model": str(model or "").strip(),
        "config_source": "image_override" if image_override else "current_llm_builtin",
        "model_source": (
            "env:TIANGONG_IMAGE_MODEL"
            if os.environ.get("TIANGONG_IMAGE_MODEL")
            else "provider_image_default"
            if default_image_model
            else "configured_llm_model"
            if configured_model
            else "provider_model_default"
        ),
    }


def _is_minimax_provider(provider: str) -> bool:
    lowered = str(provider or "").strip().lower().replace("-", "_")
    return lowered in {"minimax", "minimax_m3"} or "minimax" in lowered


def _default_image_generation_model(provider: str) -> str:
    provider = str(provider or "").strip().lower()
    defaults = {
        "openai": "gpt-image-1",
        "gpt_5_6": "gpt-image-1",
        "minimax": "image-01",
        "minimax_m3": "image-01",
    }
    return defaults.get(provider, "")


def _trim_known_generation_endpoint(value: str) -> str:
    lowered = value.lower()
    for suffix in (
        "/chat/completions",
        "/v1/chat/completions",
        "/images/generations",
        "/v1/images/generations",
        "/image_generation",
        "/v1/image_generation",
    ):
        if lowered.endswith(suffix):
            return value[: -len(suffix)].rstrip("/")
    return value


def _image_generation_endpoint(base_url: str, provider: str = "") -> str:
    value = str(base_url or "").strip().rstrip("/")
    if not value:
        return ""
    value = _trim_known_generation_endpoint(value)
    if _is_minimax_provider(provider):
        if value.endswith("/v1"):
            return f"{value}/image_generation"
        return f"{value}/v1/image_generation"
    if value.endswith("/images/generations"):
        return value
    return f"{value}/images/generations"


def _minimax_aspect_ratio(size: str) -> str:
    raw = str(size or "").strip().lower().replace(" ", "")
    allowed = {"1:1", "16:9", "4:3", "3:2", "2:3", "3:4", "9:16", "21:9"}
    if raw in allowed:
        return raw
    exact = {
        "1024x1024": "1:1",
        "1280x720": "16:9",
        "1152x864": "4:3",
        "1248x832": "3:2",
        "832x1248": "2:3",
        "864x1152": "3:4",
        "720x1280": "9:16",
        "1344x576": "21:9",
    }
    if raw in exact:
        return exact[raw]
    match = re.match(r"^(\d{2,5})x(\d{2,5})$", raw)
    if not match:
        return "1:1"
    width = max(1, int(match.group(1)))
    height = max(1, int(match.group(2)))
    ratio = width / height
    candidates = {
        "1:1": 1.0,
        "16:9": 16 / 9,
        "4:3": 4 / 3,
        "3:2": 3 / 2,
        "2:3": 2 / 3,
        "3:4": 3 / 4,
        "9:16": 9 / 16,
        "21:9": 21 / 9,
    }
    return min(candidates, key=lambda key: abs(candidates[key] - ratio))


def _safe_image_output_path(output_name: str = "") -> Path:
    from ..workspace_settings import duqu_workspace_root

    root = duqu_workspace_root()
    raw = str(output_name or "").strip()
    if raw:
        p = Path(raw).expanduser()
        if not p.is_absolute():
            p = root / "generated" / "images" / p
    else:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = root / "generated" / "images" / f"image-{stamp}.png"
    if p.suffix.lower() not in {".png", ".jpg", ".jpeg", ".webp"}:
        p = p.with_suffix(".png")
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _write_generated_image(data: dict, output_path: Path) -> tuple[Path, str]:
    candidates = data.get("data")
    if isinstance(candidates, dict):
        candidates = [candidates]
    elif not isinstance(candidates, list):
        candidates = data.get("images")
        if isinstance(candidates, dict):
            candidates = [candidates]
    if not isinstance(candidates, list):
        candidates = [data]
    for item in candidates:
        if not isinstance(item, dict):
            continue
        image_urls = item.get("image_urls")
        if isinstance(image_urls, list):
            for image_url in image_urls:
                url = str(image_url or "").strip()
                if url:
                    req = urllib.request.Request(url, headers={"User-Agent": "TiangongImage/1.0"})
                    with urllib.request.urlopen(req, timeout=120) as resp:
                        payload = resp.read()
                        content_type = resp.headers.get("Content-Type", "")
                    if output_path.suffix.lower() == ".png" and "jpeg" in content_type.lower():
                        output_path = output_path.with_suffix(".jpg")
                    elif output_path.suffix.lower() == ".png" and "webp" in content_type.lower():
                        output_path = output_path.with_suffix(".webp")
                    output_path.write_bytes(payload)
                    return output_path, "image_urls"
        b64 = str(item.get("b64_json") or item.get("base64") or item.get("image_base64") or "").strip()
        if b64:
            if "," in b64 and b64.lower().startswith("data:image/"):
                b64 = b64.split(",", 1)[1]
            output_path.write_bytes(base64.b64decode(b64))
            return output_path, "b64_json"
        url = str(item.get("url") or item.get("image_url") or "").strip()
        if url:
            req = urllib.request.Request(url, headers={"User-Agent": "TiangongImage/1.0"})
            with urllib.request.urlopen(req, timeout=120) as resp:
                payload = resp.read()
                content_type = resp.headers.get("Content-Type", "")
            if output_path.suffix.lower() == ".png" and "jpeg" in content_type.lower():
                output_path = output_path.with_suffix(".jpg")
            elif output_path.suffix.lower() == ".png" and "webp" in content_type.lower():
                output_path = output_path.with_suffix(".webp")
            output_path.write_bytes(payload)
            return output_path, "url"
    raise ValueError("image_response_missing_b64_or_url")


def _image_generation_call(prompt: str, size: str = "1024x1024", style: str = "", output_name: str = "") -> dict:
    clean_prompt = str(prompt or "").strip()
    if not clean_prompt:
        return {"zhuangtai": "cuowu", "cuowu": "prompt 为空"}
    settings = _image_generation_settings()
    if not settings.get("provider") or not settings.get("base_url") or not settings.get("api_key") or not settings.get("model"):
        return {
            "zhuangtai": "weishixian",
            "leixing": "image_generate",
            "vision_state": "unconfigured",
            "cuowu": "未配置图片生成模型的 provider/base_url/model/api_key",
            "provider": settings.get("provider", ""),
            "model": settings.get("model", ""),
            "config_source": settings.get("config_source", ""),
            "model_source": settings.get("model_source", ""),
            "tishi": "在模型设置里填写生图模型配置；MiniMax 走 /v1/image_generation，OpenAI-compatible 走 /images/generations。也可设置 TIANGONG_IMAGE_PROVIDER、TIANGONG_IMAGE_BASE_URL、TIANGONG_IMAGE_MODEL、TIANGONG_IMAGE_API_KEY。",
        }
    output_path = _safe_image_output_path(output_name)
    endpoint = _image_generation_endpoint(settings["base_url"], settings.get("provider", ""))
    if _is_minimax_provider(settings.get("provider", "")):
        if not str(settings.get("model") or "").strip().lower().startswith("image-"):
            settings["model"] = _default_image_generation_model(settings.get("provider", "")) or "image-01"
        styled_prompt = clean_prompt
        if str(style or "").strip():
            styled_prompt = f"{clean_prompt}\nStyle: {str(style).strip()}"
        payload = {
            "model": settings["model"],
            "prompt": styled_prompt[:1500],
            "aspect_ratio": _minimax_aspect_ratio(size),
            "response_format": "url",
            "n": 1,
            "prompt_optimizer": True,
        }
    else:
        payload = {
            "model": settings["model"],
            "prompt": clean_prompt,
            "size": str(size or "1024x1024"),
            "n": 1,
            "response_format": "b64_json",
        }
        if str(style or "").strip():
            payload["style"] = str(style).strip()

    def _post(body: dict) -> dict:
        req = urllib.request.Request(
            endpoint,
            data=json.dumps(body, ensure_ascii=False).encode("utf-8"),
            headers={
                "Authorization": f"Bearer {settings['api_key']}",
                "Content-Type": "application/json",
                "User-Agent": "TiangongImage/1.0",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=180) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        return loads_json_object(raw, source="image_generation_provider")

    try:
        try:
            data = _post(payload)
        except Exception as first_exc:
            lean_payload = {k: v for k, v in payload.items() if k not in {"response_format", "style", "prompt_optimizer"}}
            data = _post(lean_payload)
            data["_first_attempt_error"] = str(first_exc)[:240]
        base_resp = data.get("base_resp") if isinstance(data.get("base_resp"), dict) else {}
        if base_resp and int(base_resp.get("status_code") or 0) != 0:
            raise ValueError(f"minimax_image_generation_failed:{base_resp.get('status_code')}:{base_resp.get('status_msg')}")
        final_path, source = _write_generated_image(data, output_path)
        return {
            "zhuangtai": "wancheng",
            "leixing": "image_generate",
            "lujing": str(final_path),
            "daxiao": final_path.stat().st_size,
            "provider": settings["provider"],
            "model": settings["model"],
            "config_source": settings.get("config_source", ""),
            "model_source": settings.get("model_source", ""),
            "base_url": settings["base_url"],
            "endpoint": endpoint,
            "source": source,
            "prompt": clean_prompt,
            "size": payload.get("size") or payload.get("aspect_ratio", ""),
            "aspect_ratio": payload.get("aspect_ratio", ""),
        }
    except Exception as exc:
        detail = error_payload(exc, source="image_generation_provider", ok_key=False)
        return {
            "zhuangtai": "cuowu",
            "leixing": "image_generate",
            "cuowu": detail.get("error", str(exc))[:500],
            "error_code": detail.get("error_code", type(exc).__name__),
            "detail": detail.get("detail", str(exc))[:800],
            "raw_preview": detail.get("raw_preview", ""),
            "provider": settings.get("provider", ""),
            "model": settings.get("model", ""),
            "config_source": settings.get("config_source", ""),
            "model_source": settings.get("model_source", ""),
            "base_url": settings.get("base_url", ""),
            "endpoint": endpoint,
            "tishi": "MiniMax 生图使用 /v1/image_generation，模型 image-01/image-01-live；OpenAI-compatible 生图才使用 /images/generations。请确认 API key 已开通对应生图额度。",
        }


def _video_generation_settings() -> dict:
    from ..peizhi import (
        MOREN_PROVIDER,
        duqu_api_miyao,
        duqu_configured_model_ming,
        duqu_moren_provider,
        duqu_provider_base_url,
    )

    video_override = any(os.environ.get(key) for key in (
        "TIANGONG_VIDEO_PROVIDER",
        "TIANGONG_VIDEO_BASE_URL",
        "TIANGONG_VIDEO_MODEL",
        "TIANGONG_VIDEO_API_KEY",
    ))
    provider = os.environ.get("TIANGONG_VIDEO_PROVIDER") or duqu_moren_provider(MOREN_PROVIDER)
    provider = str(provider or "").strip()
    base_url = os.environ.get("TIANGONG_VIDEO_BASE_URL") or duqu_provider_base_url(provider) or ""
    api_key = os.environ.get("TIANGONG_VIDEO_API_KEY") or duqu_api_miyao(provider) or ""
    default_video_model = _default_video_generation_model(provider)
    configured_model = duqu_configured_model_ming(provider)
    model = os.environ.get("TIANGONG_VIDEO_MODEL") or default_video_model or configured_model
    return {
        "provider": provider,
        "base_url": str(base_url or "").rstrip("/"),
        "api_key": api_key,
        "model": str(model or "").strip(),
        "config_source": "video_override" if video_override else "current_llm_builtin",
        "model_source": (
            "env:TIANGONG_VIDEO_MODEL"
            if os.environ.get("TIANGONG_VIDEO_MODEL")
            else "provider_video_default"
            if default_video_model
            else "configured_llm_model"
            if configured_model
            else "missing"
        ),
    }


def _default_video_generation_model(provider: str) -> str:
    if _is_minimax_provider(provider):
        return "MiniMax-Hailuo-2.3"
    return ""


def _video_generation_endpoint(base_url: str, provider: str, action: str = "create") -> str:
    value = str(base_url or "").strip().rstrip("/")
    if not value:
        return ""
    value = _trim_known_generation_endpoint(value)
    if not _is_minimax_provider(provider):
        return ""
    if value.endswith("/v1"):
        base = value
    else:
        base = f"{value}/v1"
    if action == "query":
        return f"{base}/query/video_generation"
    if action == "retrieve":
        return f"{base}/files/retrieve"
    return f"{base}/video_generation"


def _safe_video_output_path(output_name: str = "") -> Path:
    from ..workspace_settings import duqu_workspace_root

    root = duqu_workspace_root()
    raw = str(output_name or "").strip()
    if raw:
        p = Path(raw).expanduser()
        if not p.is_absolute():
            p = root / "generated" / "videos" / p
    else:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = root / "generated" / "videos" / f"video-{stamp}.mp4"
    if p.suffix.lower() not in {".mp4", ".mov", ".webm", ".mkv"}:
        p = p.with_suffix(".mp4")
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _minimax_video_duration(value: Any) -> int:
    try:
        number = int(value)
    except Exception:
        number = 6
    return 10 if number >= 8 else 6


def _minimax_video_resolution(value: str) -> str:
    raw = str(value or "").strip().upper().replace(" ", "")
    aliases = {
        "720": "720P",
        "720P": "720P",
        "768": "768P",
        "768P": "768P",
        "1080": "1080P",
        "1080P": "1080P",
    }
    return aliases.get(raw, "768P")


def _image_reference_for_provider(value: str) -> str:
    raw = str(value or "").strip()
    if not raw:
        return ""
    lowered = raw.lower()
    if lowered.startswith(("http://", "https://", "data:image/")):
        return raw
    path = _workspace_path(raw)
    if not path.exists():
        raise ValueError(f"image_reference_not_found:{raw}")
    return _image_data_url(path, max_bytes=10 * 1024 * 1024)


def _minimax_post_json(endpoint: str, api_key: str, payload: dict, timeout: int = 180) -> dict:
    req = urllib.request.Request(
        endpoint,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "User-Agent": "TiangongVideo/1.0",
        },
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        raw = resp.read().decode("utf-8", errors="replace")
    return loads_json_object(raw, source="minimax_video_provider")


def _minimax_get_json(endpoint: str, api_key: str, params: dict, timeout: int = 80) -> dict:
    url = endpoint
    query = urllib.parse.urlencode({k: v for k, v in params.items() if v not in (None, "")})
    if query:
        url = f"{endpoint}?{query}"
    req = urllib.request.Request(
        url,
        headers={
            "Authorization": f"Bearer {api_key}",
            "User-Agent": "TiangongVideo/1.0",
        },
        method="GET",
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        raw = resp.read().decode("utf-8", errors="replace")
    return loads_json_object(raw, source="minimax_video_provider")


def _raise_for_minimax_base_resp(data: dict, source: str) -> None:
    base_resp = data.get("base_resp") if isinstance(data.get("base_resp"), dict) else {}
    if not base_resp:
        return
    try:
        status_code = int(base_resp.get("status_code") or 0)
    except Exception:
        status_code = 0
    if status_code != 0:
        raise ValueError(f"{source}:{status_code}:{base_resp.get('status_msg')}")


def _download_url_to_path(url: str, output_path: Path, timeout: int = 300) -> tuple[Path, int, str]:
    clean_url = str(url or "").strip()
    if clean_url.startswith("//"):
        clean_url = "https:" + clean_url
    elif clean_url.startswith("www."):
        clean_url = "https://" + clean_url
    if not clean_url:
        raise ValueError("download_url_empty")
    req = urllib.request.Request(clean_url, headers={"User-Agent": "TiangongVideo/1.0"})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        payload = resp.read()
        content_type = resp.headers.get("Content-Type", "")
    if output_path.suffix.lower() == ".mp4" and "quicktime" in content_type.lower():
        output_path = output_path.with_suffix(".mov")
    elif output_path.suffix.lower() == ".mp4" and "webm" in content_type.lower():
        output_path = output_path.with_suffix(".webm")
    output_path.write_bytes(payload)
    return output_path, len(payload), clean_url


def _video_generation_call(
    prompt: str,
    duration: int = 6,
    resolution: str = "768P",
    model: str = "",
    first_frame_image: str = "",
    last_frame_image: str = "",
    output_name: str = "",
    prompt_optimizer: bool = True,
    poll_interval: int = 10,
    max_wait: int = 600,
) -> dict:
    clean_prompt = str(prompt or "").strip()
    if not clean_prompt:
        return {"zhuangtai": "cuowu", "cuowu": "prompt 为空"}
    settings = _video_generation_settings()
    if not _is_minimax_provider(settings.get("provider", "")):
        return {
            "zhuangtai": "weishixian",
            "leixing": "video_generate",
            "cuowu": "当前仅实现 MiniMax 视频生成 native 适配",
            "provider": settings.get("provider", ""),
            "tishi": "请把视频生成 provider 设置为 minimax/minimax_m3，或配置 TIANGONG_VIDEO_PROVIDER=minimax_m3。",
        }
    if not settings.get("base_url") or not settings.get("api_key"):
        return {
            "zhuangtai": "weishixian",
            "leixing": "video_generate",
            "cuowu": "未配置 MiniMax 视频生成 base_url/api_key",
            "provider": settings.get("provider", ""),
            "model": settings.get("model", ""),
            "config_source": settings.get("config_source", ""),
            "model_source": settings.get("model_source", ""),
        }

    selected_model = str(model or "").strip() or settings.get("model") or "MiniMax-Hailuo-2.3"
    output_path = _safe_video_output_path(output_name)
    create_endpoint = _video_generation_endpoint(settings["base_url"], settings["provider"], "create")
    query_endpoint = _video_generation_endpoint(settings["base_url"], settings["provider"], "query")
    retrieve_endpoint = _video_generation_endpoint(settings["base_url"], settings["provider"], "retrieve")
    payload = {
        "model": selected_model,
        "prompt": clean_prompt[:2000],
        "duration": _minimax_video_duration(duration),
        "resolution": _minimax_video_resolution(resolution),
        "prompt_optimizer": bool(prompt_optimizer),
    }
    if first_frame_image:
        payload["first_frame_image"] = _image_reference_for_provider(first_frame_image)
    if last_frame_image:
        payload["last_frame_image"] = _image_reference_for_provider(last_frame_image)

    try:
        create_data = _minimax_post_json(create_endpoint, settings["api_key"], payload)
        _raise_for_minimax_base_resp(create_data, "minimax_video_create_failed")
        task_id = str(create_data.get("task_id") or create_data.get("id") or "").strip()
        if not task_id:
            raise ValueError(f"minimax_video_missing_task_id:{str(create_data)[:500]}")

        interval = max(3, min(60, int(poll_interval or 10)))
        deadline = time.time() + max(30, min(1800, int(max_wait or 600)))
        last_status: dict[str, Any] = {}
        while time.time() <= deadline:
            time.sleep(interval)
            status_data = _minimax_get_json(query_endpoint, settings["api_key"], {"task_id": task_id})
            _raise_for_minimax_base_resp(status_data, "minimax_video_query_failed")
            last_status = status_data
            status = str(status_data.get("status") or status_data.get("state") or "").strip()
            lowered = status.lower()
            if lowered in {"success", "succeed", "completed", "complete", "finished"}:
                file_id = str(status_data.get("file_id") or status_data.get("output_file_id") or "").strip()
                if not file_id:
                    raise ValueError(f"minimax_video_missing_file_id:{str(status_data)[:500]}")
                retrieve_data = _minimax_get_json(retrieve_endpoint, settings["api_key"], {"file_id": file_id})
                _raise_for_minimax_base_resp(retrieve_data, "minimax_video_retrieve_failed")
                file_obj = retrieve_data.get("file") if isinstance(retrieve_data.get("file"), dict) else retrieve_data
                download_url = (
                    str(file_obj.get("download_url") or file_obj.get("url") or "").strip()
                    if isinstance(file_obj, dict)
                    else ""
                )
                final_path, size_bytes, downloaded_from = _download_url_to_path(download_url, output_path)
                return {
                    "zhuangtai": "wancheng",
                    "leixing": "video_generate",
                    "lujing": str(final_path),
                    "daxiao": size_bytes,
                    "provider": settings["provider"],
                    "model": selected_model,
                    "config_source": settings.get("config_source", ""),
                    "model_source": settings.get("model_source", ""),
                    "base_url": settings["base_url"],
                    "endpoint": create_endpoint,
                    "query_endpoint": query_endpoint,
                    "retrieve_endpoint": retrieve_endpoint,
                    "task_id": task_id,
                    "file_id": file_id,
                    "status": status,
                    "downloaded_from": downloaded_from,
                    "prompt": clean_prompt,
                    "duration": payload["duration"],
                    "resolution": payload["resolution"],
                }
            if lowered in {"fail", "failed", "error", "cancelled", "canceled"}:
                raise ValueError(f"minimax_video_task_failed:{str(status_data)[:800]}")
        return {
            "zhuangtai": "jinxingzhong",
            "leixing": "video_generate",
            "task_id": task_id,
            "status": str(last_status.get("status") or ""),
            "last_status": last_status,
            "provider": settings["provider"],
            "model": selected_model,
            "endpoint": create_endpoint,
            "query_endpoint": query_endpoint,
            "tishi": "视频任务仍在生成中，可稍后按 task_id 查询或重试。",
        }
    except Exception as exc:
        detail = error_payload(exc, source="minimax_video_generation_provider", ok_key=False)
        return {
            "zhuangtai": "cuowu",
            "leixing": "video_generate",
            "cuowu": detail.get("error", str(exc))[:500],
            "error_code": detail.get("error_code", type(exc).__name__),
            "detail": detail.get("detail", str(exc))[:1000],
            "raw_preview": detail.get("raw_preview", ""),
            "provider": settings.get("provider", ""),
            "model": selected_model,
            "config_source": settings.get("config_source", ""),
            "model_source": settings.get("model_source", ""),
            "base_url": settings.get("base_url", ""),
            "endpoint": create_endpoint,
            "query_endpoint": query_endpoint,
            "retrieve_endpoint": retrieve_endpoint,
            "tishi": "MiniMax 视频生成使用 /v1/video_generation 创建任务，再用 /v1/query/video_generation 和 /v1/files/retrieve 获取成片。请确认账号已开通视频生成额度。",
        }


def _default_vision_model(provider: str) -> str:
    provider = str(provider or "").lower()
    if provider == "openai":
        return "gpt-4o-mini"
    if provider in {"minimax", "minimax_m3"}:
        return "MiniMax-M3"
    if provider == "google":
        return "gemini-2.0-flash"
    return provider or "vision-model"


def _extract_chat_content(data: dict) -> str:
    try:
        content = data["choices"][0]["message"]["content"]
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            parts = []
            for item in content:
                if isinstance(item, dict):
                    parts.append(str(item.get("text") or item.get("content") or ""))
            return "\n".join(part for part in parts if part)
    except Exception:
        pass
    try:
        return str(data["choices"][0]["delta"]["content"])
    except Exception:
        return ""


def _clean_model_text(text: str) -> str:
    value = str(text or "")
    value = re.sub(r"(?is)<think>.*?</think>", "", value)
    value = re.sub(r"(?is)<thinking>.*?</thinking>", "", value)
    value = re.sub(r"(?is)<reasoning>.*?</reasoning>", "", value)
    return value.strip()


def _vision_image_call(path: Path, prompt: str, purpose: str = "general", max_tokens: int = 2048) -> dict:
    settings = _vision_settings()
    if not settings.get("provider") or not settings.get("base_url") or not settings.get("api_key"):
        return {"ok": False, "state": "unconfigured", "error": "vision provider/base_url/api_key not configured"}
    views = []
    try:
        views = _enhanced_image_views(path, purpose=purpose)
        content = [
            {
                "type": "text",
                "text": (
                    prompt
                    + "\nIf multiple images are provided, they are enhanced or cropped views of the same source image. "
                    + "Use them together and do not describe duplicate views as separate images."
                ),
            }
        ]
        for index, view in enumerate(views, start=1):
            content.append({"type": "text", "text": f"View {index}: {view.get('label', 'image')}."})
            content.append({"type": "image_url", "image_url": {"url": view["url"]}})
        payload = {
            "model": settings["model"],
            "messages": [{"role": "user", "content": content}],
            "max_tokens": max_tokens,
        }
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        req = urllib.request.Request(
            f"{settings['base_url']}/chat/completions",
            data=body,
            headers={
                "Authorization": f"Bearer {settings['api_key']}",
                "Content-Type": "application/json",
                "User-Agent": "TiangongVision/1.0",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=80) as resp:
            raw = resp.read().decode("utf-8", errors="replace")
        data = loads_json_object(raw, source="vision_provider")
        text = _clean_model_text(_extract_chat_content(data))
        if not text:
            return {"ok": False, "state": "empty_response", "error": raw[:500], **settings}
        return {
            "ok": True,
            "state": "ok",
            "text": text,
            "provider": settings["provider"],
            "model": settings["model"],
            "vision_views": [view.get("label", "image") for view in views],
        }
    except Exception as exc:
        if purpose != "general" and len(views) > 1:
            return _vision_image_call(path, prompt, purpose="general", max_tokens=max_tokens)
        detail = error_payload(exc, source="vision_provider", ok_key=False)
        return {
            "ok": False,
            "state": "failed",
            "error": detail.get("error", str(exc))[:500],
            "error_code": detail.get("error_code", type(exc).__name__),
            "detail": detail.get("detail", str(exc))[:800],
            "raw_preview": detail.get("raw_preview", ""),
            "provider": settings.get("provider", ""),
            "model": settings.get("model", ""),
        }


class JirouCeng:
    """肌肉层：她的手脚"""

    @staticmethod
    def _zhixing_novel_tool(canshu: dict) -> dict[str, Any]:
        payload = canshu if isinstance(canshu, dict) else {}
        command = str(payload.get("command") or payload.get("action") or "").strip().lower()
        aliases = {
            "pack": "package",
            "zip": "package",
            "package_zip": "package",
            "check": "gate",
            "review": "audit",
            "contract_init": "contract-init",
            "story_contract": "contract-init",
            "chapter_card": "chapter-card",
            "card": "chapter-card",
            "contract_check": "contract-check",
            "semantic_check": "contract-check",
        }
        command = aliases.get(command, command)
        supported_commands = ["init", "contract-init", "chapter-card", "status", "gate", "audit", "contract-check", "package"]
        if command not in set(supported_commands):
            return {
                "ok": False,
                "zhuangtai": "cuowu",
                "effect": "execute",
                "cuowu": "novel_tool command must be one of init/contract-init/chapter-card/status/gate/audit/contract-check/package",
                "supported_commands": supported_commands,
            }

        project_raw = payload.get("project_dir") or payload.get("path")
        if not project_raw:
            return {
                "ok": False,
                "zhuangtai": "cuowu",
                "effect": "execute",
                "command": command,
                "cuowu": "project_dir is required",
            }
        project_dir = _workspace_path(project_raw)
        argv: list[str] = [command, "--project-dir", str(project_dir)]

        if command == "init":
            title = str(payload.get("title") or project_dir.name or "未命名小说").strip()
            genre = str(payload.get("genre") or "网文").strip()
            mode = str(payload.get("mode") or "monitor").strip().lower()
            if mode not in {"fast", "monitor", "strict"}:
                mode = "monitor"
            try:
                chapters = int(payload.get("chapters") or 30)
            except Exception:
                chapters = 30
            argv.extend(["--title", title, "--genre", genre, "--chapters", str(max(1, chapters)), "--mode", mode])
            target_reader = str(payload.get("target_reader") or payload.get("target-reader") or "").strip()
            if target_reader:
                argv.extend(["--target-reader", target_reader])
            brief = str(payload.get("brief") or "").strip()
            if brief:
                argv.extend(["--brief", brief])
        elif command == "contract-init":
            for key, cli_key in [
                ("title", "--title"),
                ("genre", "--genre"),
                ("core_promise", "--core-promise"),
                ("main_characters", "--main-characters"),
                ("active_volume", "--active-volume"),
                ("style_notes", "--style-notes"),
                ("forbidden_drift", "--forbidden-drift"),
            ]:
                value = str(payload.get(key) or payload.get(key.replace("_", "-")) or "").strip()
                if value:
                    argv.extend([cli_key, value])
        elif command == "chapter-card":
            if payload.get("chapter_num") in (None, ""):
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "chapter_num is required"}
            title = str(payload.get("title") or "").strip()
            if not title:
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "title is required"}
            argv.extend(["--chapter-num", str(int(payload.get("chapter_num"))), "--title", title])
            for key, cli_key in [
                ("pov", "--pov"),
                ("time", "--time"),
                ("location", "--location"),
                ("characters", "--characters"),
                ("must_include", "--must-include"),
                ("must_not_include", "--must-not-include"),
                ("conflict", "--conflict"),
                ("ending_hook", "--ending-hook"),
            ]:
                value = str(payload.get(key) or payload.get(key.replace("_", "-")) or "").strip()
                if value:
                    argv.extend([cli_key, value])
        elif command == "gate":
            stage = str(payload.get("stage") or "").strip().upper()
            if not stage:
                return {"ok": False, "zhuangtai": "cuowu", "effect": "read", "command": command, "cuowu": "stage is required"}
            argv.extend(["--stage", stage])
            if payload.get("chapter_num") not in (None, ""):
                argv.extend(["--chapter-num", str(int(payload.get("chapter_num")))])
        elif command == "audit":
            chapter_raw = payload.get("chapter")
            if not chapter_raw:
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "chapter is required"}
            argv.extend(["--chapter", str(_workspace_path(chapter_raw))])
            if payload.get("chapter_num") in (None, ""):
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "chapter_num is required"}
            argv.extend(["--chapter-num", str(int(payload.get("chapter_num")))])
            if payload.get("min_chars") not in (None, ""):
                argv.extend(["--min-chars", str(int(payload.get("min_chars")))])
        elif command == "contract-check":
            chapter_raw = payload.get("chapter")
            if not chapter_raw:
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "chapter is required"}
            argv.extend(["--chapter", str(_workspace_path(chapter_raw))])
            if payload.get("chapter_num") in (None, ""):
                return {"ok": False, "zhuangtai": "cuowu", "effect": "write", "command": command, "cuowu": "chapter_num is required"}
            argv.extend(["--chapter-num", str(int(payload.get("chapter_num")))])
        elif command == "package":
            output_raw = payload.get("output")
            if output_raw:
                argv.extend(["--output", str(_workspace_path(output_raw))])

        import contextlib
        import io

        stdout_buffer = io.StringIO()
        stderr_buffer = io.StringIO()
        try:
            module = _load_novel_tool_module()
            with contextlib.redirect_stdout(stdout_buffer), contextlib.redirect_stderr(stderr_buffer):
                try:
                    return_code = int(module.main(argv))
                except SystemExit as exc:
                    code = exc.code
                    return_code = int(code) if isinstance(code, int) else (0 if code in (None, "") else 1)
        except Exception as exc:
            return {
                "ok": False,
                "zhuangtai": "cuowu",
                "effect": "execute",
                "command": command,
                "cuowu": str(exc)[:500],
                "script": str(_novel_tool_script()),
            }

        stdout = stdout_buffer.getvalue()
        stderr = stderr_buffer.getvalue()
        data = _json_from_stdout(stdout)
        ok = return_code == 0 and data.get("ok", True) is not False
        effect = "read" if command in {"status", "gate"} else "write"

        path_candidates: list[str] = []
        if command in {"init", "contract-init", "chapter-card", "status", "gate"}:
            path_candidates.append(str(project_dir))
        if isinstance(data, dict):
            for key in ("project_dir", "contract", "chapter_card", "report", "status_file", "package", "sha256"):
                value = data.get(key)
                if value:
                    path_candidates.append(str(value))
        seen_paths: set[str] = set()
        paths = []
        for item in path_candidates:
            if item and item not in seen_paths:
                paths.append(item)
                seen_paths.add(item)
        primary_path = next((item for item in paths if Path(item).exists()), paths[0] if paths else "")

        result: dict[str, Any] = {
            "ok": ok,
            "zhuangtai": "wancheng" if ok else "cuowu",
            "effect": effect,
            "command": command,
            "script": str(_novel_tool_script()),
            "args": argv,
            "data": data,
            "fanhui_ma": return_code,
            "stdout": stdout[-4000:],
            "stderr": stderr[-2000:],
            "lujing": primary_path,
        }
        if paths:
            if effect == "write":
                result["updated_paths"] = paths
            result["paths"] = paths
            result["readback"] = _result_readback(primary_path)
        if not ok:
            error_text = data.get("error") or data.get("errors") or stderr.strip() or f"novel_tool exited with {return_code}"
            result["cuowu"] = _plain_text(error_text, 1000)
        return result

    @staticmethod
    def _omni_args(canshu: dict[str, Any]) -> dict[str, Any]:
        args = canshu.get("args") if isinstance(canshu, dict) else {}
        return dict(args) if isinstance(args, dict) else {}

    @staticmethod
    def _omni_value(canshu: dict[str, Any], args: dict[str, Any], *keys: str) -> Any:
        for key in keys:
            if key in args and args.get(key) not in (None, ""):
                return args.get(key)
            if isinstance(canshu, dict) and key in canshu and canshu.get(key) not in (None, ""):
                return canshu.get(key)
        return ""

    def zhixing(self, yingshe: GongjuYingshe, canshu: dict, *, call_id: str = "") -> dict[str, Any]:
        """执行工具。当前主链只允许 omni_body。"""
        mingcheng = yingshe.mingcheng
        if mingcheng != "omni_body":
            return _normalise_tool_result(
                mingcheng,
                getattr(yingshe, "effect", "unknown"),
                {
                    "ok": False,
                    "zhuangtai": "retired",
                    "cuowu": f"{mingcheng} 已从当前主链退役；请使用 omni_body。",
                    "replacement_tool": "omni_body",
                },
            )

        try:
            proposal = dict(canshu or {})
            # The model can propose action/target/args only. Workspace and all
            # authority objects are rebound by the backend and total gateway.
            proposal.pop("workspace", None)
            proposal.pop("__capability_grant", None)
            proposal.pop("__runtime", None)
            raw_args = proposal.get("args")
            if raw_args is None:
                raw_args = {}
            if not isinstance(raw_args, dict):
                raise ValueError("omni_body.args_object_required")
            proposal = {
                "action": str(proposal.get("action") or "").strip(),
                "target": str(proposal.get("target") or "").strip(),
                "args": dict(raw_args),
            }
            authority = issue_omni_grant(
                proposal,
                workspace=_workspace_root(),
                call_id=call_id,
            )
            authorized = {
                "action": str(proposal.get("action") or ""),
                "target": str(proposal.get("target") or ""),
                "args": dict(proposal.get("args") or {}),
                "workspace": str(_workspace_root()),
                "__capability_grant": authority["grant"],
                "__runtime": authority["runtime"],
            }
            result = _normalise_tool_result(
                mingcheng,
                getattr(yingshe, "effect", "unknown"),
                _run_omni_body_tool(authorized),
            )
            result["authority_receipt"] = _authority_receipt(authority)
            return result
        except Exception as e:
            detail = error_payload(e, source=f"tool:{mingcheng}", ok_key=False)
            return _normalise_tool_result(mingcheng, getattr(yingshe, "effect", "unknown"), {
                "cuowu": detail.get("error", str(e)),
                "error_code": detail.get("error_code", type(e).__name__),
                "detail": detail.get("detail", str(e)),
                "source": detail.get("source", f"tool:{mingcheng}"),
                "raw_preview": detail.get("raw_preview", ""),
                "leixing": type(e).__name__,
            })

    @staticmethod
    def _jianmulu(
        path: str = "",
        dir_path: str = "",
        directory: str = "",
        folder: str = "",
        target: str = "",
        confirm: bool = False,
        **_ignored: Any,
    ) -> dict:
        path_value = target or dir_path or directory or folder or path
        if not str(path_value or "").strip():
            return {"ok": False, "zhuangtai": "cuowu", "cuowu": "MISSING_PATH"}
        p = _workspace_path(path_value)
        try:
            p.mkdir(parents=True, exist_ok=True)
            exists = p.exists()
            is_dir = p.is_dir()
            readback = {"ok": bool(exists and is_dir), "exists": exists, "is_dir": is_dir}
            if not readback["ok"]:
                return {
                    "ok": False,
                    "zhuangtai": "cuowu",
                    "effect": "write",
                    "lujing": str(p),
                    "path": str(p),
                    "cuowu": "DIR_VERIFY_FAILED",
                    "readback": readback,
                }
            return {
                "ok": True,
                "zhuangtai": "wancheng",
                "effect": "write",
                "lujing": str(p),
                "path": str(p),
                "updated_paths": [str(p)],
                "readback": readback,
                "zhaiyao": f"Directory ready: {p}",
            }
        except Exception as e:
            return {
                "ok": False,
                "zhuangtai": "cuowu",
                "effect": "write",
                "lujing": str(p),
                "path": str(p),
                "cuowu": str(e),
                "readback": {"ok": False, "error": str(e)[:200]},
            }

    @staticmethod
    def _duwenjian(path: str, offset: int = 1, limit: int = 500, confirm: bool = False, **_ignored: Any) -> dict:
        """读取文件"""
        p = _workspace_path(path)
        if not p.exists():
            return {"cuowu": f"文件不存在: {path}", "lujing": str(p)}
        if p.is_dir():
            return {"cuowu": f"路径是目录: {path}", "lujing": str(p)}
        try:
            hang = p.read_text(encoding="utf-8").split("\n")
            zong = len(hang)
            start = max(0, offset - 1)
            end = min(zong, start + limit)
            neirong = hang[start:end]
            return {
                "ok": True,
                "zhuangtai": "wancheng",
                "effect": "read",
                "neirong": "\n".join(neirong),
                "zong_hangshu": zong,
                "xianshi_hangshu": len(neirong),
                "lujing": str(p)
            }
        except UnicodeDecodeError:
            return {"cuowu": "二进制文件或编码不兼容", "lujing": str(p)}
        except Exception as e:
            return {"cuowu": str(e), "lujing": str(p)}

    @staticmethod
    def _xiewenjian(path: str, content: str, encoding: str = "utf-8", confirm: bool = False, **_ignored: Any) -> dict:
        """写入文件"""
        p = _workspace_path(path)
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            text = str(content)
            p.write_text(text, encoding=encoding or "utf-8")
            readback = p.read_text(encoding=encoding or "utf-8")
            if readback != text:
                return {
                    "ok": False,
                    "zhuangtai": "cuowu",
                    "lujing": str(p),
                    "cuowu": "WRITE_VERIFY_FAILED",
                    "readback": {"ok": False, "bytes": p.stat().st_size if p.exists() else 0},
                }
            digest = hashlib.sha256(readback.encode(encoding or "utf-8")).hexdigest()
            return {
                "ok": True,
                "zhuangtai": "yixieru",
                "effect": "write",
                "lujing": str(p),
                "zishu": len(text),
                "sha256": digest,
                "updated_paths": [str(p)],
                "readback": {"ok": True, "bytes": p.stat().st_size, "sha256": digest},
            }
        except Exception as e:
            return {"cuowu": str(e), "lujing": str(p)}

    @staticmethod
    def _chuangjian_docx(
        path: str = "",
        title: str = "",
        paragraphs: Any = None,
        sections: Any = None,
        tables: Any = None,
    ) -> dict:
        """Create a real .docx file using python-docx."""
        p = _office_output_path(path, ".docx", "documents")
        try:
            from docx import Document
        except ImportError:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 python-docx"}
        try:
            doc = Document()
            if title:
                doc.add_heading(str(title), level=0)
            paragraph_count = 0
            for paragraph in _as_list(paragraphs):
                text = _plain_text(paragraph, limit=6000)
                if text:
                    doc.add_paragraph(text)
                    paragraph_count += 1
            for section in _as_list(sections):
                if not isinstance(section, dict):
                    text = _plain_text(section, limit=6000)
                    if text:
                        doc.add_paragraph(text)
                        paragraph_count += 1
                    continue
                heading = _plain_text(section.get("heading") or section.get("title"), limit=300)
                if heading:
                    level = JirouCeng._clamp_int(section.get("level", 1), minimum=1, maximum=4, default=1)
                    doc.add_heading(heading, level=level)
                for paragraph in _as_list(section.get("paragraphs") or section.get("content") or section.get("body")):
                    text = _plain_text(paragraph, limit=6000)
                    if text:
                        doc.add_paragraph(text)
                        paragraph_count += 1
                for bullet in _as_list(section.get("bullets")):
                    text = _plain_text(bullet, limit=1000)
                    if text:
                        doc.add_paragraph(text, style="List Bullet")
                        paragraph_count += 1
                if section.get("table"):
                    JirouCeng._docx_add_table(doc, section.get("table"))
            table_count = 0
            for table in _as_list(tables):
                if JirouCeng._docx_add_table(doc, table):
                    table_count += 1
            if not title and paragraph_count == 0 and table_count == 0:
                doc.add_paragraph("")
            doc.save(str(p))
            Document(str(p))
            return _tool_success(p, "docx", paragraph_count=paragraph_count, table_count=table_count)
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _docx_add_table(doc: Any, table_payload: Any) -> bool:
        headers, rows = _normalise_table_payload(table_payload)
        if not headers and not rows:
            return False
        width = max(len(headers), max((len(row) for row in rows), default=0), 1)
        table = doc.add_table(rows=1 if headers else 0, cols=width)
        table.style = "Table Grid"
        if headers:
            for index, header in enumerate(headers):
                table.rows[0].cells[index].text = _plain_text(header, limit=300)
        for row in rows:
            cells = table.add_row().cells
            for index in range(width):
                cells[index].text = _plain_text(row[index] if index < len(row) else "", limit=1000)
        return True

    @staticmethod
    def _chuangjian_xlsx(
        path: str = "",
        sheets: Any = None,
        rows: Any = None,
        headers: Any = None,
        sheet_name: str = "Sheet1",
        title: str = "",
    ) -> dict:
        """Create a real .xlsx workbook using openpyxl."""
        p = _office_output_path(path, ".xlsx", "spreadsheets")
        try:
            from openpyxl import Workbook, load_workbook
            from openpyxl.styles import Font
        except ImportError:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 openpyxl"}
        try:
            wb = Workbook()
            wb.remove(wb.active)
            sheet_payloads = JirouCeng._normalise_xlsx_sheets(sheets, rows, headers, sheet_name, title)
            used_sheet_names: set[str] = set()
            total_data_rows = 0
            total_nonempty_cells = 0
            for index, payload in enumerate(sheet_payloads, 1):
                name = _safe_xlsx_sheet_name(payload.get("name"), f"Sheet{index}", used_sheet_names)
                ws = wb.create_sheet(title=name)
                headers2, rows2 = _normalise_table_rows(payload.get("rows") or [], payload.get("headers"))
                current_row = 1
                if payload.get("title"):
                    ws.cell(row=current_row, column=1, value=_plain_text(payload.get("title"), limit=500))
                    ws.cell(row=current_row, column=1).font = Font(bold=True, size=14)
                    current_row += 2
                if headers2:
                    for col, header in enumerate(headers2, 1):
                        cell = ws.cell(row=current_row, column=col, value=_plain_text(header, limit=300))
                        cell.font = Font(bold=True)
                    current_row += 1
                for row in rows2:
                    total_data_rows += 1
                    for col, value in enumerate(row, 1):
                        if value not in (None, ""):
                            total_nonempty_cells += 1
                        ws.cell(row=current_row, column=col, value=value)
                    current_row += 1
                if headers2:
                    ws.freeze_panes = f"A{3 if payload.get('title') else 2}"
                JirouCeng._xlsx_autofit(ws)
            if not wb.sheetnames:
                wb.create_sheet(title="Sheet1")
            if total_data_rows > 0 and total_nonempty_cells == 0:
                return {
                    "zhuangtai": "cuowu",
                    "lujing": str(p),
                    "cuowu": "XLSX_ROWS_EMPTY_AFTER_NORMALIZATION",
                    "row_count": total_data_rows,
                    "nonempty_cell_count": total_nonempty_cells,
                }
            wb.save(str(p))
            verify = load_workbook(str(p), read_only=True, data_only=True)
            sheet_count = len(verify.sheetnames)
            verify.close()
            return _tool_success(
                p,
                "xlsx",
                sheet_count=sheet_count,
                sheets=wb.sheetnames,
                row_count=total_data_rows,
                nonempty_cell_count=total_nonempty_cells,
            )
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _normalise_xlsx_sheets(sheets: Any, rows: Any, headers: Any, sheet_name: str, title: str) -> list[dict]:
        if sheets:
            if isinstance(sheets, dict):
                payloads = []
                for name, value in sheets.items():
                    if isinstance(value, dict):
                        payload = dict(value)
                        payload.setdefault("name", name)
                    else:
                        payload = {"name": name, "rows": value}
                    payloads.append(payload)
                return payloads
            payloads = []
            for index, item in enumerate(_as_list(sheets), 1):
                if isinstance(item, dict):
                    payload = dict(item)
                    payload.setdefault("name", f"Sheet{index}")
                else:
                    payload = {"name": f"Sheet{index}", "rows": item}
                payloads.append(payload)
            return payloads
        return [{"name": sheet_name or "Sheet1", "title": title, "headers": headers, "rows": rows or []}]

    @staticmethod
    def _xlsx_autofit(ws: Any) -> None:
        for column_cells in ws.columns:
            max_len = 8
            col_letter = column_cells[0].column_letter
            for cell in column_cells:
                value = "" if cell.value is None else str(cell.value)
                max_len = max(max_len, min(len(value) + 2, 42))
            ws.column_dimensions[col_letter].width = max_len

    @staticmethod
    def _chuangjian_pptx(
        path: str = "",
        title: str = "",
        subtitle: str = "",
        slides: Any = None,
    ) -> dict:
        """Create a real .pptx file using python-pptx."""
        p = _office_output_path(path, ".pptx", "presentations")
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 python-pptx"}
        try:
            prs = Presentation()
            slide_count = 0
            if title or subtitle:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = str(title or "")
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = str(subtitle or "")
                slide_count += 1
            for item in _as_list(slides):
                if isinstance(item, dict):
                    slide_title = _plain_text(item.get("title") or item.get("heading"), limit=200)
                    bullets = _as_list(item.get("bullets") or item.get("content") or item.get("body"))
                    table_payload = item.get("table")
                else:
                    slide_title = _plain_text(item, limit=200)
                    bullets = []
                    table_payload = None
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                body = slide.placeholders[1].text_frame
                body.clear()
                for index, bullet in enumerate(bullets):
                    text = _plain_text(bullet, limit=1000)
                    if not text:
                        continue
                    paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                    paragraph.text = text
                    paragraph.level = 0
                if table_payload:
                    JirouCeng._pptx_add_table(slide, table_payload, Inches)
                slide_count += 1
            if slide_count == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = ""
                slide_count = 1
            prs.save(str(p))
            verify = Presentation(str(p))
            return _tool_success(p, "pptx", slide_count=len(verify.slides))
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _pptx_add_table(slide: Any, table_payload: Any, Inches: Any) -> bool:
        headers, rows = _normalise_table_payload(table_payload)
        if not headers and not rows:
            return False
        data = ([headers] if headers else []) + rows
        cols = max((len(row) for row in data), default=1)
        table_shape = slide.shapes.add_table(len(data), cols, Inches(0.7), Inches(3.3), Inches(8.6), Inches(2.5))
        table = table_shape.table
        for r_index, row in enumerate(data):
            for c_index in range(cols):
                table.cell(r_index, c_index).text = _plain_text(row[c_index] if c_index < len(row) else "", limit=300)
        return True

    @staticmethod
    def _zhongduan(
        command: str = "",
        workdir: str = "",
        chaoshi: int = 120,
        timeout: int | None = None,
        confirm: bool = False,
    ) -> dict:
        """Execute a terminal command only inside the authoritative sandbox.

        ``confirm`` is retained only for call compatibility. It never bypasses
        an A5 classification; destructive or remote-script commands are
        rejected before a process is created.
        """
        try:
            command_text = str(command or "").strip()
            if not command_text:
                return {"cuowu": "[BAD_ARGS] missing command"}
            if len(command_text) > 8000:
                return {"cuowu": "[BAD_ARGS] command too long"}
            if JirouCeng._zhongduan_xuyao_queren(command_text):
                return {
                    "cuowu": "[A5_REJECTED] destructive or remote-script terminal command is forbidden",
                    "a5_rejected": True,
                    "requires_confirm": False,
                    "confirm_ignored": bool(confirm),
                    "command_preview": command_text[:500],
                }
            if timeout is not None:
                try:
                    chaoshi = int(timeout)
                except Exception:
                    pass
            chaoshi = max(1, min(int(chaoshi), 24 * 60 * 60))
            workspace = _workspace_root()
            wd = _workspace_path(workdir) if workdir else workspace
            try:
                wd.relative_to(workspace)
            except ValueError:
                return {"cuowu": "[SANDBOX] workdir must remain inside workspace", "workdir": str(wd)}
            if not wd.exists() or not wd.is_dir():
                return {"cuowu": "[SANDBOX] workdir is missing or not a directory", "workdir": str(wd)}

            from omni_body_skill.tools.sandbox_runtime import (
                SandboxError,
                SandboxLimits,
                SandboxRunner,
                WINDOWS_POWERSHELL_SHELL_MARKER,
            )

            if os.name == "nt":
                # Keep the source visible until SandboxRunner rewrites any
                # workspace paths. The marker then lets cmd enter the private
                # AppContainer workspace before launching encoded PowerShell;
                # PowerShell's own Set-Location provider rejects this otherwise
                # valid AppContainer path on Windows.
                shell_command = [
                    WINDOWS_POWERSHELL_SHELL_MARKER,
                    command_text,
                ]
            else:
                shell_command = ["/bin/sh", "-lc", command_text]
            runner = SandboxRunner(
                workspace=workspace,
                state_root=workspace / ".tiangong_sandboxes",
                trash_root=workspace / ".omni_trash",
                limits=SandboxLimits(timeout_seconds=chaoshi),
            )
            result = runner.run(
                shell_command,
                cwd=wd,
                timeout_seconds=chaoshi,
                op_id=f"legacy-terminal-{time.time_ns()}",
            )
            return_code = int(result.get("returncode", 1))
            stderr_text = str(result.get("stderr") or "")
            response = {
                "biaozhun_shuchu": str(result.get("stdout") or "")[-5000:],
                "biaozhun_cuowu": stderr_text[-2000:],
                "fanhui_ma": return_code,
                "chaoshi": False,
                "workdir": str(wd),
                "workdir_fallback": "",
                "sandbox": str(result.get("containment") or ""),
                "network": str(result.get("network") or ""),
                "changed_files": list(result.get("changed_files") or []),
                "deleted_files": list(result.get("deleted_files") or []),
            }
            if return_code != 0:
                response["cuowu"] = (
                    stderr_text.strip()[-2000:]
                    or f"[SANDBOX] command exited with status {return_code}"
                )
            return response
        except Exception as exc:
            if type(exc).__name__ == "SandboxError":
                return {"cuowu": f"[SANDBOX] {exc}", "chaoshi": str(exc) == "sandbox_timeout"}
            return {"cuowu": str(exc)}

    @staticmethod
    def _codex_shell_command(
        command: str = "",
        workdir: str = "",
        timeout_ms: int | None = None,
        timeout: int | None = None,
        chaoshi: int = 120,
        login: bool | None = None,
        confirm: bool = False,
        justification: str = "",
        prefix_rule: Any = None,
        sandbox_permissions: str = "",
        **_ignored: Any,
    ) -> dict:
        """Codex-compatible wrapper over the same fail-closed sandbox."""
        command_text = str(command or "").strip()
        if not command_text:
            return {"ok": False, "zhuangtai": "cuowu", "cuowu": "[BAD_ARGS] missing command"}
        if len(command_text) > 8000:
            return {"ok": False, "zhuangtai": "cuowu", "cuowu": "[BAD_ARGS] command too long"}
        if timeout_ms is not None:
            try:
                chaoshi = max(1, int(timeout_ms) // 1000)
            except Exception:
                pass
        elif timeout is not None:
            try:
                chaoshi = int(timeout)
            except Exception:
                pass
        result = JirouCeng._zhongduan(
            command=command_text,
            workdir=workdir,
            chaoshi=chaoshi,
            confirm=confirm,
        )
        ok = not result.get("cuowu") and int(result.get("fanhui_ma", 0) or 0) == 0
        result.update({
            "ok": bool(ok),
            "zhuangtai": "wancheng" if ok else "cuowu",
            "effect": "execute",
            "schema": "tiangong.v3.codex.shell_command.v1",
            "source": "openai/codex-compatible-shell_command",
            "timeout_ms": int(timeout_ms) if timeout_ms is not None else int(chaoshi) * 1000,
            "login_ignored": login is not None,
            "sandbox_permissions": str(sandbox_permissions or ""),
            "justification": str(justification or "")[:500],
            "prefix_rule": prefix_rule if isinstance(prefix_rule, list) else [],
        })
        return result

    @staticmethod
    def _zhongduan_xuyao_queren(command: str) -> bool:
        lowered = str(command or "").lower()
        return any(re.search(pattern, lowered, flags=re.IGNORECASE) for pattern in TERMINAL_DANGEROUS_COMMAND_PATTERNS)

    @staticmethod
    def _sousuowenjian(pattern: str, path: str = ".") -> dict:
        """搜索文件"""
        p = _workspace_path(path)
        try:
            jieguo = list(p.rglob(pattern))[:50]
            return {
                "pipei_shu": len(jieguo),
                "wenjian": [str(f) for f in jieguo[:20]],
                "lujing": str(p)
            }
        except Exception as e:
            return {"cuowu": str(e)}

    @staticmethod
    def _wangluosousuo(query: str, max_results: int = 10, mode: str = "") -> dict:
        """网络搜索 — 多引擎按序降级：百度 > 搜狗 > DDG > Bing HTML > Google > Bing RSS"""
        query = str(query or "").strip()
        if not query:
            return {"zhuangtai": "cuowu", "cuowu": "query 为空"}

        max_results = JirouCeng._clamp_int(max_results, minimum=1, maximum=15, default=10)
        explicit_mode = str(mode or "").strip().lower()
        news_mode = explicit_mode in {"news", "current", "xinwen", "新闻", "实时"} or JirouCeng._is_news_query(query)
        image_mode = explicit_mode in {"image", "images", "tupian", "图片", "图搜"}
        
        # 图片搜索走单独通道
        if image_mode:
            return JirouCeng._image_search(query, max_results=max_results, timeout=15)

        errors: list[str] = []
        relevance_rejections: list[str] = []

        providers: list[tuple[str, Any]] = []
        if news_mode:
            providers.append(("china_news_rss", JirouCeng._china_news_rss_search))
            providers.append(("global_news_rss", JirouCeng._global_news_rss_search))
        # 全引擎降级链：百度 → 搜狗 → DDG → Bing HTML → Google → Bing RSS
        providers.append(("baidu_html", JirouCeng._baidu_html_search))
        providers.append(("sogou_html", JirouCeng._sogou_html_search))
        providers.append(("duckduckgo_html", JirouCeng._duckduckgo_html_search))
        providers.append(("bing_html", JirouCeng._bing_html_search))
        providers.append(("google_html", JirouCeng._google_html_search))
        providers.append(("bing_rss", JirouCeng._bing_rss_search))
        if news_mode:
            providers.append(("google_news_rss", JirouCeng._google_news_rss_search))

        for provider_name, provider in providers:
            try:
                results = provider(query, max_results=max_results, timeout=12)
            except Exception as exc:
                errors.append(f"{provider_name}: {exc}")
                continue
            if results:
                results = JirouCeng._prefer_specific_results(results)
                results, rejected = JirouCeng._filter_relevant_results(query, results, max_results=max_results)
                if rejected:
                    relevance_rejections.append(f"{provider_name}: {len(rejected)} 条结果未达到相关性阈值")
                if not results:
                    continue
                content = JirouCeng._search_results_to_content(query, results, provider_name)
                return {
                    "zhuangtai": "wancheng",
                    "chaxun": query,
                    "mode": "news" if news_mode else "general",
                    "laiyuan": provider_name,
                    "jieguo": results,
                    "urls": [item["url"] for item in results if item.get("url")],
                    "content": content,
                    "zhaiyao": content[:4000],
                    "relevance_rejections": relevance_rejections[-5:],
                }

        duck = JirouCeng._duckduckgo_instant_answer(query)
        if duck.get("zhuangtai") == "wancheng":
            duck_item = {
                "title": str(duck.get("zhaiyao") or ""),
                "snippet": " ".join(str(item.get("miaoshu") or "") for item in duck.get("xiangguan", []) if isinstance(item, dict)),
                "url": str(duck.get("laiyuan") or ""),
                "source": "duckduckgo_instant_answer",
            }
            duck_score, _ = JirouCeng._search_item_relevance(query, duck_item)
            if duck_score > 0:
                duck["mode"] = "instant_answer"
                duck["rss_errors"] = errors[-3:]
                duck["relevance_rejections"] = relevance_rejections[-5:]
                return duck
            relevance_rejections.append("duckduckgo_instant_answer: 未达到相关性阈值或缺少可核验 URL")

        if errors:
            hint = "搜索服务暂时没有返回可用来源；可以换更具体关键词，或让用户提供新闻 URL 后用 web_readability_extract 读取正文。"
            return {
                "zhuangtai": "wu_jieguo",
                "chaxun": query,
                "cuowu": hint,
                "recoverable": True,
                "provider_errors": errors[-5:],
                "relevance_rejections": relevance_rejections[-5:],
                "duckduckgo": duck.get("cuowu", ""),
                "content": f"联网搜索未取得可核验来源。\n检索问题：{query}\n建议：{hint}\n注意：这代表工具/网络/关键词未命中，不代表模型知识截止。",
                "zhaiyao": hint,
            }
        return {
            "zhuangtai": "wu_jieguo",
            "chaxun": query,
            "recoverable": True,
            "relevance_rejections": relevance_rejections[-5:],
            "content": f"联网搜索没有命中可核验来源。\n检索问题：{query}\n建议换更具体关键词、指定站点或提供 URL。",
        }

    @staticmethod
    def _clamp_int(value: object, *, minimum: int, maximum: int, default: int) -> int:
        try:
            number = int(value)
        except (TypeError, ValueError):
            return default
        return max(minimum, min(maximum, number))

    @staticmethod
    def _is_news_query(text: str) -> bool:
        lowered = str(text or "").lower()
        return any(marker in lowered for marker in _NEWS_QUERY_MARKERS)

    @staticmethod
    def _news_window(text: str) -> str:
        lowered = str(text or "").lower()
        fresh_markers = ("今天", "今日", "刚刚", "实时", "快讯", "today", "breaking", "current")
        if any(marker in lowered for marker in fresh_markers):
            return "1d"
        return "7d"

    @staticmethod
    def _clean_html_text(value: object, *, limit: int = 1000) -> str:
        text = html.unescape(str(value or ""))
        text = re.sub(r"(?is)<script[^>]*>.*?</script>", " ", text)
        text = re.sub(r"(?is)<style[^>]*>.*?</style>", " ", text)
        text = re.sub(r"(?is)<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return text[: max(80, int(limit))]

    @staticmethod
    def _domain(url: str) -> str:
        try:
            return urllib.parse.urlparse(str(url or "")).netloc.replace("www.", "")
        except Exception:
            return ""

    @staticmethod
    def _search_site_filters(query: str) -> list[str]:
        domains: list[str] = []
        for match in re.findall(r"(?:^|\s)site:([^\s]+)", str(query or ""), flags=re.I):
            raw = match.strip().strip('"\'').lstrip("*.")
            parsed = urllib.parse.urlsplit(raw if "://" in raw else "//" + raw)
            host = str(parsed.hostname or raw.split("/", 1)[0]).rstrip(".").lower()
            if host and host not in domains:
                domains.append(host)
        return domains[:3]

    @staticmethod
    def _search_item_relevance(query: str, item: dict[str, str]) -> tuple[int, str]:
        url = str(item.get("url") or "").strip()
        parsed = urllib.parse.urlsplit(url)
        if parsed.scheme.lower() not in {"http", "https"} or not parsed.hostname:
            return 0, "missing_public_http_url"
        host = str(parsed.hostname or "").rstrip(".").lower()
        if host == "localhost" or host.endswith((".localhost", ".local", ".lan", ".home", ".internal")):
            return 0, "unsafe_result_host"
        try:
            result_ip = ipaddress.ip_address(host)
        except ValueError:
            result_ip = None
        if result_ip is not None and not result_ip.is_global:
            return 0, "unsafe_result_host"

        site_filters = JirouCeng._search_site_filters(query)
        if site_filters and not any(host == domain or host.endswith("." + domain) for domain in site_filters):
            return 0, "site_filter_mismatch"

        path = (parsed.path or "/").rstrip("/").lower() or "/"
        search_pages = (
            ("baidu.com", "/s"), ("sogou.com", "/web"), ("bing.com", "/search"),
            ("google.com", "/search"), ("google.com.hk", "/search"),
            ("duckduckgo.com", "/html"),
        )
        if any((host == domain or host.endswith("." + domain)) and path == page for domain, page in search_pages):
            return 0, "search_page_result"

        title = JirouCeng._clean_html_text(item.get("title"), limit=500).casefold()
        snippet = JirouCeng._clean_html_text(item.get("snippet"), limit=1200).casefold()
        source_name = JirouCeng._clean_html_text(item.get("source_name"), limit=300).casefold()
        haystack = " ".join((title, snippet, source_name, host, urllib.parse.unquote(url).casefold()))
        if any(marker.casefold() in haystack for marker in _SEARCH_AD_MARKERS):
            return 0, "probable_ad"

        terms = JirouCeng._query_terms(query)
        if not terms:
            source = str(item.get("source") or "").casefold()
            if JirouCeng._is_news_query(query) and ("rss" in source or bool(item.get("published"))):
                return 1, "broad_news_query"
            if site_filters:
                return 2, "site_filter_match"
            return 0, "no_discriminating_query_terms"

        matched = [term for term in terms if term in haystack]
        required_hits = 1 if len(terms) <= 2 else 2
        if len(matched) < required_hits:
            return 0, "query_term_mismatch"
        score = sum(min(len(term), 12) for term in matched)
        score += 2 * sum(min(len(term), 12) for term in matched if term in title)
        if site_filters:
            score += 8
        return max(score, 1), "relevant"

    @staticmethod
    def _filter_relevant_results(
        query: str,
        results: list[dict[str, str]],
        *,
        max_results: int,
    ) -> tuple[list[dict[str, str]], list[str]]:
        scored: list[tuple[int, int, dict[str, str]]] = []
        rejected: list[str] = []
        seen: set[str] = set()
        for index, item in enumerate(results):
            if not isinstance(item, dict):
                rejected.append("invalid_result")
                continue
            score, reason = JirouCeng._search_item_relevance(query, item)
            key = str(item.get("url") or item.get("title") or "").strip()
            if score <= 0 or not key or key in seen:
                rejected.append(reason if score <= 0 else "duplicate_result")
                continue
            seen.add(key)
            scored.append((score, -index, dict(item)))
        scored.sort(key=lambda row: (row[0], row[1]), reverse=True)
        return [row[2] for row in scored[:max_results]], rejected

    @staticmethod
    def _prefer_specific_results(results: list[dict[str, str]]) -> list[dict[str, str]]:
        def score(item: dict[str, str]) -> int:
            url = str(item.get("url") or "")
            title = str(item.get("title") or "")
            parsed = urllib.parse.urlparse(url)
            path = (parsed.path or "/").strip("/")
            value = 0
            if not path:
                value += 8
            elif len(path) <= 8:
                value += 3
            if re.search(r"(\d{4}|/20\d{2}|article|content|news|detail|doc-|\.s?html?$)", path, re.I):
                value -= 3
            if title.strip() in {"今日头条", "新闻", "中国新闻", "央视网", "新华网"}:
                value += 5
            return value

        return sorted(results, key=score)

    @staticmethod
    def _rss_items(url: str, *, timeout: float) -> ET.Element:
        request = urllib.request.Request(url, headers=_SEARCH_HEADERS)
        with urllib.request.urlopen(request, timeout=timeout) as response:
            raw = response.read(1024 * 1024)
        return ET.fromstring(raw.decode("utf-8", "replace"))

    @staticmethod
    def _query_terms(query: str) -> list[str]:
        text = html.unescape(str(query or "")).casefold()
        text = re.sub(r"(?:^|\s)(?:site|filetype|inurl|intitle):(?:\"[^\"]+\"|\S+)", " ", text, flags=re.I)
        stopwords = set(_SEARCH_INTENT_STOPWORDS) | {str(word).casefold() for word in _NEWS_STOPWORDS}
        for word in sorted(stopwords, key=len, reverse=True):
            if not word:
                continue
            if re.fullmatch(r"[a-z0-9]+", word):
                text = re.sub(rf"\b{re.escape(word)}\b", " ", text)
            else:
                text = text.replace(word, " ")

        terms: list[str] = []
        for token in re.findall(r"[a-z0-9][a-z0-9._+-]{1,63}", text):
            token = token.strip("._+-")
            if token and token not in {"com", "http", "https", "www"} and token not in terms:
                terms.append(token)
        for chunk in re.findall(r"[\u4e00-\u9fff]{2,}", text):
            if len(chunk) <= 12 and chunk not in terms:
                terms.append(chunk)
            if len(chunk) > 2:
                for index in range(len(chunk) - 1):
                    token = chunk[index:index + 2]
                    if token not in stopwords and token not in terms:
                        terms.append(token)
        return terms[:12]

    @staticmethod
    def _rss_feed_search(feed_name: str, url: str, query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        root = JirouCeng._rss_items(url, timeout=timeout)
        terms = JirouCeng._query_terms(query)
        results: list[dict[str, str]] = []
        for item in root.findall("./channel/item"):
            title = JirouCeng._clean_html_text(item.findtext("title"), limit=240)
            link = str(item.findtext("link") or item.findtext("guid") or "").strip()
            snippet = JirouCeng._clean_html_text(item.findtext("description"), limit=700)
            published = JirouCeng._clean_html_text(item.findtext("pubDate"), limit=120)
            haystack = f"{title} {snippet}".lower()
            if terms and not any(term in haystack for term in terms):
                continue
            if not title or not link.startswith("http"):
                continue
            if JirouCeng._is_stale_news_item(published, link):
                continue
            results.append({
                "title": title,
                "url": link,
                "domain": JirouCeng._domain(link),
                "published": published,
                "snippet": snippet,
                "source": "china_news_rss",
                "source_name": feed_name,
            })
            if len(results) >= max_results:
                break
        return results

    @staticmethod
    def _parse_news_datetime(published: str, url: str) -> datetime | None:
        text = str(published or "").strip()
        if text:
            try:
                value = parsedate_to_datetime(text)
                return value.replace(tzinfo=None)
            except Exception:
                pass
        raw_url = str(url or "")
        patterns = (
            r"/(20\d{2})/(\d{2})-(\d{2})/",
            r"/(20\d{2})-(\d{2})/(\d{2})/",
            r"/(20\d{2})/(\d{2})/(\d{2})/",
        )
        for pattern in patterns:
            match = re.search(pattern, raw_url)
            if not match:
                continue
            try:
                return datetime(int(match.group(1)), int(match.group(2)), int(match.group(3)))
            except Exception:
                return None
        return None

    @staticmethod
    def _is_stale_news_item(published: str, url: str, *, max_age_days: int = 14) -> bool:
        value = JirouCeng._parse_news_datetime(published, url)
        if value is None:
            return False
        return value < datetime.now() - timedelta(days=max_age_days)

    @staticmethod
    def _china_news_rss_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        results: list[dict[str, str]] = []
        errors: list[str] = []
        seen: set[str] = set()
        per_feed_timeout = max(3.0, min(6.0, float(timeout) / 2))
        for feed_name, url in _CHINA_NEWS_RSS_FEEDS:
            try:
                rows = JirouCeng._rss_feed_search(feed_name, url, query, max_results=max_results, timeout=per_feed_timeout)
            except Exception as exc:
                errors.append(f"{feed_name}: {exc}")
                continue
            for item in rows:
                key = item.get("url") or item.get("title") or ""
                if not key or key in seen:
                    continue
                seen.add(key)
                results.append(item)
                if len(results) >= max_results:
                    return results
        if not results and errors:
            raise RuntimeError("; ".join(errors[-4:]))
        return results

    @staticmethod
    def _global_news_rss_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """搜索国际新闻 RSS 源（Reuters, BBC, AP, NPR, Al Jazeera, CNN）"""
        results: list[dict[str, str]] = []
        seen: set[str] = set()
        per_feed_timeout = max(3.0, min(6.0, float(timeout) / 2))
        for feed_name, url in _GLOBAL_NEWS_RSS_FEEDS:
            try:
                rows = JirouCeng._rss_feed_search(feed_name, url, query, max_results=max_results, timeout=per_feed_timeout)
            except Exception:
                continue
            for item in rows:
                key = item.get("url") or item.get("title") or ""
                if not key or key in seen:
                    continue
                seen.add(key)
                results.append(item)
                if len(results) >= max_results:
                    return results
        return results

    @staticmethod
    def _bing_rss_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        params = {
            "format": "rss",
            "q": str(query or "").strip(),
            "setlang": "zh-CN",
            "mkt": "zh-CN",
            "cc": "CN",
        }
        url = "https://www.bing.com/search?" + urllib.parse.urlencode(params)
        root = JirouCeng._rss_items(url, timeout=timeout)
        results: list[dict[str, str]] = []
        for item in root.findall("./channel/item"):
            title = JirouCeng._clean_html_text(item.findtext("title"), limit=240)
            link = str(item.findtext("link") or "").strip()
            snippet = JirouCeng._clean_html_text(item.findtext("description"), limit=700)
            published = JirouCeng._clean_html_text(item.findtext("pubDate"), limit=120)
            if not title or not link.startswith("http"):
                continue
            results.append({
                "title": title,
                "url": link,
                "domain": JirouCeng._domain(link),
                "published": published,
                "snippet": snippet,
                "source": "bing_rss",
            })
            if len(results) >= max_results:
                break
        return results

    @staticmethod
    def _google_news_rss_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        search_query = str(query or "").strip()
        if "when:" not in search_query.lower():
            search_query = f"{search_query} when:{JirouCeng._news_window(search_query)}".strip()
        params = {
            "q": search_query,
            "hl": "zh-CN",
            "gl": "CN",
            "ceid": "CN:zh-Hans",
        }
        url = "https://news.google.com/rss/search?" + urllib.parse.urlencode(params)
        root = JirouCeng._rss_items(url, timeout=timeout)
        results: list[dict[str, str]] = []
        for item in root.findall("./channel/item"):
            title = JirouCeng._clean_html_text(item.findtext("title"), limit=240)
            link = str(item.findtext("link") or "").strip()
            snippet = JirouCeng._clean_html_text(item.findtext("description"), limit=700)
            published = JirouCeng._clean_html_text(item.findtext("pubDate"), limit=120)
            source_node = item.find("source")
            source_name = JirouCeng._clean_html_text(source_node.text if source_node is not None else "", limit=160)
            source_url = str(source_node.attrib.get("url") if source_node is not None else "").strip()
            if not title or not link.startswith("http"):
                continue
            results.append({
                "title": title,
                "url": link,
                "domain": JirouCeng._domain(source_url or link),
                "published": published,
                "snippet": snippet,
                "source": "google_news_rss",
                "source_name": source_name,
                "source_url": source_url,
            })
            if len(results) >= max_results:
                break
        return results

    @staticmethod
    def _search_results_to_content(query: str, results: list[dict[str, str]], provider: str) -> str:
        today = datetime.now().strftime("%Y-%m-%d")
        lines = [
            f"联网搜索证据包（生成日期：{today}）",
            f"检索问题：{query}",
            f"检索路径：{provider}",
            "",
            "关键来源：",
        ]
        for index, item in enumerate(results, start=1):
            lines.append(f"{index}. {item.get('title') or 'Untitled'}")
            if item.get("domain") or item.get("published"):
                lines.append(f"   来源：{item.get('domain') or 'unknown'}；时间：{item.get('published') or 'unknown'}")
            if item.get("source_name") or item.get("source_url"):
                lines.append(f"   原始来源：{item.get('source_name') or 'unknown'} {item.get('source_url') or ''}".rstrip())
            lines.append(f"   URL: {item.get('url') or ''}")
            if item.get("snippet"):
                lines.append(f"   摘要：{item.get('snippet')}")
        lines.append("")
        lines.append("这些内容来自外部网页，请交叉核验后再回答用户。")
        return "\n".join(lines)

    @staticmethod
    def _duckduckgo_instant_answer(query: str) -> dict:
        try:
            url = f"https://api.duckduckgo.com/?q={urllib.parse.quote(query)}&format=json&no_html=1&skip_disambig=1"
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = loads_json_object(resp.read().decode("utf-8", "replace"), source="duckduckgo_provider")

            jieguo = {
                "zhuangtai": "wancheng",
                "chaxun": query,
            }
            if data.get("AbstractText"):
                jieguo["zhaiyao"] = data["AbstractText"][:500]
                jieguo["laiyuan"] = data.get("AbstractURL", "")
            topics = data.get("RelatedTopics", [])
            if topics:
                jieguo["xiangguan"] = [
                    {"miaoshu": t.get("Text", "")[:200]}
                    for t in topics[:5] if isinstance(t, dict) and t.get("Text")
                ]
            if not jieguo.get("zhaiyao") and not jieguo.get("xiangguan"):
                jieguo["zhuangtai"] = "wu_jieguo"
            return jieguo
        except Exception as e:
            return {"zhuangtai": "cuowu", "chaxun": query, "cuowu": str(e)[:200]}

    @staticmethod
    def _duckduckgo_html_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """DuckDuckGo HTML 搜索 — 比 Bing RSS 更可靠的通用搜索"""
        try:
            url = "https://html.duckduckgo.com/html?q=" + urllib.parse.quote(str(query or "").strip())
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"})
            with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                html_text = resp.read().decode("utf-8", "replace")
        except Exception:
            return []

        results: list[dict[str, str]] = []
        blocks = re.split(r'<div\s+class="[^"]*result[^"]*"[^>]*>', html_text)
        for block in blocks[1:]:
            if len(results) >= max_results:
                break
            link_match = re.search(r'<a\s+[^>]*class="[^"]*result__a[^"]*"\s+href="([^"]+)"[^>]*>([^<]+)</a>', block, re.I)
            if not link_match:
                continue
            link = html.unescape(link_match.group(1).strip())
            title = html.unescape(re.sub(r'<[^>]+>', '', link_match.group(2)).strip())
            if not title or not link.startswith("http"):
                continue
            snippet_match = re.search(r'<a\s+[^>]*class="[^"]*result__snippet[^"]*"[^>]*>(.*?)</a>', block, re.I | re.DOTALL)
            snippet = html.unescape(re.sub(r'<[^>]+>', '', snippet_match.group(1).strip())) if snippet_match else ""
            domain = JirouCeng._domain(link)
            results.append({
                "title": title[:240],
                "url": link,
                "domain": domain,
                "snippet": snippet[:700],
                "published": "",
                "source": "duckduckgo_html",
            })
        return results

    @staticmethod
    def _bing_html_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """Bing HTML 搜索 — 解析搜索结果页的 h2+link 结构"""
        try:
            url = "https://www.bing.com/search?q=" + urllib.parse.quote(str(query or "").strip()) + "&setlang=zh-CN"
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"})
            with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                html_text = resp.read().decode("utf-8", "replace")
        except Exception:
            return []

        results: list[dict[str, str]] = []
        seen_urls: set[str] = set()
        # Bing 搜索结果: <h2>...<a href="url">title</a>...</h2> + <p class="b_caption">snippet</p>
        h2_matches = list(re.finditer(
            r'<h2[^>]*>.*?<a[^>]*href="(https?://[^"]+)"[^>]*>(.*?)</a>',
            html_text, re.I | re.DOTALL
        ))
        # Find snippets after each h2
        snippet_matches = list(re.finditer(
            r'<p[^>]*class="[^"]*b_caption[^"]*"[^>]*>(.*?)</p>',
            html_text, re.I | re.DOTALL
        ))
        snippets = [html.unescape(re.sub(r'<[^>]+>', '', m.group(1)).strip()) for m in snippet_matches]

        for i, m in enumerate(h2_matches):
            link = html.unescape(m.group(1).strip())
            title = html.unescape(re.sub(r'<[^>]+>', '', m.group(2)).strip())
            if not title or not link.startswith("http") or link in seen_urls:
                continue
            seen_urls.add(link)
            domain = JirouCeng._domain(link)
            snippet = snippets[i][:700] if i < len(snippets) else ""
            results.append({
                "title": title[:240],
                "url": link,
                "domain": domain,
                "snippet": snippet,
                "published": "",
                "source": "bing_html",
            })
            if len(results) >= max_results:
                break
        return results

    # ── 新增搜索引擎：百度 / 搜狗 / Google ──────────────────────────

    @staticmethod
    def _baidu_html_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """百度 HTML 搜索"""
        try:
            url = "https://www.baidu.com/s?wd=" + urllib.parse.quote(str(query or "").strip()) + "&rn=10"
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml",
                "Accept-Language": "zh-CN,zh;q=0.9",
                "Cookie": "BAIDUID=FAKE;",
            })
            with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                html_text = resp.read().decode("utf-8", "replace")
        except Exception:
            return []

        results: list[dict[str, str]] = []
        seen_urls: set[str] = set()
        blocks = re.split(r'<div[^>]*class="[^"]*(?:result|c-container)[^"]*"[^>]*>', html_text)
        for block in blocks[1:]:
            if len(results) >= max_results:
                break
            link = ""
            mu_match = re.search(r'data-mu="(https?://[^"]+)"', block, re.I)
            if mu_match:
                link = html.unescape(mu_match.group(1).strip())
            if not link:
                h3_match = re.search(r'<h3[^>]*>.*?<a[^>]*href="(https?://[^"]+)"', block, re.I | re.DOTALL)
                if h3_match:
                    link = html.unescape(h3_match.group(1).strip())
            if not link or not link.startswith("http") or link in seen_urls:
                continue
            title = ""
            title_match = re.search(r'<h3[^>]*>(.*?)</h3>', block, re.I | re.DOTALL)
            if title_match:
                title = html.unescape(re.sub(r'<[^>]+>', '', title_match.group(1)).strip())
            if not title:
                continue
            seen_urls.add(link)
            snippet = ""
            sn_match = re.search(r'<(?:span|div)[^>]*class="[^"]*(?:content-right|c-abstract|content-summary)[^"]*"[^>]*>(.*?)</(?:span|div)>', block, re.I | re.DOTALL)
            if sn_match:
                snippet = html.unescape(re.sub(r'<[^>]+>', '', sn_match.group(1)).strip())
            domain = JirouCeng._domain(link)
            results.append({
                "title": title[:240],
                "url": link,
                "domain": domain,
                "snippet": snippet[:700],
                "published": "",
                "source": "baidu_html",
            })
        return results

    @staticmethod
    def _sogou_html_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """搜狗 HTML 搜索"""
        try:
            url = "https://www.sogou.com/web?query=" + urllib.parse.quote(str(query or "").strip())
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml",
                "Accept-Language": "zh-CN,zh;q=0.9",
            })
            with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                html_text = resp.read().decode("utf-8", "replace")
        except Exception:
            return []

        results: list[dict[str, str]] = []
        seen_urls: set[str] = set()
        blocks = re.split(r'<div[^>]*class="[^"]*(?:vrwrap|rb)[^"]*"[^>]*>', html_text)
        for block in blocks[1:]:
            if len(results) >= max_results:
                break
            link_match = re.search(r'<a[^>]*href="(https?://[^"]+)"[^>]*id="[^"]*"', block, re.I)
            if not link_match:
                link_match = re.search(r'<a[^>]*href="(https?://[^"]+)"[^>]*>.*?</a>', block, re.I)
            if not link_match:
                continue
            link = html.unescape(link_match.group(1).strip())
            if not link.startswith("http") or link in seen_urls:
                continue
            title = html.unescape(re.sub(r'<[^>]+>', '', link_match.group(0)).strip())[:240]
            if not title:
                title_match = re.search(r'<h3[^>]*>(.*?)</h3>', block, re.I | re.DOTALL)
                title = html.unescape(re.sub(r'<[^>]+>', '', title_match.group(1)).strip())[:240] if title_match else ""
            if not title:
                continue
            seen_urls.add(link)
            snippet = ""
            sn_match = re.search(r'<(?:p|div)[^>]*class="[^"]*(?:str_info|space-txt|abstract)[^"]*"[^>]*>(.*?)</(?:p|div)>', block, re.I | re.DOTALL)
            if sn_match:
                snippet = html.unescape(re.sub(r'<[^>]+>', '', sn_match.group(1)).strip())[:700]
            domain = JirouCeng._domain(link)
            results.append({
                "title": title,
                "url": link,
                "domain": domain,
                "snippet": snippet,
                "published": "",
                "source": "sogou_html",
            })
        return results

    @staticmethod
    def _google_html_search(query: str, *, max_results: int, timeout: float) -> list[dict[str, str]]:
        """Google HTML 搜索 — 多镜像尝试"""
        mirrors = [
            "https://www.google.com/search?q=",
            "https://www.google.com.hk/search?q=",
        ]
        html_text = ""
        for base in mirrors:
            try:
                url = base + urllib.parse.quote(str(query or "").strip()) + "&hl=zh-CN"
                req = urllib.request.Request(url, headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                    "Accept": "text/html,application/xhtml+xml",
                    "Accept-Language": "zh-CN,zh;q=0.9",
                })
                with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                    html_text = resp.read().decode("utf-8", "replace")
                break
            except Exception:
                continue
        if not html_text:
            return []

        results: list[dict[str, str]] = []
        seen_urls: set[str] = set()
        blocks = re.split(r'<div[^>]*class="[^"]*g[^"]*"[^>]*>', html_text)
        for block in blocks[1:]:
            if len(results) >= max_results:
                break
            a_match = re.search(r'<a[^>]*href="(/url\\?q=)?(https?://[^"&]+)', block, re.I)
            if not a_match:
                a_match = re.search(r'<a[^>]*href="(https?://[^"]+)"', block, re.I)
            if not a_match:
                continue
            link = html.unescape(a_match.group(2) if a_match.lastindex and a_match.lastindex >= 2 else a_match.group(1)).strip()
            if not link.startswith("http") or link in seen_urls or "google" in link:
                continue
            seen_urls.add(link)
            title = ""
            h3_match = re.search(r'<h3[^>]*>(.*?)</h3>', block, re.I | re.DOTALL)
            if h3_match:
                title = html.unescape(re.sub(r'<[^>]+>', '', h3_match.group(1)).strip())[:240]
            if not title:
                continue
            snippet = ""
            sn_match = re.search(r'<(?:span|div)[^>]*class="[^"]*(?:st|VwiC3b)[^"]*"[^>]*>(.*?)</(?:span|div)>', block, re.I | re.DOTALL)
            if sn_match:
                snippet = html.unescape(re.sub(r'<[^>]+>', '', sn_match.group(1)).strip())[:700]
            domain = JirouCeng._domain(link)
            results.append({
                "title": title,
                "url": link,
                "domain": domain,
                "snippet": snippet,
                "published": "",
                "source": "google_html",
            })
        return results

    @staticmethod
    def _image_search(query: str, *, max_results: int = 10, timeout: float = 15) -> dict:
        """图片搜索 — Bing Images + 百度图片"""
        results: list[dict] = []
        errors: list[str] = []

        # Bing Images
        try:
            url = "https://www.bing.com/images/search?q=" + urllib.parse.quote(str(query or "").strip())
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            })
            with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                html_text = resp.read().decode("utf-8", "replace")
            for m in re.finditer(r'<a[^>]*class="[^"]*iusc[^"]*"[^>]*m=\'({[^}]+})\'', html_text, re.I):
                try:
                    meta = json.loads(html.unescape(m.group(1)))
                    img_url = meta.get("murl") or meta.get("turl") or ""
                    if img_url.startswith("http"):
                        results.append({
                            "url": img_url,
                            "thumbnail": meta.get("turl", ""),
                            "title": str(meta.get("t") or "")[:200],
                            "source": "bing_images",
                            "width": meta.get("w", 0),
                            "height": meta.get("h", 0),
                        })
                except Exception:
                    continue
                if len(results) >= max_results:
                    break
        except Exception as e:
            errors.append(f"bing_images: {e}")

        # 百度图片
        if len(results) < max_results:
            try:
                url = "https://image.baidu.com/search/flip?tn=baiduimage&word=" + urllib.parse.quote(str(query or "").strip())
                req = urllib.request.Request(url, headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                })
                with urllib.request.urlopen(req, timeout=max(timeout, 10)) as resp:
                    html_text = resp.read().decode("utf-8", "replace")
                for m in re.finditer(r'"objURL"\s*:\s*"(https?://[^"]+)"', html_text):
                    img_url = m.group(1)
                    if img_url not in {r["url"] for r in results}:
                        results.append({
                            "url": img_url,
                            "thumbnail": "",
                            "title": query,
                            "source": "baidu_images",
                            "width": 0,
                            "height": 0,
                        })
                    if len(results) >= max_results:
                        break
            except Exception as e:
                errors.append(f"baidu_images: {e}")

        if results:
            return {
                "zhuangtai": "wancheng",
                "chaxun": query,
                "tupian_shu": len(results),
                "tupian": results,
                "urls": [r["url"] for r in results],
            }
        return {"zhuangtai": "wu_jieguo", "chaxun": query, "cuowu": "; ".join(errors[-3:]) if errors else "未找到图片"}

    @staticmethod
    def _wangye_zhengwen(
        url: str = "",
        html_or_text: str = "",
        text: str = "",
        content: str = "",
        max_chars: int = 12000,
        timeout: int = 15,
    ) -> dict:
        """网页正文清洗 — 对齐 v2 的 web_readability_extract 基础能力"""
        source_text = str(html_or_text or text or content or "")
        requested_url = str(url or "").strip()
        final_url = requested_url
        if not requested_url and not source_text:
            return {"zhuangtai": "cuowu", "cuowu": "缺少要清洗的网页 URL 或正文"}

        max_chars = JirouCeng._clamp_int(max_chars, minimum=500, maximum=50000, default=12000)
        timeout = JirouCeng._clamp_int(timeout, minimum=3, maximum=60, default=15)
        content_type = "text/plain"

        try:
            if requested_url:
                try:
                    safe_url = _normalise_public_web_url(requested_url, resolve=True)
                except _WebUrlSecurityError as exc:
                    return _web_url_error_payload(requested_url, exc)
                req = urllib.request.Request(safe_url, headers=_SEARCH_HEADERS)
                with _public_web_opener().open(req, timeout=timeout) as resp:
                    final_url = _normalise_public_web_url(str(resp.geturl() or safe_url), resolve=True)
                    raw = resp.read(2_000_000)
                    content_type = str(resp.headers.get("content-type") or "")
                charset = "utf-8"
                for part in content_type.split(";"):
                    part = part.strip()
                    if part.lower().startswith("charset="):
                        charset = part.split("=", 1)[1].strip().strip('"').strip("'") or "utf-8"
                        break
                try:
                    source_text = raw.decode(charset, errors="replace")
                except LookupError:
                    source_text = raw.decode("utf-8", errors="replace")

            cleaned = re.sub(
                r"(?is)<(script|style|noscript|iframe|svg|canvas|nav|footer|header|aside)[^>]*>.*?</\1>",
                " ",
                source_text,
            )
            cleaned = re.sub(r"(?is)<!--.*?-->", " ", cleaned)
            cleaned = re.sub(r"(?is)<br\s*/?>", "\n", cleaned)
            cleaned = re.sub(r"(?is)</p\s*>|</div\s*>|</li\s*>|</h[1-6]\s*>", "\n", cleaned)
            plain = re.sub(r"(?is)<[^>]+>", " ", cleaned)
            plain = html.unescape(plain)
            plain = re.sub(r"[ \t\r\f\v]+", " ", plain)
            plain = re.sub(r"\n\s*\n\s*\n+", "\n\n", plain)
            plain = "\n".join(line.strip() for line in plain.splitlines() if line.strip())
            original_len = len(plain)
            blocked_marker = _web_content_block_reason(source_text, plain)
            if not plain or blocked_marker:
                return {
                    "ok": False,
                    "zhuangtai": "content_blocked",
                    "url": requested_url,
                    "final_url": final_url,
                    "content_type": content_type,
                    "text": plain[:500],
                    "zhengwen": "",
                    "zishu": 0,
                    "yuanwen_zishu": original_len,
                    "recoverable": False,
                    "error_code": "web_content_unavailable",
                    "cuowu": f"网页返回错误页或验证页：{blocked_marker or 'empty_content'}",
                }
            if original_len > max_chars:
                plain = plain[:max_chars] + f"\n\n...(截断，原文共 {original_len} 字符)"
            return {
                "ok": True,
                "zhuangtai": "wancheng",
                "url": requested_url,
                "final_url": final_url,
                "content_type": content_type,
                "text": plain,
                "zhengwen": plain,
                "zishu": len(plain),
                "yuanwen_zishu": original_len,
            }
        except _WebUrlSecurityError as exc:
            return _web_url_error_payload(requested_url, exc)
        except Exception as e:
            return _web_fetch_error_payload(requested_url, final_url, e)

    @staticmethod
    def _jiyisousuo(query: str) -> dict:
        """搜索自己的记忆 — 调用记忆引擎检索"""
        try:
            from ..jiyi.yinqing import JiyiYinqing
            from ..shenti_zhuangtai import ShentiZhuangtai
            jiyi = JiyiYinqing()
            st = ShentiZhuangtai()
            neirong = jiyi.jiansuo(st, query)
            return {
                "zhuangtai": "wancheng",
                "chaxun": query,
                "jieguo": neirong[:3000] if neirong else "（无匹配记忆）",
                "tiaoshu": len(neirong.split("\n")) if neirong else 0,
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "chaxun": query, "cuowu": str(e)[:200]}

    @staticmethod
    def _xuexi_liucheng(
        trigger: str = "user_request",
        topic: str = "",
        content: str = "",
        path: str = "",
        search_results: Any = None,
        skill_draft: Any = None,
        tool_blueprint: Any = None,
        allow_network: bool = True,
        release_skill_tool: bool = True,
        actor: str = "model",
    ) -> dict:
        """Run the v3 learning pipeline and connect it into the L0-L6 closed loop."""
        try:
            from ..jingyan_xuexi.xuexi_lian import XuexiLian

            def _search(query: str) -> dict:
                return JirouCeng._wangluosousuo(query, max_results=6, mode="general")

            available_tools: list[str] = []
            try:
                from .guge_ceng import GUGE

                available_tools = [
                    str(item.get("name") or "").strip()
                    for item in GUGE.suoyou_gongju()
                    if isinstance(item, dict) and str(item.get("name") or "").strip()
                ]
            except Exception:
                available_tools = []

            pipeline = XuexiLian(sousuo_han_shu=_search)
            result = pipeline.xuexi_neirong(
                trigger=trigger,
                topic=topic,
                content=content,
                path=path,
                search_results=search_results,
                skill_draft=skill_draft,
                tool_blueprint=tool_blueprint,
                actor=actor,
                allow_network=allow_network,
                release_skill_tool=release_skill_tool,
                available_tools=available_tools,
            )

            try:
                from ..jinhua.bihuan_yinqing import JinhuaBihuanYinqing

                closed_loop = JinhuaBihuanYinqing().yunxing(xiaoxi=topic or content[:200] or path, reason="learning_pipeline")
                result["l0_l6_closed_loop"] = closed_loop
                report = result.get("frontend_report") if isinstance(result.get("frontend_report"), dict) else {}
                report["l0L6ClosedLoop"] = closed_loop
                result["frontend_report"] = report
                if isinstance(result.get("ability"), dict):
                    result["ability"]["l0_l6_closed_loop"] = closed_loop
                    result["ability"]["frontend_report"] = report
                    rows = [result["ability"]]
                    if isinstance(result.get("tool_candidate"), dict):
                        rows.append(result["tool_candidate"])
                    try:
                        pipeline._zhuce_nengli_liebiao(rows)
                    except Exception:
                        pass
            except Exception as exc:
                result["l0_l6_closed_loop"] = {
                    "ok": False,
                    "error": str(exc)[:300],
                    "note": "学习能力已注册，但 L0-L6 闭环报告生成失败。",
                }
            return result
        except Exception as e:
            return {"zhuangtai": "cuowu", "leixing": "learning_pipeline", "cuowu": str(e)[:500]}

    @staticmethod
    def _xiazaiwenjian(url: str, target: str = "", overwrite: bool = False) -> dict:
        """从URL下载文件到工作区"""
        import urllib.request
        import urllib.error
        import os
        try:
            if not target:
                fname = url.rstrip("/").split("/")[-1] or "download"
                target = f"downloads/{fname}"
            p = _workspace_path(target)
            if p.exists() and not overwrite:
                return {"zhuangtai": "yicunzai", "lujing": str(p), "tishi": "文件已存在，设 overwrite=true 覆盖"}
            p.parent.mkdir(parents=True, exist_ok=True)
            req = urllib.request.Request(url, headers={"User-Agent": "Tiangong-v3/1.0"})
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = resp.read()
            p.write_bytes(data)
            return {
                "zhuangtai": "yixiazai",
                "lujing": str(p),
                "daxiao": len(data),
                "url": url,
            }
        except urllib.error.HTTPError as e:
            return {"zhuangtai": "cuowu", "url": url, "cuowu": f"HTTP {e.code}: {e.reason}"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "url": url, "cuowu": str(e)[:300]}

    @staticmethod
    def _liemulu(path: str = ".", dir_path: str = "", directory: str = "", folder: str = "", confirm: bool = False, **_ignored: Any) -> dict:
        """列出目录内容"""
        path = dir_path or directory or folder or path
        p = _workspace_path(path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "路径不存在"}
        if not p.is_dir():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "不是目录"}
        try:
            items = []
            for item in sorted(p.iterdir()):
                t = "dir" if item.is_dir() else "file"
                size = item.stat().st_size if item.is_file() else None
                entry = {"name": item.name, "type": t}
                if size is not None:
                    entry["size"] = size
                items.append(entry)
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "shuliang": len(items),
                "neirong": items[:100],
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)}

    @staticmethod
    def _yidong(source: str, target: str, overwrite: bool = False, confirm: bool = False, **_ignored: Any) -> dict:
        """移动文件或目录"""
        import shutil
        src = _workspace_path(source)
        dst = _workspace_path(target)
        if not src.exists():
            return {"zhuangtai": "cuowu", "source": str(src), "cuowu": "源路径不存在"}
        try:
            dst.parent.mkdir(parents=True, exist_ok=True)
            if dst.exists() and not overwrite:
                return {
                    "ok": False,
                    "zhuangtai": "target_exists",
                    "effect": "write",
                    "source": str(src),
                    "target": str(dst),
                    "updated_paths": [],
                    "readback": {"ok": False, "source_exists": src.exists(), "target_exists": True, "target_is_dir": dst.is_dir()},
                    "cuowu": "target already exists; move was not performed",
                    "error": "target already exists; move was not performed",
                    "note": "目标已存在，未覆盖，移动未执行。",
                }
            if dst.exists() and overwrite:
                if dst.is_dir():
                    shutil.rmtree(dst)
                else:
                    dst.unlink()
            shutil.move(str(src), str(dst))
            target_exists = dst.exists()
            source_exists = src.exists()
            ok = target_exists and not source_exists
            return {
                "ok": ok,
                "zhuangtai": "yiyidong" if ok else "cuowu",
                "effect": "write",
                "source": str(src),
                "target": str(dst),
                "updated_paths": [str(dst), str(src)],
                "readback": {"ok": ok, "target_exists": target_exists, "source_exists": source_exists, "is_dir": dst.is_dir() if target_exists else False},
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "source": str(src), "target": str(dst), "cuowu": str(e)}

    @staticmethod
    def _fuzhi(source: str, target: str, overwrite: bool = False, confirm: bool = False, **_ignored: Any) -> dict:
        """复制文件或目录"""
        import shutil
        src = _workspace_path(source)
        dst = _workspace_path(target)
        if not src.exists():
            return {"zhuangtai": "cuowu", "source": str(src), "cuowu": "源路径不存在"}
        try:
            dst.parent.mkdir(parents=True, exist_ok=True)
            if dst.exists() and not overwrite:
                return {
                    "ok": True,
                    "zhuangtai": "yicunzai",
                    "effect": "write",
                    "source": str(src),
                    "target": str(dst),
                    "updated_paths": [str(dst)],
                    "readback": {"ok": True, "exists": True, "is_dir": dst.is_dir()},
                    "note": "目标已存在，未覆盖；目标路径可用。",
                }
            if src.is_dir():
                shutil.copytree(src, dst, dirs_exist_ok=True)
            else:
                shutil.copy2(src, dst)
            return {
                "ok": dst.exists(),
                "zhuangtai": "yifuzhi" if dst.exists() else "cuowu",
                "effect": "write",
                "source": str(src),
                "target": str(dst),
                "updated_paths": [str(dst)],
                "readback": {"ok": dst.exists(), "is_dir": dst.is_dir() if dst.exists() else False},
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "source": str(src), "target": str(dst), "cuowu": str(e)}

    @staticmethod
    def _shanchu(path: str, recursive: bool = False, confirm: bool = False, **_ignored: Any) -> dict:
        """删除文件或目录"""
        import shutil
        p = _workspace_path(path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "路径不存在"}
        try:
            if p.is_dir():
                if not recursive:
                    return {"zhuangtai": "jujue", "lujing": str(p), "tishi": "是目录，设 recursive=true 递归删除"}
                shutil.rmtree(p)
            else:
                p.unlink()
            return {
                "ok": not p.exists(),
                "zhuangtai": "yishanchu" if not p.exists() else "cuowu",
                "effect": "write",
                "lujing": str(p),
                "updated_paths": [str(p)],
                "readback": {"ok": not p.exists(), "exists": p.exists()},
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)}

    @staticmethod
    def _gongzuoqu_sousuo(query: str, path: str = ".") -> dict:
        """在工作区文件中搜索文本（grep）"""
        try:
            p = _workspace_path(path)
            if not p.exists():
                return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "路径不存在"}
            jieguo = []
            for f in p.rglob("*"):
                if not f.is_file():
                    continue
                try:
                    text = f.read_text(encoding="utf-8", errors="ignore")
                except Exception:
                    continue
                for i, line in enumerate(text.split("\n"), 1):
                    if query in line:
                        jieguo.append(f"{f}:{i}:{line.strip()[:200]}")
                        if len(jieguo) >= 50:
                            break
                if len(jieguo) >= 50:
                    break
            return {
                "zhuangtai": "wancheng",
                "chaxun": query,
                "lujing": str(p),
                "pipei_shu": len(jieguo),
                "jieguo": jieguo,
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "chaxun": query, "cuowu": str(e)[:200]}

    @staticmethod
    def _wendangjiexi(path: str) -> dict:
        """解析文档内容，支持 txt/md/csv/json/docx/xlsx/pptx/pdf"""
        p = _workspace_path(path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        suffix = p.suffix.lower()
        try:
            if suffix in {".txt", ".md", ".py", ".yaml", ".yml", ".toml", ".log"}:
                text = p.read_text(encoding="utf-8", errors="ignore")
                return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "text", "neirong": text[:10000], "zishu": len(text)}
            if suffix == ".csv":
                import csv
                with open(p, encoding="utf-8", errors="ignore") as f:
                    rows = list(csv.reader(f))[:200]
                return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "csv", "hangshu": len(rows), "neirong": str(rows)[:10000]}
            if suffix == ".json":
                raw_json = p.read_text(encoding="utf-8")
                try:
                    data = json.loads(raw_json)
                except Exception as exc:
                    detail = error_payload(exc, source="document_json", ok_key=False)
                    return {
                        "zhuangtai": "cuowu",
                        "lujing": str(p),
                        "cuowu": detail.get("error", str(exc)),
                        "error_code": detail.get("error_code", type(exc).__name__),
                        "detail": detail.get("detail", str(exc)),
                    }
                text = json.dumps(data, ensure_ascii=False, indent=2)
                return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "json", "neirong": text[:10000]}
            if suffix == ".docx":
                try:
                    from docx import Document
                    doc = Document(str(p))
                    text = "\n".join(para.text for para in doc.paragraphs)
                    return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "docx", "neirong": text[:10000], "zishu": len(text)}
                except ImportError:
                    return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 python-docx"}
            if suffix == ".xlsx":
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(str(p), data_only=True)
                    sheets = {}
                    for name in wb.sheetnames:
                        ws = wb[name]
                        rows = [[str(c.value) if c.value is not None else "" for c in row] for row in ws.iter_rows()][:100]
                        sheets[name] = rows
                    return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "xlsx", "biaoge": sheets}
                except ImportError:
                    return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 openpyxl"}
            if suffix == ".pdf":
                try:
                    import pdfplumber
                    with pdfplumber.open(str(p)) as pdf:
                        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
                    return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "pdf", "neirong": text[:10000], "yeshu": len(pdf.pages)}
                except ImportError:
                    return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 pdfplumber"}
            if suffix == ".pptx":
                try:
                    from pptx import Presentation
                    prs = Presentation(str(p))
                    slides = []
                    for slide in prs.slides:
                        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                        slides.append(" | ".join(texts))
                    return {"zhuangtai": "wancheng", "lujing": str(p), "leixing": "pptx", "huandengpian": slides[:50]}
                except ImportError:
                    return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 python-pptx"}
            return {"zhuangtai": "buzhichi", "lujing": str(p), "leixing": suffix, "tishi": f"不支持的格式 {suffix}，尝试用 terminal 调用外部工具"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _httpqingqiu(url: str, method: str = "GET", headers: dict = None, body: str = "", timeout: int = 15) -> dict:
        """发送经过公网地址固定和标准证书校验的 HTTP 请求。"""
        requested_url = str(url or "").strip()
        try:
            safe_url = _normalise_public_web_url(requested_url, resolve=True)
        except _WebUrlSecurityError as exc:
            return _web_url_error_payload(requested_url, exc)
        hdrs = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
                "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8"}
        if headers:
            hdrs.update(headers)
        data = body.encode("utf-8") if body else None
        try:
            req = urllib.request.Request(safe_url, data=data, headers=hdrs, method=method.upper())
            with _public_web_opener().open(req, timeout=timeout) as resp:
                final_url = _normalise_public_web_url(str(resp.geturl() or safe_url), resolve=True)
                raw = resp.read(2_000_001)
                truncated = len(raw) > 2_000_000
                raw = raw[:2_000_000]
                rbody = raw.decode("utf-8", errors="replace")
            return {
                "zhuangtai": "wancheng", "url": requested_url, "final_url": final_url, "method": method.upper(),
                "status": resp.status, "headers": dict(resp.headers),
                "body": rbody[:8000], "daxiao": len(raw), "truncated": truncated,
            }
        except urllib.error.HTTPError as e:
            try:
                err_body = e.read().decode("utf-8", errors="replace")[:2000]
            except Exception:
                err_body = ""
            return {"ok": False, "zhuangtai": "cuowu", "url": requested_url, "status": e.code, "cuowu": str(e.reason), "body": err_body}
        except _WebUrlSecurityError as exc:
            return _web_url_error_payload(requested_url, exc)
        except Exception as e:
            result = _web_fetch_error_payload(requested_url, requested_url, e)
            result["method"] = method.upper()
            return result

    @staticmethod
    def _dnsjiexi(host: str) -> dict:
        """DNS解析"""
        import socket
        try:
            ips = socket.getaddrinfo(host, None)
            v4 = list(set(a[4][0] for a in ips if a[0] == socket.AF_INET))
            v6 = list(set(a[4][0] for a in ips if a[0] == socket.AF_INET6))
            return {
                "zhuangtai": "wancheng",
                "host": host,
                "ipv4": v4,
                "ipv6": v6,
            }
        except socket.gaierror as e:
            return {"zhuangtai": "cuowu", "host": host, "cuowu": f"解析失败: {e}"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "host": host, "cuowu": str(e)[:200]}

    @staticmethod
    def _xiangmusaomiao(path: str = ".") -> dict:
        """扫描项目结构"""
        p = _workspace_path(path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "路径不存在"}
        try:
            py_files = []
            other_files = []
            dirs = []
            for item in sorted(p.rglob("*")):
                if item.name.startswith(".") or "__pycache__" in str(item):
                    continue
                rel = str(item.relative_to(p))
                if item.is_dir():
                    dirs.append(rel)
                elif item.suffix == ".py":
                    py_files.append(rel)
                else:
                    other_files.append(rel)
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "mulu_shu": len(dirs),
                "py_wenjian_shu": len(py_files),
                "qita_wenjian_shu": len(other_files),
                "mulu": dirs[:30],
                "py_wenjian": py_files[:50],
                "qita_wenjian": other_files[:30],
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:200]}

    # ═══════════════════════════════════════════
    # 多媒体
    # ═══════════════════════════════════════════

    @staticmethod
    def _tupian_shengcheng(prompt: str, size: str = "1024x1024", style: str = "", output_name: str = "") -> dict:
        """图片生成。支持 OpenAI-compatible /images/generations。"""
        return _image_generation_call(prompt=prompt, size=size, style=style, output_name=output_name)

    @staticmethod
    def _shipin_shengcheng(
        prompt: str,
        duration: int = 6,
        resolution: str = "768P",
        model: str = "",
        first_frame_image: str = "",
        last_frame_image: str = "",
        output_name: str = "",
        prompt_optimizer: bool = True,
        poll_interval: int = 10,
        max_wait: int = 600,
    ) -> dict:
        """视频生成。MiniMax 使用异步 video_generation 任务。"""
        return _video_generation_call(
            prompt=prompt,
            duration=duration,
            resolution=resolution,
            model=model,
            first_frame_image=first_frame_image,
            last_frame_image=last_frame_image,
            output_name=output_name,
            prompt_optimizer=prompt_optimizer,
            poll_interval=poll_interval,
            max_wait=max_wait,
        )

    @staticmethod
    def _tupianjiance(image_path: str, question: str = "") -> dict:
        """识别图片内容。优先走视觉模型；未配置时稳定返回本地元数据。"""
        p = _workspace_path(image_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        info = _image_metadata(p)
        prompt = question or "请用中文简洁描述这张图片的主要内容、可见文字、关键物体和可能的用途。"
        vision = _vision_image_call(p, prompt, purpose="general", max_tokens=2048)
        if vision.get("ok"):
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "xinxi": info,
                "neirong": vision.get("text", ""),
                "miaoshu": vision.get("text", ""),
                "vision_state": "ok",
                "provider": vision.get("provider", ""),
                "model": vision.get("model", ""),
                "vision_views": vision.get("vision_views", []),
                "question": question,
            }
        return {
            "zhuangtai": "wancheng",
            "lujing": str(p),
            "xinxi": info,
            "vision_state": vision.get("state", "unavailable"),
            "vision_error": vision.get("error", ""),
            "tishi": "未配置可用视觉模型，当前返回图片元数据；如需真正看图，请配置支持 image_url 的 OpenAI 兼容视觉模型。",
            "question": question,
        }

    @staticmethod
    def _tupian_ocr(image_path: str, language_hint: str = "auto") -> dict:
        """图片文字提取"""
        p = _workspace_path(image_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        tesseract_error = ""
        try:
            from PIL import Image
            img = Image.open(str(p))
            try:
                import pytesseract
                lang = "chi_sim+eng" if language_hint in ("auto", "zh", "mixed") else "eng"
                text = pytesseract.image_to_string(img, lang=lang)
                if str(text or "").strip():
                    return {"zhuangtai": "wancheng", "lujing": str(p), "wenzi": text[:5000], "yuyan": lang}
                tesseract_error = "tesseract returned empty text"
            except ImportError as e:
                tesseract_error = f"需要安装 pytesseract 和 tesseract-ocr: {e}"
            except Exception as e:
                tesseract_error = str(e)[:300]
        except ImportError:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "需要安装 Pillow"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:200]}
        vision = _vision_image_call(
            p,
            "Extract all visible text from the image. Preserve reading order and line breaks. "
            "Use enhanced/cropped views to verify small or low-contrast text. "
            "Return only the extracted text; return an empty string if no text is visible.",
            purpose="ocr",
            max_tokens=4096,
        )
        if vision.get("ok"):
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "wenzi": str(vision.get("text") or "")[:5000],
                "yuyan": language_hint,
                "vision_state": "ok",
                "fallback": "vision_model",
                "tesseract_error": tesseract_error,
                "vision_views": vision.get("vision_views", []),
            }
        return {
            "zhuangtai": "cuowu",
            "lujing": str(p),
            "cuowu": tesseract_error or "OCR 引擎不可用",
            "vision_state": vision.get("state", "unavailable"),
            "vision_error": vision.get("error", ""),
        }

    @staticmethod
    def _tupian_biaoge(image_path: str) -> dict:
        """提取图片中的表格"""
        p = _workspace_path(image_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        vision = _vision_image_call(
            p,
            "Extract tables from the image. Prefer Markdown tables. "
            "Use enhanced/cropped views to verify row and column boundaries. "
            "If no table is present, say no table was found.",
            purpose="table",
            max_tokens=4096,
        )
        if vision.get("ok"):
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "biaoge": str(vision.get("text") or "")[:12000],
                "vision_state": "ok",
                "provider": vision.get("provider", ""),
                "model": vision.get("model", ""),
                "vision_views": vision.get("vision_views", []),
            }
        return {
            "zhuangtai": "weishixian",
            "lujing": str(p),
            "xinxi": _image_metadata(p),
            "vision_state": vision.get("state", "unavailable"),
            "vision_error": vision.get("error", ""),
            "tishi": "未配置可用视觉模型，暂不能做图片表格识别。",
        }

    @staticmethod
    def _shipinjiance(video_path: str, question: str = "") -> dict:
        """分析视频内容"""
        p = _workspace_path(video_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        try:
            jieguo = subprocess.run(
                [_bundled_binary("ffprobe"), "-v", "quiet", "-print_format", "json", "-show_format", "-show_streams", str(p)],
                capture_output=True, text=True, timeout=15
            )
            if jieguo.returncode != 0:
                return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "ffprobe 无法解析此文件"}
            data = loads_json_object(jieguo.stdout, source="ffprobe")
            fmt = data.get("format", {})
            streams = data.get("streams", [])
            video_streams = [s for s in streams if s.get("codec_type") == "video"]
            audio_streams = [s for s in streams if s.get("codec_type") == "audio"]
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "shichang": fmt.get("duration"),
                "daxiao": fmt.get("size"),
                "geshi": fmt.get("format_name"),
                "shipin_liu": [{"bianma": s.get("codec_name"), "fenbianlv": f"{s.get('width')}x{s.get('height')}", "zhenlv": s.get("r_frame_rate")} for s in video_streams],
                "yinpin_liu": [{"bianma": s.get("codec_name"), "caiyanglv": s.get("sample_rate"), "shengdao": s.get("channels")} for s in audio_streams],
                "question": question,
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:200]}

    @staticmethod
    def _shipin_guanjianzhen(video_path: str, interval_sec: int = 5, max_frames: int = 20) -> dict:
        """提取视频关键帧"""
        p = _workspace_path(video_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        import tempfile
        tmpdir = Path(tempfile.mkdtemp())
        try:
            cmd = [
                _bundled_binary("ffmpeg"), "-i", str(p), "-vf", f"fps=1/{interval_sec}",
                "-frames:v", str(max_frames), "-q:v", "2",
                f"{tmpdir}/frame_%03d.jpg", "-y"
            ]
            jg = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            frames = sorted(tmpdir.glob("*.jpg"))
            return {
                "zhuangtai": "wancheng" if frames else "cuowu",
                "lujing": str(p),
                "zhen_shu": len(frames),
                "zhen_mulu": str(tmpdir),
                "wenjian": [str(f) for f in frames],
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _shipin_changjing(video_path: str) -> dict:
        """视频场景切分"""
        p = _workspace_path(video_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        import tempfile
        tmpdir = Path(tempfile.mkdtemp())
        try:
            cmd = [
                _bundled_binary("ffmpeg"), "-i", str(p),
                "-filter:v", "select='gt(scene,0.3)',showinfo",
                "-f", "null", "-", "-y"
            ]
            jg = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            # 解析 showinfo 输出找场景切换时间点
            timestamps = []
            for line in jg.stderr.split("\n"):
                if "pts_time:" in line:
                    try:
                        ts = line.split("pts_time:")[1].strip().split()[0]
                        timestamps.append(float(ts))
                    except Exception:
                        pass
            return {
                "zhuangtai": "wancheng",
                "lujing": str(p),
                "changjing_shu": len(timestamps),
                "shijian_dian": timestamps[:30],
                "tishi": "场景切换时间点已列出，可用 video_trim 按时间点裁剪",
            }
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _shipin_zimu(video_path: str) -> dict:
        """提取视频字幕"""
        p = _workspace_path(video_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        try:
            # 先看有哪些字幕流
            jg = subprocess.run(
                [_bundled_binary("ffprobe"), "-v", "quiet", "-print_format", "json", "-show_streams", str(p)],
                capture_output=True, text=True, timeout=10
            )
            sub_streams = []
            if jg.returncode == 0:
                data = loads_json_object(jg.stdout, source="ffprobe")
                sub_streams = [s for s in data.get("streams", []) if s.get("codec_type") == "subtitle"]
            if not sub_streams:
                return {"zhuangtai": "wu_zimu", "lujing": str(p), "tishi": "视频中没有内嵌字幕流"}
            # 提取第一个字幕流
            import tempfile
            tmp = Path(tempfile.mkdtemp()) / "subtitle.srt"
            cmd = [_bundled_binary("ffmpeg"), "-i", str(p), "-map", f"0:s:0", str(tmp), "-y"]
            jg2 = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            if jg2.returncode == 0 and tmp.exists():
                text = tmp.read_text(encoding="utf-8", errors="ignore")
                return {"zhuangtai": "wancheng", "lujing": str(p), "zimu_wenjian": str(tmp), "neirong": text[:5000]}
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "字幕提取失败"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _yuyin_zhuanwenzi(audio_path: str, language_hint: str = "auto") -> dict:
        """语音转文字 — 用 OpenAI whisper-1 API"""
        p = _workspace_path(audio_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        try:
            from ..peizhi import duqu_api_miyao, PROVIDER_BASE_URL
            import urllib.request
            import json as _json

            miyao = duqu_api_miyao("openai")
            if not miyao:
                return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "未配置 openai API密钥，请将密钥写入 ~/.tiangong/api_keys.json 的 openai 字段"}

            base = PROVIDER_BASE_URL.get("openai", "https://api.openai.com/v1")
            url = f"{base}/audio/transcriptions"

            # 读文件内容，构建 multipart
            boundary = "----TiangongAudioBoundary"
            body = (
                f"--{boundary}\r\n"
                f'Content-Disposition: form-data; name="model"\r\n\r\nwhisper-1\r\n'
                f"--{boundary}\r\n"
                f'Content-Disposition: form-data; name="file"; filename="{p.name}"\r\n'
                f"Content-Type: application/octet-stream\r\n\r\n"
            ).encode("utf-8")
            body += p.read_bytes()
            lang = language_hint if language_hint != "auto" else "zh"
            body += f"\r\n--{boundary}\r\nContent-Disposition: form-data; name=\"language\"\r\n\r\n{lang}\r\n--{boundary}--\r\n".encode("utf-8")

            req = urllib.request.Request(
                url, data=body,
                headers={
                    "Authorization": f"Bearer {miyao}",
                    "Content-Type": f"multipart/form-data; boundary={boundary}",
                },
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = loads_json_object(resp.read().decode("utf-8", errors="replace"), source="whisper_provider")
            return {"zhuangtai": "wancheng", "lujing": str(p), "wenzi": data.get("text", "")[:10000], "yuyan": lang}
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}

    @staticmethod
    def _yuyin_shuohuaren(audio_path: str) -> dict:
        """说话人分离"""
        p = _workspace_path(audio_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        return {"zhuangtai": "weishixian", "lujing": str(p), "tishi": "说话人分离需要 pyannote-audio 等模型。当前仅返回占位。"}

    @staticmethod
    def _yuyin_zhaiyao(audio_path: str, transcript: str = "") -> dict:
        """音频内容摘要"""
        p = _workspace_path(audio_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        if transcript:
            return {"zhuangtai": "wancheng", "lujing": str(p), "zhaiyao_tishi": "已提供文本，LLM可直接总结", "wenben": transcript[:5000]}
        return {"zhuangtai": "weishixian", "lujing": str(p), "tishi": "请先调用 audio_transcribe 获取文本，再传入 transcript 参数"}

    @staticmethod
    def _wenzi_zhuanyuyin(text: str, voice: str = "zh-CN-XiaoxiaoNeural", output_name: str = "") -> dict:
        """文字转语音 — 用 edge-tts"""
        if not text:
            return {"zhuangtai": "cuowu", "cuowu": "text 为空"}
        try:
            from edge_tts import Communicate
            import asyncio
            out = output_name or "tts_output.mp3"
            p = _workspace_path(out)
            async def _hecheng():
                comm = Communicate(text, voice)
                await comm.save(str(p))
            asyncio.run(_hecheng())
            return {"zhuangtai": "wancheng", "lujing": str(p), "daxiao": p.stat().st_size, "zishu": len(text), "voice": voice}
        except ImportError:
            return {"zhuangtai": "cuowu", "cuowu": "需要安装 edge-tts (pip install edge-tts)"}
        except Exception as e:
            return {"zhuangtai": "cuowu", "cuowu": str(e)[:300]}

    @staticmethod
    def _yuyin_jiangzao(audio_path: str) -> dict:
        """音频降噪 — ffmpeg anlmdn"""
        p = _workspace_path(audio_path)
        if not p.exists():
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": "文件不存在"}
        try:
            out_path = p.parent / f"{p.stem}_denoised{p.suffix}"
            cmd = [_bundled_binary("ffmpeg"), "-i", str(p), "-af", "anlmdn", str(out_path), "-y"]
            jg = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if jg.returncode == 0:
                return {"zhuangtai": "wancheng", "lujing": str(p), "shuchu": str(out_path)}
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": jg.stderr[-300:]}
        except Exception as e:
            return {"zhuangtai": "cuowu", "lujing": str(p), "cuowu": str(e)[:300]}


# 全局单例
JIROU = JirouCeng()
