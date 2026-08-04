from __future__ import annotations

import base64
import csv
import hashlib
import html
import json
import mimetypes
import os
import re
import shutil
import time
from datetime import datetime
from pathlib import Path
from typing import Any, Callable


TEXT_SUFFIXES = {
    ".txt", ".md", ".markdown", ".log", ".csv", ".json", ".jsonl",
    ".html", ".htm", ".xml", ".yaml", ".yml", ".toml", ".py", ".js",
    ".ts", ".tsx", ".jsx", ".css", ".scss", ".less", ".vue", ".svelte",
    ".java", ".c", ".cc", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs",
    ".php", ".rb", ".swift", ".kt", ".kts", ".sql", ".ini", ".conf",
}
DOCUMENT_SUFFIXES = TEXT_SUFFIXES | {".docx", ".xlsx", ".pptx", ".pdf"}
MEDIA_SUFFIXES = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".ico", ".svg",
    ".avif", ".tif", ".tiff", ".mp4", ".webm", ".ogv", ".mov", ".mkv",
    ".avi", ".m4v", ".wmv", ".flv", ".mpeg", ".mpg", ".3gp", ".ts",
    ".m2ts", ".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac", ".opus",
    ".wma",
}
ARCHIVE_SUFFIXES = {".zip", ".rar", ".7z", ".tar", ".gz", ".tgz", ".bz2", ".xz", ".zst", ".iso", ".dmg"}
APPLICATION_SUFFIXES = {
    ".exe", ".msi", ".msix", ".msixbundle", ".appx", ".appxbundle", ".apk",
    ".ipa", ".bat", ".cmd", ".ps1", ".vbs", ".lnk", ".jar",
}

MAX_INLINE_BYTES = int(os.environ.get("TIANGONG_KNOWLEDGE_INLINE_MAX_BYTES", 25 * 1024 * 1024))
MAX_TRANSFER_BYTES = int(os.environ.get("TIANGONG_FILE_TRANSFER_MAX_BYTES", 200 * 1024 * 1024))
MAX_PARSE_BYTES = int(os.environ.get("TIANGONG_KNOWLEDGE_PARSE_MAX_BYTES", 80 * 1024 * 1024))
DEFAULT_MAX_CHARS = int(os.environ.get("TIANGONG_KNOWLEDGE_MAX_CHARS", 240_000))
CHUNK_SIZE = int(os.environ.get("TIANGONG_KNOWLEDGE_CHUNK_SIZE", 1400))
CHUNK_OVERLAP = int(os.environ.get("TIANGONG_KNOWLEDGE_CHUNK_OVERLAP", 180))
CARD_INPUT_MAX_CHARS = int(os.environ.get("TIANGONG_KNOWLEDGE_CARD_INPUT_MAX_CHARS", 24_000))
_CARD_ENRICHER: Callable[[dict[str, Any]], dict[str, Any]] | None = None


def set_card_enricher(enricher: Callable[[dict[str, Any]], dict[str, Any]] | None) -> None:
    """Bind the runtime-owned LLM bridge without coupling storage to a provider."""
    global _CARD_ENRICHER
    if enricher is not None and not callable(enricher):
        raise TypeError("knowledge card enricher must be callable")
    _CARD_ENRICHER = enricher


def _workspace_root() -> Path:
    try:
        from .workspace_settings import duqu_workspace_root

        return duqu_workspace_root()
    except Exception:
        raw = (
            os.environ.get("TIANGONG_DESKTOP_WORKSPACE_ROOT")
            or os.environ.get("TIANGONG_WORKSPACE_ROOT")
            or os.getcwd()
        )
        return Path(raw).expanduser().resolve(strict=False)


def _workspace_path(value: Any) -> Path:
    path = Path(str(value or "")).expanduser()
    if not path.is_absolute():
        path = _workspace_root() / path
    return path.resolve(strict=False)


def _knowledge_config_path() -> Path:
    state_dir = os.environ.get("TIANGONG_DESKTOP_STATE_DIR")
    if state_dir:
        return Path(state_dir).expanduser().parent / "knowledge-settings.json"
    return Path.home() / ".tiangong" / "v3" / "knowledge-settings.json"


def _configured_knowledge_root() -> str:
    path = _knowledge_config_path()
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return str(data.get("knowledge_root") or "").strip() if isinstance(data, dict) else ""
    except (FileNotFoundError, OSError, ValueError, TypeError):
        return ""


def knowledge_settings() -> dict[str, Any]:
    """Return the durable knowledge-root projection without scanning documents."""
    path = _knowledge_config_path()
    root = ""
    updated_at = ""
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            root = str(data.get("knowledge_root") or "").strip()
            updated_at = str(data.get("updated_at") or "").strip()
    except (FileNotFoundError, OSError, ValueError, TypeError):
        pass
    return {
        "ok": True,
        "schema": "tiangong.v3.knowledge.settings.v1",
        "knowledgeRoot": root,
        "configured": bool(root),
        "updated_at": updated_at,
    }


def configure_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    """Persist the user-selected root as the backend retrieval authority."""
    payload = payload if isinstance(payload, dict) else {}
    raw = payload.get("knowledgeRoot") or payload.get("knowledge_root")
    if not raw:
        return {"ok": False, "error": "missing_knowledge_root"}
    root = Path(str(raw)).expanduser()
    if not root.is_absolute():
        root = _workspace_root() / root
    root = root.resolve(strict=False)
    if root.exists() and not root.is_dir():
        return {"ok": False, "error": "knowledge_root_not_directory"}
    root.mkdir(parents=True, exist_ok=True)
    config_path = _knowledge_config_path()
    config_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path = config_path.with_suffix(".tmp")
    temp_path.write_text(
        json.dumps(
            {"schema": "tiangong.v3.knowledge.settings.v1", "knowledge_root": str(root), "updated_at": _now_iso()},
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    temp_path.replace(config_path)
    listing = knowledge_list({"knowledgeRoot": str(root)})
    return {**listing, "configured": True, "config_path": str(config_path)}


def knowledge_root(payload: dict[str, Any] | None = None) -> Path:
    payload = payload if isinstance(payload, dict) else {}
    explicit = (
        payload.get("knowledgeRoot")
        or payload.get("knowledge_root")
        or os.environ.get("TIANGONG_KNOWLEDGE_DIR")
        or os.environ.get("TIANGONG_DESKTOP_KNOWLEDGE_ROOT")
        or _configured_knowledge_root()
    )
    if explicit:
        root = Path(str(explicit)).expanduser()
    else:
        state_dir = os.environ.get("TIANGONG_DESKTOP_STATE_DIR")
        if state_dir:
            root = Path(state_dir).expanduser().parent / "knowledge"
        else:
            root = Path.home() / ".tiangong" / "v3" / "knowledge"
    root.mkdir(parents=True, exist_ok=True)
    _contexts_dir(root).mkdir(parents=True, exist_ok=True)
    _files_dir(root).mkdir(parents=True, exist_ok=True)
    _exports_dir(root).mkdir(parents=True, exist_ok=True)
    _handoffs_dir(root).mkdir(parents=True, exist_ok=True)
    return root.resolve()


def _contexts_dir(root: Path) -> Path:
    return root / "contexts"


def _files_dir(root: Path) -> Path:
    return root / "files"


def _exports_dir(root: Path) -> Path:
    return root / "exports"


def _handoffs_dir(root: Path) -> Path:
    return root / "handoffs"


def _index_path(root: Path) -> Path:
    return root / "index.json"


def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _safe_name(value: Any, fallback: str = "file") -> str:
    name = Path(str(value or fallback)).name
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name).strip(" .")
    return (name or fallback)[:160]


def _safe_text(value: Any, limit: int = 2000) -> str:
    text = str(value or "")
    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit]


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _load_index(root: Path) -> dict[str, Any]:
    path = _index_path(root)
    if not path.exists():
        return {"schema": "tiangong.v3.knowledge.index.v1", "updated_at": "", "documents": {}, "last_document_id": ""}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            data.setdefault("schema", "tiangong.v3.knowledge.index.v1")
            data.setdefault("documents", {})
            data.setdefault("last_document_id", "")
            return data
    except Exception:
        pass
    return {"schema": "tiangong.v3.knowledge.index.v1", "updated_at": "", "documents": {}, "last_document_id": ""}


def _save_index(root: Path, index: dict[str, Any]) -> None:
    index["updated_at"] = _now_iso()
    tmp = _index_path(root).with_suffix(".json.tmp")
    tmp.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(_index_path(root))


def _context_path(root: Path, document_id: str) -> Path:
    clean = re.sub(r"[^a-zA-Z0-9_.-]", "_", str(document_id or ""))[:100] or "document"
    return _contexts_dir(root) / f"{clean}.json"


def _load_context(root: Path, document_id: str) -> dict[str, Any] | None:
    path = _context_path(root, document_id)
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else None
    except Exception:
        return None


def _save_context(root: Path, ctx: dict[str, Any]) -> None:
    path = _context_path(root, str(ctx.get("document_id") or "document"))
    tmp = path.with_suffix(".json.tmp")
    tmp.write_text(json.dumps(ctx, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(path)


def _kind_from_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in DOCUMENT_SUFFIXES:
        return "document"
    if suffix in MEDIA_SUFFIXES:
        if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".ico", ".svg", ".avif", ".tif", ".tiff"}:
            return "image"
        if suffix in {".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac", ".opus", ".wma"}:
            return "audio"
        return "video"
    if suffix in ARCHIVE_SUFFIXES:
        return "archive"
    if suffix in APPLICATION_SUFFIXES:
        return "application"
    return "file"


def _ext_from_mime(mime_type: str, fallback: str = ".bin") -> str:
    mime = str(mime_type or "").lower()
    if "jpeg" in mime or "jpg" in mime:
        return ".jpg"
    guessed = mimetypes.guess_extension(mime.split(";")[0].strip()) if mime else ""
    return guessed or fallback


def _decode_data_url(item: dict[str, Any], target_dir: Path) -> Path | None:
    data_url = str(item.get("dataUrl") or item.get("data_url") or "")
    match = re.match(r"^data:([^;,]+)?(?:;[^,]*)?;base64,(.+)$", data_url, re.S)
    if not match:
        return None
    mime_type = str(item.get("type") or match.group(1) or "application/octet-stream")
    raw_b64 = re.sub(r"\s+", "", match.group(2) or "")
    data = base64.b64decode(raw_b64, validate=False)
    if not data or len(data) > MAX_INLINE_BYTES:
        return None
    raw_name = _safe_name(item.get("name") or "pasted_file")
    if not Path(raw_name).suffix:
        raw_name += _ext_from_mime(mime_type)
    stem = Path(raw_name).stem[:80] or "pasted_file"
    suffix = Path(raw_name).suffix[:16] or ".bin"
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / f"{int(time.time() * 1000)}_{hashlib.sha1(data).hexdigest()[:10]}_{stem}{suffix}"
    target.write_bytes(data)
    return target


def _strip_html(text: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    return html.unescape(re.sub(r"\s+", " ", text)).strip()


def _read_text_file(path: Path, max_chars: int) -> str:
    suffix = path.suffix.lower()
    if suffix == ".csv":
        rows: list[str] = []
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            for index, row in enumerate(csv.reader(f)):
                if index >= 2000:
                    break
                rows.append("\t".join(str(cell) for cell in row))
        return "\n".join(rows)[:max_chars]
    raw = path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".json":
        try:
            raw = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
        except Exception:
            pass
    if suffix in {".html", ".htm", ".xml"}:
        raw = _strip_html(raw)
    return raw[:max_chars]


def _read_docx(path: Path, max_chars: int) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)[:max_chars]


def _read_xlsx(path: Path, max_chars: int) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    for name in wb.sheetnames[:20]:
        ws = wb[name]
        parts.append(f"# sheet: {name}")
        for index, row in enumerate(ws.iter_rows(values_only=True)):
            if index >= 500:
                break
            cells = ["" if value is None else str(value) for value in row]
            if any(cells):
                parts.append("\t".join(cells))
            if sum(len(part) for part in parts) >= max_chars:
                return "\n".join(parts)[:max_chars]
    return "\n".join(parts)[:max_chars]


def _read_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
        if texts:
            parts.append(f"# slide {idx}\n" + "\n".join(texts))
        if sum(len(part) for part in parts) >= max_chars:
            break
    return "\n\n".join(parts)[:max_chars]


def _read_pdf(path: Path, max_chars: int) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"# page {idx}\n{text}")
            if sum(len(part) for part in parts) >= max_chars:
                break
    return "\n\n".join(parts)[:max_chars]


def extract_text(path: Path, max_chars: int = DEFAULT_MAX_CHARS) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix not in DOCUMENT_SUFFIXES:
        raise ValueError(f"unsupported_file_type:{suffix or 'none'}")
    size = path.stat().st_size
    if size > MAX_PARSE_BYTES:
        raise ValueError(f"file_too_large_for_parse:{size}>{MAX_PARSE_BYTES}")
    if suffix in TEXT_SUFFIXES:
        return _read_text_file(path, max_chars), "text"
    if suffix == ".docx":
        return _read_docx(path, max_chars), "docx"
    if suffix == ".xlsx":
        return _read_xlsx(path, max_chars), "xlsx"
    if suffix == ".pptx":
        return _read_pptx(path, max_chars), "pptx"
    if suffix == ".pdf":
        return _read_pdf(path, max_chars), "pdf"
    raise ValueError(f"unsupported_file_type:{suffix}")


def _chunk_text(text: str) -> list[dict[str, Any]]:
    text = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
    if not text:
        return []
    blocks: list[dict[str, Any]] = []
    start = 0
    idx = 1
    while start < len(text):
        end = min(len(text), start + CHUNK_SIZE)
        if end < len(text):
            boundary = max(text.rfind("\n\n", start, end), text.rfind("。", start, end), text.rfind(".", start, end))
            if boundary > start + max(300, CHUNK_SIZE // 2):
                end = boundary + 1
        chunk = text[start:end].strip()
        if chunk:
            blocks.append({
                "local_id": f"chunk-{idx:04d}",
                "citation_id": f"C{idx:04d}",
                "title": f"Chunk {idx}",
                "text": chunk,
                "start": start,
                "end": end,
            })
            idx += 1
        if end >= len(text):
            break
        start = max(end - CHUNK_OVERLAP, start + 1)
    return blocks


def _summary(text: str, limit: int = 420) -> str:
    cleaned = _safe_text(text, limit * 2)
    return cleaned[:limit]


def _bounded_strings(value: Any, *, count: int, item_limit: int) -> list[str]:
    raw = value if isinstance(value, list) else []
    output: list[str] = []
    seen: set[str] = set()
    for item in raw:
        text = _safe_text(item, item_limit)
        key = text.casefold()
        if not text or key in seen:
            continue
        seen.add(key)
        output.append(text)
        if len(output) >= count:
            break
    return output


def _representative_text(text: str, limit: int = CARD_INPUT_MAX_CHARS) -> str:
    """Keep beginning, middle and ending evidence for one bounded LLM call."""
    value = str(text or "").strip()
    if len(value) <= limit:
        return value
    part = max(1000, limit // 3)
    middle_start = max(0, (len(value) - part) // 2)
    return "\n\n[...middle sample...]\n\n".join((
        value[:part],
        value[middle_start:middle_start + part],
        value[-part:],
    ))[:limit]


def _fallback_card(text: str, source: Path) -> dict[str, Any]:
    paragraphs = [
        _safe_text(item, 280)
        for item in re.split(r"\n\s*\n|(?<=[。！？.!?])\s+", str(text or ""))
        if _safe_text(item, 280)
    ]
    headings = [
        _safe_text(match.group(1), 120)
        for match in re.finditer(r"(?m)^\s*#{1,6}\s+(.+?)\s*$", str(text or ""))
    ]
    words = re.findall(r"[A-Za-z][A-Za-z0-9_.-]{2,}|[\u4e00-\u9fff]{2,8}", str(text or "").lower())
    counts: dict[str, int] = {}
    for word in words:
        counts[word] = counts.get(word, 0) + 1
    keywords = [
        word for word, _count in sorted(counts.items(), key=lambda item: (-item[1], item[0]))
        if word not in {"the", "and", "this", "that", "with", "from", "文件", "内容"}
    ][:12]
    return {
        "schema": "tiangong.v3.knowledge.card.v1",
        "title": source.stem[:160] or source.name[:160],
        "summary": _summary(text),
        "key_points": paragraphs[:6],
        "keywords": keywords,
        "outline": headings[:10],
        "content_extract": _safe_text(text, 1200),
        "extractor": "deterministic",
        "extraction_status": "fallback",
        "generated_at": _now_iso(),
    }


def _knowledge_card(text: str, source: Path, parser: str, payload: dict[str, Any]) -> dict[str, Any]:
    fallback = _fallback_card(text, source)
    enricher = _CARD_ENRICHER
    if not callable(enricher) or payload.get("llm_extract") is False:
        return fallback
    material = {
        "schema": "tiangong.v3.knowledge.card-input.v1",
        "file_name": source.name,
        "suffix": source.suffix.lower(),
        "parser": parser,
        "content": _representative_text(text),
    }
    try:
        value = enricher(material)
        if not isinstance(value, dict):
            raise ValueError("knowledge_card_enricher_invalid")
        summary = _safe_text(value.get("summary"), 900)
        if not summary:
            raise ValueError("knowledge_card_summary_missing")
        return {
            "schema": "tiangong.v3.knowledge.card.v1",
            "title": _safe_text(value.get("title") or fallback["title"], 160),
            "summary": summary,
            "key_points": _bounded_strings(value.get("key_points"), count=8, item_limit=360),
            "keywords": _bounded_strings(value.get("keywords"), count=16, item_limit=80),
            "outline": _bounded_strings(value.get("outline"), count=12, item_limit=160),
            "content_extract": _safe_text(value.get("content_extract") or fallback["content_extract"], 1600),
            "extractor": "llm",
            "extraction_status": "completed",
            "generated_at": _now_iso(),
        }
    except Exception as exc:
        fallback["extraction_error"] = f"{type(exc).__name__}: {_safe_text(exc, 240)}"
        return fallback


def _stored_copy(root: Path, source: Path, digest: str) -> Path:
    target = _files_dir(root) / f"{digest[:16]}_{_safe_name(source.name)}"
    if not target.exists() or target.stat().st_size != source.stat().st_size:
        shutil.copy2(source, target)
    return target


def _public_doc(ctx: dict[str, Any]) -> dict[str, Any]:
    meta = ctx.get("metadata") if isinstance(ctx.get("metadata"), dict) else {}
    blocks = ctx.get("blocks") if isinstance(ctx.get("blocks"), list) else []
    card = ctx.get("card") if isinstance(ctx.get("card"), dict) else {}
    return {
        "document_id": ctx.get("document_id", ""),
        "file_name": meta.get("file_name") or ctx.get("document_id", ""),
        "file_type": meta.get("file_type") or "",
        "suffix": meta.get("suffix") or "",
        "parser": meta.get("parser") or "",
        "created_at": ctx.get("created_at") or "",
        "updated_at": ctx.get("updated_at") or "",
        "citation_count": len(blocks),
        "size_bytes": int(meta.get("size_bytes") or 0),
        "summary": card.get("summary") or meta.get("summary") or "",
        "card": card,
        "card_title": card.get("title") or meta.get("file_name") or "",
        "key_points": list(card.get("key_points") or []),
        "keywords": list(card.get("keywords") or []),
        "content_extract": card.get("content_extract") or "",
        "extraction_status": card.get("extraction_status") or "legacy",
        "extractor": card.get("extractor") or "legacy",
        "file_path": meta.get("file_path") or meta.get("stored_path") or "",
        "stored_path": meta.get("stored_path") or "",
        "file_path_digest": meta.get("sha256") or "",
        "safe_projection_only": True,
        "raw_bytes_hidden": True,
    }


def _import_one(root: Path, path: Path, payload: dict[str, Any] | None = None, *, original_path: str = "") -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    source = _workspace_path(path)
    if not source.exists() or not source.is_file():
        raise FileNotFoundError(str(source))
    size = source.stat().st_size
    if size > MAX_TRANSFER_BYTES:
        raise ValueError(f"file_too_large:{size}>{MAX_TRANSFER_BYTES}")
    digest = _sha256_file(source)
    text, parser = extract_text(source, int(payload.get("max_chars") or DEFAULT_MAX_CHARS))
    blocks = _chunk_text(text)
    if not blocks:
        raise ValueError("no_extractable_text")
    stored = _stored_copy(root, source, digest)
    document_id = f"doc_{digest[:18]}"
    now = _now_iso()
    card = _knowledge_card(text, source, parser, payload)
    ctx = {
        "schema": "tiangong.v3.knowledge.context.v1",
        "document_id": document_id,
        "created_at": now,
        "updated_at": now,
        "metadata": {
            "file_name": source.name,
            "file_type": _kind_from_path(source),
            "suffix": source.suffix.lower(),
            "parser": parser,
            "size_bytes": size,
            "sha256": digest,
            "summary": card.get("summary") or _summary(text),
            "file_path": original_path or str(source),
            "stored_path": str(stored),
        },
        "card": card,
        "blocks": blocks,
    }
    previous = _load_context(root, document_id)
    if previous and previous.get("created_at"):
        ctx["created_at"] = previous["created_at"]
    _save_context(root, ctx)
    index = _load_index(root)
    documents = index.setdefault("documents", {})
    documents[document_id] = _public_doc(ctx)
    index["last_document_id"] = document_id
    _save_index(root, index)
    return _public_doc(ctx)


def _coerce_paths(payload: dict[str, Any]) -> list[Path]:
    raw_paths = payload.get("paths") or payload.get("filePaths") or payload.get("files") or []
    if isinstance(raw_paths, (str, os.PathLike)):
        raw_paths = [raw_paths]
    paths = []
    for raw in raw_paths if isinstance(raw_paths, list) else []:
        if raw:
            paths.append(_workspace_path(raw))
    return paths


def _inline_paths(root: Path, payload: dict[str, Any], *, handoff: bool = False) -> list[Path]:
    items = payload.get("items") or payload.get("inlineItems") or []
    out: list[Path] = []
    if not isinstance(items, list):
        return out
    target = _handoffs_dir(root) if handoff else _files_dir(root)
    for item in items[:20]:
        if not isinstance(item, dict):
            continue
        saved = _decode_data_url(item, target)
        if saved:
            out.append(saved)
    return out


def knowledge_list(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    root = knowledge_root(payload)
    index = _load_index(root)
    docs: list[dict[str, Any]] = []
    for document_id, entry in (index.get("documents") or {}).items():
        ctx = _load_context(root, str(document_id))
        if ctx:
            docs.append(_public_doc(ctx))
        elif isinstance(entry, dict):
            docs.append({"document_id": document_id, **entry})
    docs.sort(key=lambda item: str(item.get("updated_at") or item.get("created_at") or ""), reverse=True)
    return {
        "ok": True,
        "workspace": str(root),
        "knowledgeRoot": str(root),
        "index_path": str(_index_path(root)),
        "count": len(docs),
        "last_document_id": index.get("last_document_id") or "",
        "documents": docs,
    }


def import_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    paths = _coerce_paths(payload) + _inline_paths(root, payload, handoff=False)
    imported: list[dict[str, Any]] = []
    failed: list[dict[str, Any]] = []
    for source in paths:
        try:
            doc = _import_one(root, source, payload, original_path=str(source.resolve()))
            imported.append(doc)
        except Exception as exc:
            failed.append({"path": str(source), "error": f"{type(exc).__name__}: {_safe_text(exc, 300)}"})
    listing = knowledge_list({"knowledgeRoot": str(root)})
    return {
        **listing,
        "ok": bool(imported) or not failed,
        "imported": imported,
        "failed": failed,
    }


def _query_terms(query: str) -> list[str]:
    text = str(query or "").lower()
    terms: list[str] = []
    for term in re.findall(r"[a-zA-Z0-9_./:-]{2,}", text):
        terms.append(term)
    cjk = "".join(re.findall(r"[\u4e00-\u9fff]+", text))
    if cjk:
        terms.append(cjk)
        terms.extend(cjk[i:i + 2] for i in range(max(0, len(cjk) - 1)))
    compact = re.sub(r"\s+", " ", text).strip()
    if compact and len(compact) <= 80:
        terms.append(compact)
    seen = set()
    uniq = []
    for term in terms:
        if term and term not in seen:
            seen.add(term)
            uniq.append(term)
    return uniq[:80]


def _score_text(text: str, query: str, terms: list[str]) -> int:
    lower = str(text or "").lower()
    score = 0
    q = str(query or "").lower().strip()
    if q and q in lower:
        score += 30
    for term in terms:
        hits = lower.count(term)
        if hits:
            score += hits * (6 if len(term) >= 3 else 3)
    return score


def _query_context(ctx: dict[str, Any], query: str, top_k: int = 6) -> dict[str, Any]:
    terms = _query_terms(query)
    blocks = [block for block in (ctx.get("blocks") or []) if isinstance(block, dict)]
    matches = []
    card = ctx.get("card") if isinstance(ctx.get("card"), dict) else {}
    card_text = "\n".join([
        str(card.get("title") or ""),
        str(card.get("summary") or ""),
        "\n".join(str(item) for item in card.get("key_points") or []),
        " ".join(str(item) for item in card.get("keywords") or []),
        "\n".join(str(item) for item in card.get("outline") or []),
        str(card.get("content_extract") or ""),
    ])
    card_score = _score_text(card_text, query, terms)
    if card_score > 0:
        matches.append({
            "local_id": "knowledge-card",
            "citation_id": "CARD",
            "title": str(card.get("title") or "Knowledge card"),
            "text": "\n".join([
                str(card.get("summary") or ""),
                *[f"- {item}" for item in card.get("key_points") or []],
            ])[:1800],
            "score": card_score,
            "kind": "knowledge_card",
        })
    for block in blocks:
        score = _score_text(block.get("text", ""), query, terms)
        if score <= 0:
            continue
        matches.append({
            "local_id": block.get("local_id") or "",
            "citation_id": block.get("citation_id") or "",
            "title": block.get("title") or "",
            "text": str(block.get("text") or "")[:1400],
            "score": score,
        })
    matches.sort(key=lambda item: int(item.get("score") or 0), reverse=True)
    matches = matches[:max(1, min(int(top_k or 6), 20))]
    return {
        "answer_summary": f"匹配到 {len(matches)} 个内容片段。" if matches else "没有匹配到内容片段。",
        "matches": matches,
        "score": sum(int(item.get("score") or 0) for item in matches),
    }


def query_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    query = _safe_text(payload.get("query"), 1000)
    if not query:
        return {"ok": False, "error": "missing_query"}
    document_id = _safe_text(payload.get("document_id") or payload.get("documentId"), 120)
    if not document_id:
        cards = search_knowledge({**payload, "knowledgeRoot": str(root), "query": query, "top_k": payload.get("top_k") or 6})
        return {"ok": cards.get("ok", False), "query": query, "result": {"answer_summary": f"匹配到 {cards.get('count', 0)} 篇文档。", "matches": cards.get("cards", [])}}
    ctx = _load_context(root, document_id)
    if not ctx:
        return {"ok": False, "error": "document_not_found", "document_id": document_id}
    result = _query_context(ctx, query, int(payload.get("top_k") or 6))
    return {"ok": True, "workspace": str(root), "document_id": document_id, "query": query, "result": result}


def search_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    query = _safe_text(payload.get("query"), 1000)
    if not query:
        return {"ok": False, "error": "missing_query", "cards": []}
    limit = max(1, min(int(payload.get("top_k") or payload.get("limit") or 8), 30))
    per_doc = max(1, min(int(payload.get("per_doc") or 3), 10))
    index = _load_index(root)
    cards: list[dict[str, Any]] = []
    for document_id in (index.get("documents") or {}).keys():
        ctx = _load_context(root, str(document_id))
        if not ctx:
            continue
        result = _query_context(ctx, query, per_doc)
        if int(result.get("score") or 0) <= 0:
            continue
        meta = ctx.get("metadata") if isinstance(ctx.get("metadata"), dict) else {}
        card = ctx.get("card") if isinstance(ctx.get("card"), dict) else {}
        cards.append({
            "document_id": ctx.get("document_id", ""),
            "title": meta.get("file_name") or ctx.get("document_id", ""),
            "file_name": meta.get("file_name") or ctx.get("document_id", ""),
            "file_type": meta.get("file_type") or "",
            "suffix": meta.get("suffix") or "",
            "parser": meta.get("parser") or "",
            "summary": card.get("summary") or meta.get("summary") or "",
            "card": card,
            "key_points": list(card.get("key_points") or []),
            "keywords": list(card.get("keywords") or []),
            "content_extract": card.get("content_extract") or "",
            "extraction_status": card.get("extraction_status") or "legacy",
            "created_at": ctx.get("created_at") or "",
            "citation_count": len(ctx.get("blocks") or []),
            "file_path": meta.get("file_path") or meta.get("stored_path") or "",
            "stored_path": meta.get("stored_path") or "",
            "score": int(result.get("score") or 0),
            "matches": result.get("matches") or [],
            "safe_projection_only": True,
            "raw_bytes_hidden": True,
        })
    cards.sort(key=lambda item: (-int(item.get("score") or 0), str(item.get("created_at") or "")))
    cards = cards[:limit]
    return {
        "ok": True,
        "workspace": str(root),
        "knowledgeRoot": str(root),
        "index_path": str(_index_path(root)),
        "query": query,
        "count": len(cards),
        "cards": cards,
        "safe_projection_only": True,
        "raw_bytes_hidden": True,
    }


def organize_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    root = knowledge_root(payload)
    documents: dict[str, Any] = {}
    broken: list[str] = []
    for path in sorted(_contexts_dir(root).glob("*.json")):
        try:
            ctx = json.loads(path.read_text(encoding="utf-8"))
            if isinstance(ctx, dict) and ctx.get("document_id"):
                documents[str(ctx["document_id"])] = _public_doc(ctx)
            else:
                broken.append(str(path))
        except Exception:
            broken.append(str(path))
    index = {
        "schema": "tiangong.v3.knowledge.index.v1",
        "updated_at": _now_iso(),
        "documents": documents,
        "last_document_id": next(iter(documents.keys()), ""),
    }
    _save_index(root, index)
    listing = knowledge_list({"knowledgeRoot": str(root)})
    return {
        **listing,
        "organized": True,
        "report": {
            "ok": True,
            "workspace": str(root),
            "document_count": len(documents),
            "broken_count": len(broken),
            "broken": broken[:20],
            "policy": "index_only_no_source_file_mutation",
        },
    }


def _render_markdown(ctx: dict[str, Any]) -> str:
    meta = ctx.get("metadata") if isinstance(ctx.get("metadata"), dict) else {}
    card = ctx.get("card") if isinstance(ctx.get("card"), dict) else {}
    lines = [
        f"# {card.get('title') or meta.get('file_name') or ctx.get('document_id')}",
        "",
        f"- document_id: {ctx.get('document_id')}",
        f"- source: {meta.get('file_path') or ''}",
        f"- created_at: {ctx.get('created_at') or ''}",
        "",
        "## Summary",
        "",
        card.get("summary") or meta.get("summary") or "",
        "",
        "## Key points",
        "",
        *[f"- {item}" for item in card.get("key_points") or []],
        "",
        "## Keywords",
        "",
        ", ".join(str(item) for item in card.get("keywords") or []),
        "",
        "## Content extract",
        "",
        str(card.get("content_extract") or ""),
        "",
        "## Chunks",
        "",
    ]
    for block in ctx.get("blocks") or []:
        if not isinstance(block, dict):
            continue
        lines.extend([f"### {block.get('citation_id') or block.get('local_id')}", "", str(block.get("text") or ""), ""])
    return "\n".join(lines)


def export_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    document_id = _safe_text(payload.get("document_id") or payload.get("documentId"), 120)
    ctx = _load_context(root, document_id)
    if not ctx:
        return {"ok": False, "error": "document_not_found", "document_id": document_id}
    fmt = str(payload.get("format") or "md").lower().lstrip(".")
    if fmt not in {"md", "markdown", "txt", "json"}:
        fmt = "md"
    suffix = "json" if fmt == "json" else ("txt" if fmt == "txt" else "md")
    target = _exports_dir(root) / f"{document_id}_summary.{suffix}"
    if suffix == "json":
        content = json.dumps(ctx, ensure_ascii=False, indent=2)
    elif suffix == "txt":
        content = "\n\n".join(str(block.get("text") or "") for block in ctx.get("blocks") or [] if isinstance(block, dict))
    else:
        content = _render_markdown(ctx)
    target.write_text(content, encoding="utf-8")
    return {"ok": True, "workspace": str(root), "document_id": document_id, "target": str(target), "format": suffix}


def remove_knowledge(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    document_id = _safe_text(payload.get("document_id") or payload.get("documentId"), 120)
    if not document_id:
        return {"ok": False, "error": "missing_document_id"}
    index = _load_index(root)
    docs = dict(index.get("documents") or {})
    entry = docs.get(document_id) if isinstance(docs.get(document_id), dict) else {}
    stored_path = Path(str(entry.get("stored_path") or "")).resolve(strict=False) if entry.get("stored_path") else None
    existed = document_id in docs
    docs.pop(document_id, None)
    index["documents"] = docs
    if index.get("last_document_id") == document_id:
        index["last_document_id"] = next(iter(docs.keys()), "")
    _save_index(root, index)
    removed_context = False
    with_context = _context_path(root, document_id)
    if with_context.exists():
        with_context.unlink()
        removed_context = True
    removed_stored_file = False
    files_root = _files_dir(root).resolve(strict=False)
    if stored_path and stored_path.is_relative_to(files_root) and stored_path.is_file():
        stored_path.unlink()
        removed_stored_file = True
    listing = knowledge_list({"knowledgeRoot": str(root)})
    return {
        **listing,
        "ok": existed or removed_context or removed_stored_file,
        "removed_document_id": document_id,
        "removed_context": removed_context,
        "removed_stored_file": removed_stored_file,
    }


def import_files(payload: dict[str, Any] | None = None) -> dict[str, Any]:
    payload = payload if isinstance(payload, dict) else {}
    root = knowledge_root(payload)
    source_paths = _coerce_paths(payload)
    inline_paths = _inline_paths(root, payload, handoff=True)
    paths = source_paths + inline_paths
    import_to_knowledge = payload.get("import_to_knowledge", payload.get("importToKnowledge", True)) is not False
    attachments: list[dict[str, Any]] = []
    imported: list[dict[str, Any]] = []
    failed: list[dict[str, Any]] = []

    total = 0
    for path in paths:
        try:
            if path.exists() and path.is_file():
                total += path.stat().st_size
        except Exception:
            pass
    if total > MAX_TRANSFER_BYTES:
        for path in paths:
            failed.append({"path": str(path), "error": f"total_upload_too_large:{total}>{MAX_TRANSFER_BYTES}"})
            attachments.append(_attachment(path, status="failed", error=failed[-1]["error"]))
        return {"ok": False, "attachments": attachments, "imported": imported, "failed": failed, "knowledgeRoot": str(root)}

    for raw_path in paths:
        path = Path(raw_path).expanduser()
        if not path.exists() or not path.is_file():
            item = _attachment(path, status="failed", error="file_not_found")
            attachments.append(item)
            failed.append({"path": str(path), "error": "file_not_found"})
            continue
        doc = None
        err = ""
        if import_to_knowledge and path.suffix.lower() in DOCUMENT_SUFFIXES:
            try:
                doc = _import_one(root, path, payload, original_path=str(path.resolve()))
                imported.append(doc)
            except Exception as exc:
                err = f"{type(exc).__name__}: {_safe_text(exc, 240)}"
                failed.append({"path": str(path), "error": err})
        status = "imported" if doc else "attached"
        if err and path.suffix.lower() in DOCUMENT_SUFFIXES:
            status = "attached"
        attachments.append(_attachment(path, doc, status=status, error="", import_error=err))
    return {
        "ok": any(item.get("status") in {"imported", "attached"} for item in attachments),
        "attachments": attachments,
        "imported": imported,
        "failed": failed,
        "knowledgeRoot": str(root),
    }


def _attachment(path: Path, doc: dict[str, Any] | None = None, *, status: str = "attached", error: str = "", import_error: str = "") -> dict[str, Any]:
    try:
        resolved = Path(path).expanduser().resolve()
    except Exception:
        resolved = Path(str(path))
    stat = None
    try:
        stat = resolved.stat()
    except Exception:
        pass
    return {
        "path": str(resolved),
        "name": resolved.name,
        "ext": resolved.suffix.lower().lstrip("."),
        "kind": _kind_from_path(resolved),
        "size": int(stat.st_size) if stat else 0,
        "documentId": doc.get("document_id", "") if doc else "",
        "document_id": doc.get("document_id", "") if doc else "",
        "status": status,
        "summary": doc.get("summary", "") if doc else "",
        "citationCount": int(doc.get("citation_count") or 0) if doc else 0,
        "error": error,
        "importError": import_error,
    }
