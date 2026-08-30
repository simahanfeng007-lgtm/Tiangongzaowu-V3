"""P19-R2 M2.2 Gateway ArtifactContentOracle — RECORD ONLY.

Evaluates one explicit ``AcceptancePredicate`` against one bound
``ArtifactManifest`` by reading immutable object-store bytes through the
existing trusted chain (``VerifiedArtifactContentSource``). The oracle
never opens host paths and never derives predicates from the raw user
message — that heuristic layer stays in the V3 local preflight
(``v3/simple_chain/content_preflight.py``) with no final authority.

Status: implementation-present / descriptor-registered /
production-unwired. Nothing consumes its records to gate completion.

M2.2 hardening (review 2026-08-30):
* predicate trust boundary runs the FULL semantic identity check
  (hashes + canonical params + type/subject domain), never a bare hash;
* AUTHORITY BEFORE ANY STATUS: manifest/QC-facts/object-reference
  verification happens before applicability/capability/format/size
  decisions — even NOT_APPLICABLE and INCONCLUSIVE subjects carry an
  authority-verified manifest;
* resource discipline: the verified manifest's size_bytes is checked
  against the descriptor limit BEFORE any blob read;
* determinism: xlsx formula text counts as content everywhere (single
  data_only=False pass); pptx meaningful content uses the official shape
  enum (picture/media/group/table/chart) plus table-cell text; text
  marker matching normalizes BOTH sides (NFKC + casefold);
* parser taxonomy: known zip/XML container corruption -> INCONCLUSIVE,
  ImportError/dependency -> ERROR, unknown runtime bugs -> ERROR;
* reason codes are fixed machine codes (no user strings); missing-item
  details live only in the observation digest (missing_count +
  missing_items_sha256);
* the pinned descriptor is validated EXACTLY against the oracle config —
  same id/version with a different config or capability is rejected.

RECORD ONLY: produced records change no CompletionDecision, no request
state, no delivery.
"""

from __future__ import annotations

import csv as _csv
import io
import zipfile
import xml.etree.ElementTree as _etree
from typing import Any

from contracts import ArtifactManifest, canonical_sha256
from contracts.verification import (
    AcceptancePredicate,
    RegistrySnapshot,
    VerificationRecord,
    derive_verification_record_id,
    normalize_predicate_text,
)
from total_gateway.artifact_content import (
    ArtifactContentError,
    VerifiedArtifactContentSource,
)
from total_gateway.fact_ledger import FactLedger
from total_gateway.object_store import ContentAddressedObjectStore
from total_gateway.verification_oracle_config import (
    ARTIFACT_APPLICABLE_FORMATS,
    ARTIFACT_DESCRIPTOR_EXPECTATIONS,
    ARTIFACT_IMPLEMENTED_PREDICATE_TYPES,
    ARTIFACT_INSPECTABLE_FORMATS,
    ARTIFACT_INSPECTOR_SEMANTIC_VERSION,
    ARTIFACT_MAX_INPUT_BYTES,
    IMPLEMENTATION_REF,
    VERIFIER_ID,
    artifact_oracle_config_sha256,
)
from total_gateway.verification_registry import (
    UnknownVerifierError,
    VerifierRegistry,
)


class OracleSnapshotInvalid(ValueError):
    """Raised when the oracle cannot be bound to the given snapshot."""


class ContentUnparseable(Exception):
    """Known container/zip/XML corruption — cannot reliably inspect."""


class _InspectorDependencyMissing(Exception):
    """ImportError from an inspector dependency (ERROR, not INCONCLUSIVE)."""


#: Exception types that mean "known corrupted container" -> INCONCLUSIVE.
_KNOWN_CORRUPTION_TYPES: tuple[type[BaseException], ...] = (
    zipfile.BadZipFile,
    _etree.ParseError,
)


# ---------------------------------------------------------------------------
# Descriptor binding (§2)
# ---------------------------------------------------------------------------

def _validate_artifact_descriptor(descriptor) -> None:
    """Exact-match the pinned descriptor against the oracle config.

    A same-id/same-version descriptor whose config hash, capability set,
    limits or authority claims differ from what THIS implementation reads
    is rejected — no silent drift.
    """
    expectations = ARTIFACT_DESCRIPTOR_EXPECTATIONS
    problems: list[str] = []
    if descriptor.verifier_id != VERIFIER_ID:
        problems.append("verifier_id")
    if descriptor.verifier_version != ARTIFACT_INSPECTOR_SEMANTIC_VERSION:
        problems.append("verifier_version")
    if descriptor.config_sha256 != artifact_oracle_config_sha256():
        problems.append("config_sha256")
    if list(descriptor.supported_predicate_types) != sorted(
        ARTIFACT_IMPLEMENTED_PREDICATE_TYPES
    ):
        problems.append("supported_predicate_types")
    if tuple(descriptor.supported_subject_kinds) != expectations[
        "supported_subject_kinds"
    ]:
        problems.append("supported_subject_kinds")
    if set(descriptor.accepted_authorities) != set(
        expectations["accepted_authorities"]
    ):
        problems.append("accepted_authorities")
    if descriptor.max_input_bytes != ARTIFACT_MAX_INPUT_BYTES:
        problems.append("max_input_bytes")
    if descriptor.implementation_ref != IMPLEMENTATION_REF:
        problems.append("implementation_ref")
    if descriptor.producer_component_id != expectations["producer_component_id"]:
        problems.append("producer_component_id")
    if descriptor.layer != expectations["layer"]:
        problems.append("layer")
    if descriptor.deterministic is not expectations["deterministic"]:
        problems.append("deterministic")
    if descriptor.default_enforcement != expectations["default_enforcement"]:
        problems.append("default_enforcement")
    if problems:
        raise OracleSnapshotInvalid(
            "artifact descriptor does not match the oracle config:"
            f" {', '.join(problems)}"
        )


# ---------------------------------------------------------------------------
# Inspectors (pure readers over verified bytes; no host paths)
# ---------------------------------------------------------------------------

def _inspect_docx(data: bytes) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise _InspectorDependencyMissing("docx") from exc
    try:
        document = Document(io.BytesIO(data))
        paragraph_chars = sum(
            len(paragraph.text.strip()) for paragraph in document.paragraphs
        )
        table_chars = 0
        table_cells = 0
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        table_cells += 1
                        table_chars += len(text)
        return {
            "measured_format": "docx",
            "visible_text_chars": paragraph_chars + table_chars,
        }
    except _InspectorDependencyMissing:
        raise
    except BaseException as exc:
        if isinstance(exc, _KNOWN_CORRUPTION_TYPES):
            raise ContentUnparseable("docx container corrupted") from exc
        raise  # unknown runtime bug -> ERROR by the dispatcher


def _inspect_xlsx(data: bytes) -> dict[str, Any]:
    """Single data_only=False pass: formula text is content everywhere.

    Formula cells carry no cached value in this mode, so their formula
    string is what we see — it counts toward visible chars, non-empty
    cells AND data rows (M2.2 review rule).
    """
    try:
        import openpyxl
    except ImportError as exc:
        raise _InspectorDependencyMissing("openpyxl") from exc
    workbook = None
    try:
        workbook = openpyxl.load_workbook(
            io.BytesIO(data), read_only=True, data_only=False
        )
        sheet = workbook.active
        if sheet is None:
            raise ContentUnparseable("xlsx workbook has no active sheet")
        header: list[str] = []
        data_row_count = 0
        nonempty_cell_count = 0
        visible_text_chars = 0
        for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
            row_nonempty = 0
            for cell in row:
                if cell is None:
                    continue
                text = str(cell).strip()
                if not text:
                    continue
                nonempty_cell_count += 1
                row_nonempty += 1
                visible_text_chars += len(text)
            if row_index == 0:
                header = ["" if cell is None else str(cell).strip() for cell in row]
                continue  # header is not a data row (oracle semantics)
            if row_nonempty:
                data_row_count += 1
        sheet_names = list(workbook.sheetnames)
        return {
            "measured_format": "xlsx",
            "header": header,
            "data_row_count": data_row_count,
            "nonempty_cell_count": nonempty_cell_count,
            "sheet_names": sheet_names,
            "visible_text_chars": visible_text_chars,
        }
    except _InspectorDependencyMissing:
        raise
    except ContentUnparseable:
        raise
    except BaseException as exc:
        if isinstance(exc, _KNOWN_CORRUPTION_TYPES):
            raise ContentUnparseable("xlsx container corrupted") from exc
        raise  # unknown runtime bug -> ERROR by the dispatcher
    finally:
        if workbook is not None:
            workbook.close()


def _inspect_pptx(data: bytes) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise _InspectorDependencyMissing("pptx") from exc
    # Shape types that are meaningful content even without text.
    meaningful_shape_types = {
        member
        for member in (
            getattr(MSO_SHAPE_TYPE, "PICTURE", None),
            getattr(MSO_SHAPE_TYPE, "LINKED_PICTURE", None),
            getattr(MSO_SHAPE_TYPE, "MEDIA", None),
            getattr(MSO_SHAPE_TYPE, "GROUP", None),
        )
        if member is not None
    }
    try:
        presentation = Presentation(io.BytesIO(data))
        slide_count = 0
        meaningful_slide_count = 0
        text_bearing_slide_count = 0
        slide_text_chars = 0
        total_shape_count = 0
        for slide in presentation.slides:
            slide_count += 1
            has_text = False
            has_other_content = False
            for shape in slide.shapes:
                total_shape_count += 1
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        has_text = True
                        slide_text_chars += len(text)
                if getattr(shape, "has_table", False):
                    has_other_content = True
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                slide_text_chars += len(cell_text)
                if getattr(shape, "has_chart", False):
                    has_other_content = True
                if shape.shape_type in meaningful_shape_types:
                    has_other_content = True
            if has_text:
                text_bearing_slide_count += 1
            if has_text or has_other_content:
                meaningful_slide_count += 1
        return {
            "measured_format": "pptx",
            "slide_count": slide_count,
            "text_bearing_slide_count": text_bearing_slide_count,
            "meaningful_slide_count": meaningful_slide_count,
            "slide_text_chars": slide_text_chars,
            "total_shape_count": total_shape_count,
        }
    except _InspectorDependencyMissing:
        raise
    except BaseException as exc:
        if isinstance(exc, _KNOWN_CORRUPTION_TYPES):
            raise ContentUnparseable("pptx container corrupted") from exc
        raise  # unknown runtime bug -> ERROR by the dispatcher


def _inspect_text(data: bytes, *, filename: str) -> dict[str, Any]:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ContentUnparseable("text not decodable as utf-8") from exc
    metrics: dict[str, Any] = {
        "measured_format": "text",
        # Marker matching normalizes BOTH sides: haystack here, needles in
        # the predicate params (both NFKC + casefold).
        "_text_norm": normalize_predicate_text(text),
        "char_count": len(text),
        "non_whitespace_chars": sum(1 for ch in text if not ch.isspace()),
        "line_count": sum(1 for line in text.splitlines() if line.strip()),
    }
    if filename.lower().endswith(".csv"):
        rows = list(_csv.reader(io.StringIO(text)))
        metrics["header"] = [cell.strip() for cell in rows[0]] if rows else []
        metrics["data_row_count"] = sum(
            1 for row in rows[1:] if any(cell.strip() for cell in row)
        )
    return metrics


_INSPECTORS = {
    "docx": _inspect_docx,
    "xlsx": _inspect_xlsx,
    "pptx": _inspect_pptx,
}


def _inspect(format_id: str, data: bytes, *, filename: str) -> dict[str, Any]:
    if format_id == "text":
        return _inspect_text(data, filename=filename)
    inspector = _INSPECTORS.get(format_id)
    if inspector is None:  # defensive: guarded by _INSPECTABLE_FORMATS
        raise ContentUnparseable(f"no inspector for format: {format_id}")
    return inspector(data)


# ---------------------------------------------------------------------------
# Deterministic checks — fixed machine reason codes only (§4)
# ---------------------------------------------------------------------------

#: Per-format metric key used by artifact.min_visible_text_chars. A format
#: without a reliable metric key raises ContentUnparseable (INCONCLUSIVE)
#: — never a fake 0.
_VISIBLE_TEXT_METRIC_KEYS: dict[str, str] = {
    "docx": "visible_text_chars",
    "xlsx": "visible_text_chars",
    "pptx": "slide_text_chars",
    "text": "non_whitespace_chars",
}


def _visible_text_metric(format_id: str, metrics: dict[str, Any]) -> int:
    key = _VISIBLE_TEXT_METRIC_KEYS.get(format_id)
    if key is None or key not in metrics:
        raise ContentUnparseable(
            f"no reliable visible-text metric for format: {format_id}"
        )
    value = metrics[key]
    if not isinstance(value, int) or value < 0:
        raise ContentUnparseable(f"unusable visible-text metric for format: {format_id}")
    return value


def _pptx_emptiness_is_unreliable(metrics: dict[str, Any]) -> bool:
    """Slides carry shapes we cannot classify: not provably empty."""
    return (
        metrics["slide_count"] > 0
        and metrics["meaningful_slide_count"] == 0
        and metrics["total_shape_count"] > 0
    )


def _check_predicate(
    predicate: AcceptancePredicate,
    format_id: str,
    metrics: dict[str, Any],
) -> tuple[bool, tuple[str, ...], dict[str, Any]]:
    """Return (holds, reason_codes, observation_extra).

    Reason codes are FIXED machine codes — user strings (column/marker
    names) never appear; their detail goes into observation_extra which
    feeds only the observation digest (missing_count + items digest).
    """
    kind = predicate.predicate_type
    params = predicate.param_mapping()
    if kind == "artifact.nonempty":
        if format_id == "xlsx":
            holds = metrics["nonempty_cell_count"] > 0
        elif format_id == "pptx":
            if _pptx_emptiness_is_unreliable(metrics):
                raise ContentUnparseable(
                    "pptx slides carry unclassifiable content; emptiness"
                    " cannot be judged reliably"
                )
            holds = metrics["meaningful_slide_count"] > 0
        else:
            holds = _visible_text_metric(format_id, metrics) > 0
        return holds, () if holds else ("artifact.content_empty",), {}
    if kind == "artifact.min_visible_text_chars":
        minimum = params["min_chars"]
        assert isinstance(minimum, int)
        visible = _visible_text_metric(format_id, metrics)
        holds = visible >= minimum
        return holds, () if holds else ("artifact.visible_text_below_minimum",), {}
    if kind == "xlsx.required_columns":
        requested = params["columns"]
        assert isinstance(requested, tuple)
        header = {
            normalize_predicate_text(cell)
            for cell in metrics["header"]
            if isinstance(cell, str) and cell.strip()
        }
        missing = [column for column in requested if column not in header]
        return (
            not missing,
            () if not missing else ("xlsx.required_columns_missing",),
            {}
            if not missing
            else {
                "missing_count": len(missing),
                "missing_items_sha256": canonical_sha256(list(missing)),
            },
        )
    if kind == "xlsx.min_data_rows":
        minimum = params["min_rows"]
        assert isinstance(minimum, int)
        holds = metrics["data_row_count"] >= minimum
        return holds, () if holds else ("xlsx.data_rows_below_minimum",), {}
    if kind == "text.required_markers":
        markers = params["markers"]
        assert isinstance(markers, tuple)
        haystack = metrics.get("_text_norm", "")
        missing = [marker for marker in markers if marker not in haystack]
        return (
            not missing,
            () if not missing else ("text.required_markers_missing",),
            {}
            if not missing
            else {
                "missing_count": len(missing),
                "missing_items_sha256": canonical_sha256(list(missing)),
            },
        )
    if kind == "pptx.min_nonempty_slides":
        minimum = params["min_slides"]
        assert isinstance(minimum, int)
        if _pptx_emptiness_is_unreliable(metrics):
            raise ContentUnparseable(
                "pptx slides carry unclassifiable content; meaningful"
                " slide count cannot be judged reliably"
            )
        holds = metrics["meaningful_slide_count"] >= minimum
        return (
            holds,
            () if holds else ("pptx.meaningful_slides_below_minimum",),
            {},
        )
    raise AssertionError(
        f"predicate type declared implemented but has no check: {kind}"
    )


# ---------------------------------------------------------------------------
# Oracle
# ---------------------------------------------------------------------------

class ArtifactContentOracle:
    """Deterministic content oracle over immutable, QC-passed artifacts.

    Bound to one validated ``RegistrySnapshot`` whose artifact descriptor
    is EXACTLY the one this implementation reads its config from.
    """

    def __init__(
        self,
        *,
        snapshot: RegistrySnapshot,
        object_store: ContentAddressedObjectStore,
        fact_ledger: FactLedger,
    ) -> None:
        if not snapshot.has_valid_identity():
            raise OracleSnapshotInvalid("registry snapshot identity binding is invalid")
        try:
            registry = VerifierRegistry(snapshot.verifiers)
        except ValueError as exc:
            raise OracleSnapshotInvalid(
                "registry snapshot contains invalid verifier descriptors"
            ) from exc
        try:
            descriptor = registry.find(
                VERIFIER_ID, ARTIFACT_INSPECTOR_SEMANTIC_VERSION
            )
        except UnknownVerifierError as exc:
            raise OracleSnapshotInvalid(
                f"snapshot does not carry {VERIFIER_ID}@"
                f"{ARTIFACT_INSPECTOR_SEMANTIC_VERSION}"
            ) from exc
        _validate_artifact_descriptor(descriptor)
        object.__setattr__(self, "_snapshot", snapshot)
        object.__setattr__(self, "_descriptor", descriptor)
        object.__setattr__(self, "_object_store", object_store)
        object.__setattr__(self, "_fact_ledger", fact_ledger)

    @property
    def descriptor(self):
        return self._descriptor  # type: ignore[attr-defined]

    def evaluate(
        self,
        manifest: ArtifactManifest,
        predicate: AcceptancePredicate,
        *,
        evaluated_at_ms: int,
        evaluation_phase: str = "POST_EXECUTION",
    ) -> VerificationRecord:
        # 1. Predicate semantic identity — FULL check, never a bare hash.
        if predicate.subject_kind != "artifact" or not predicate.has_valid_identity():
            raise ValueError("predicate failed full semantic identity validation")
        status, reason_codes, observation = self._evaluate_to_status(
            manifest, predicate
        )
        return self._build_record(
            manifest=manifest,
            predicate=predicate,
            evaluated_at_ms=evaluated_at_ms,
            evaluation_phase=evaluation_phase,
            status=status,
            reason_codes=reason_codes,
            observation=observation,
        )

    # -- internals ---------------------------------------------------------

    def _evaluate_to_status(
        self, manifest: ArtifactManifest, predicate: AcceptancePredicate
    ) -> tuple[str, tuple[str, ...], dict[str, Any]]:
        format_id = manifest.format_id

        # 2-3. AUTHORITY BEFORE ANY STATUS: manifest PASSED + hash, QC
        # facts, object reference binding. Every downstream verdict —
        # including NOT_APPLICABLE and INCONCLUSIVE — is bound to a
        # subject that passed this verification. ArtifactContentError
        # messages ARE stable machine codes (artifact.content.*).
        try:
            source = VerifiedArtifactContentSource(
                self._object_store, self._fact_ledger, (manifest,)  # type: ignore[attr-defined]
            )
            source.verify_artifact_revision(manifest.artifact_revision_id)
        except ArtifactContentError as exc:
            return "ERROR", (f"authority:{exc}",), {"authority_error": str(exc)}
        except Exception as exc:  # object-store corruption, IO, runtime
            return "ERROR", ("authority_failure",), {
                "authority_error": type(exc).__name__
            }

        # 4. Applicability (subject already authority-verified above).
        prefix = predicate.predicate_type.split(".", 1)[0]
        applicable = ARTIFACT_APPLICABLE_FORMATS.get(prefix)
        if applicable is not None and format_id not in applicable:
            return "NOT_APPLICABLE", (), {
                "predicate_prefix": prefix,
                "format_id": format_id,
            }

        # 5. Descriptor capability — from the pinned descriptor.
        if predicate.predicate_type not in self._descriptor.supported_predicate_types:  # type: ignore[attr-defined]
            return "INCONCLUSIVE", ("predicate_not_implemented",), {
                "predicate_type": predicate.predicate_type
            }

        # 6. Format inspectability.
        if format_id not in ARTIFACT_INSPECTABLE_FORMATS:
            return "INCONCLUSIVE", ("format_not_inspectable",), {
                "format_id": format_id
            }

        # 7-8. Size pre-check on the VERIFIED manifest — oversize never
        # touches object-store bytes.
        max_input_bytes = self._descriptor.max_input_bytes  # type: ignore[attr-defined]
        if manifest.size_bytes > max_input_bytes:
            return "INCONCLUSIVE", ("input_too_large",), {
                "size_bytes": manifest.size_bytes,
                "max_input_bytes": max_input_bytes,
            }

        # 9. Immutable readback (re-verifies on the way in).
        try:
            data = source.read_verified_artifact(manifest.artifact_revision_id)
        except ArtifactContentError as exc:
            return "ERROR", (f"authority:{exc}",), {"authority_error": str(exc)}
        except Exception as exc:
            return "ERROR", ("authority_failure",), {
                "authority_error": type(exc).__name__
            }

        # 10. Inspection + deterministic check.
        try:
            metrics = _inspect(format_id, data, filename=manifest.filename)
        except _InspectorDependencyMissing:
            return "ERROR", ("inspector_dependency_missing",), {}
        except ContentUnparseable:
            return "INCONCLUSIVE", ("content_unparseable",), {}
        except Exception:  # unknown runtime bug — never silent
            return "ERROR", ("inspector_failure",), {}

        try:
            holds, reason_codes, extra = _check_predicate(
                predicate, format_id, metrics
            )
        except ContentUnparseable:
            return "INCONCLUSIVE", ("content_unparseable",), {}
        except Exception:  # unknown runtime bug — never silent
            return "ERROR", ("check_failure",), {}

        observation = {
            key: value
            for key, value in metrics.items()
            if key not in ("_text_norm", "_text")
        }
        observation.update(extra)
        observation["inspector_version"] = self._descriptor.verifier_version  # type: ignore[attr-defined]
        return ("PASS" if holds else "FAIL"), reason_codes, observation

    def _build_record(
        self,
        *,
        manifest: ArtifactManifest,
        predicate: AcceptancePredicate,
        evaluated_at_ms: int,
        evaluation_phase: str,
        status: str,
        reason_codes: tuple[str, ...],
        observation: dict[str, Any],
    ) -> VerificationRecord:
        descriptor = self._descriptor  # type: ignore[attr-defined]
        qc_evidence_set = sorted(
            (
                evidence.check_id,
                evidence.check_version,
                evidence.tool_fact_id,
                evidence.evidence_sha256,
            )
            for evidence in manifest.qc_evidence
        )
        observation_payload = {
            "measured_format": observation.get("measured_format"),
            "inspector_version": observation.get("inspector_version"),
            "metrics": {
                key: (list(value) if isinstance(value, (list, tuple)) else value)
                for key, value in sorted(observation.items())
                if key not in ("measured_format", "inspector_version")
            },
        }
        evidence_refs = (
            f"artifact_revision:{manifest.artifact_revision_id}",
            f"manifest_sha256:{manifest.manifest_sha256}",
            f"content_object:{manifest.content_object_id}",
            f"content_sha256:{manifest.sha256}",
            f"qc_evidence_set_sha256:{canonical_sha256(qc_evidence_set)}",
            f"predicate_sha256:{predicate.predicate_sha256}",
            f"observation_sha256:{canonical_sha256(observation_payload)}",
        )
        payload = dict(
            verification_record_id="vrs_" + "0" * 64,
            # Lineage is taken from the manifest itself.
            request_id=manifest.request_id,
            run_id=manifest.run_id,
            generation=manifest.generation,
            verifier_id=descriptor.verifier_id,
            verifier_version=descriptor.verifier_version,
            registry_snapshot_sha256=self._snapshot.snapshot_sha256,  # type: ignore[attr-defined]
            predicate_id=predicate.predicate_id,
            predicate_type=predicate.predicate_type,
            subject_kind="artifact",
            subject_identity=manifest.artifact_revision_id,
            evaluation_phase=evaluation_phase,
            status=status,
            enforcement="RECORD",
            reason_codes=reason_codes,
            evidence_refs=evidence_refs,
            evidence_sha256=canonical_sha256(list(evidence_refs)),
            producer_component_id=descriptor.producer_component_id,
            model_generated=False,
            evaluated_at_ms=evaluated_at_ms,
            result_sha256="0" * 64,
        )
        record = VerificationRecord(**payload).with_computed_sha256()
        return record.model_copy(
            update={
                "verification_record_id": derive_verification_record_id(
                    result_sha256=record.result_sha256
                )
            }
        )


__all__ = [
    "ArtifactContentOracle",
    "ContentUnparseable",
    "OracleSnapshotInvalid",
]
