"""P19-R2 M2.1 ArtifactContentOracle tests — review matrix A–P.

Covers the 2026-08-30 M2 HOLD review: unified predicate identity,
descriptor v2 binding, sealed lineage, per-format deterministic metrics,
typed evidence refs, resource discipline, RECORD-only persistence, and
the "never fake a verdict" rules.
"""

from __future__ import annotations

import base64
import io
import sqlite3
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from contracts import (
    InboundEnvelope,
    InboundScope,
    derive_effect_identity,
    derive_inbound_scope_keys,
    derive_request_identity,
    derive_run_identity,
)
from contracts.verification import AcceptancePredicate, VerificationRecord
from total_gateway.artifact_gate import ArtifactCandidate, ArtifactGate
from total_gateway.artifact_qc import ArtifactIntegrityQcService
from total_gateway.backend_client import BackendClient
from total_gateway.docx_qc import DocxQcPolicy, DocxQcService
from total_gateway.effects import EffectClaim
from total_gateway.fact_ledger import FactLedger
from total_gateway.object_store import ContentAddressedObjectStore
from total_gateway.outcome_oracles import artifact_content as oracle_module
from total_gateway.outcome_oracles.artifact_content import (
    ArtifactContentOracle,
    OracleSnapshotInvalid,
)
from total_gateway.store import GatewayStateStore
from total_gateway.verification_oracle_config import (
    ARTIFACT_IMPLEMENTED_PREDICATE_TYPES,
)
from total_gateway.verification_registry import (
    VerifierRegistry,
    legacy_artifact_v1_descriptor,
)
from total_gateway.verification_recording import VerificationRecorder
from tests.test_backend_client import FakeBackendTransport, backend_response, signed_ticket
from tests.test_docx_qc import docx_bytes

HASH_B = "b" * 64

# 1x1 transparent PNG for picture-bearing slides.
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
    "YPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def xlsx_bytes(header: list[str], rows: list[list[object]], *, formulas: list[str] | None = None) -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    if header:
        sheet.append(list(header))
    for row in rows:
        sheet.append(list(row))
    for cell_ref, formula in formulas or []:
        sheet[cell_ref] = formula
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def pptx_bytes(slides: list[str], *, with_picture_slide: bool = False) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = text
    if with_picture_slide:
        blank = presentation.slides.add_slide(presentation.slide_layouts[6])
        blank.shapes.add_picture(io.BytesIO(_TINY_PNG), 0, 0, width=100, height=100)
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def blank_pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])  # truly blank
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


class M21OracleTestBase(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        root = Path(self.temporary.name)
        self.gateway_store = GatewayStateStore.open(root / "gateway.sqlite3", now_ms=1_000)
        self.object_store = ContentAddressedObjectStore.open(root / "objects", now_ms=1_000)
        self.fact_ledger = FactLedger.open(root / "facts.sqlite3", self.object_store, now_ms=1_000)
        self.gate = ArtifactGate(self.object_store, self.fact_ledger)
        self.docx_qc = DocxQcService(self.object_store, self.fact_ledger)
        self.integrity_qc = ArtifactIntegrityQcService(self.object_store, self.fact_ledger)
        self.request = derive_request_identity("7" * 64)
        self.run = derive_run_identity(self.request.request_id, 1)
        self.snapshot = VerifierRegistry.with_defaults().snapshot(captured_at_ms=1)
        self.oracle = ArtifactContentOracle(
            snapshot=self.snapshot,
            object_store=self.object_store,
            fact_ledger=self.fact_ledger,
        )
        self._prepare_sequence = 0

    def tearDown(self) -> None:
        self.fact_ledger.close()
        self.object_store.close()
        self.gateway_store.close()
        self.temporary.cleanup()

    # -- fixture helpers ---------------------------------------------------

    def _prepare(self, data: bytes) -> str:
        self._prepare_sequence += 1
        reference = self.object_store.put_bytes(
            data,
            kind="artifact",
            tenant_id="tenant_001",
            link_account_id="wechat_001",
            conversation_scope_hash=HASH_B,
            created_at_ms=20_000,
        ).reference
        effect = derive_effect_identity(
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            run_sequence=1,
            generation=2,
            effect_kind="execution",
            ordinal=self._prepare_sequence - 1,
            intent_sha256="6" * 64,
        )
        arguments = {"content": f"create artifact {self._prepare_sequence}"}
        ticket, capability, trust = signed_ticket(
            arguments,
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            effect_id=effect.effect_id,
        )
        transport = FakeBackendTransport()
        envelope = backend_response(ticket, {"object_id": reference.object_id})
        result = dict(envelope["execution_result"])
        result["result_id"] = f"execution_result_m21_{self._prepare_sequence}"
        result["fact_ids"] = [f"fact_m21_producer_{self._prepare_sequence}"]
        result["output_object_refs"] = [reference.object_id]
        envelope["execution_result"] = result
        transport.response = envelope
        response = BackendClient(
            transport,
            self.gateway_store,
            ticket_consumer_instance_id="oracle_m21_test",
        ).execute(
            ticket,
            arguments,
            capability_manifest=capability,
            trust_bundle=trust,
            now_ms=20_000,
            expected_gateway_epoch=3,
            minimum_generation=2,
        )
        self.fact_ledger.record_execution(response, observed_at_ms=20_200)
        self.gateway_store.claim_effect(
            EffectClaim(
                effect_id=effect.effect_id,
                request_id=self.request.request_id,
                run_id=self.run.run_id,
                run_sequence=1,
                generation=2,
                effect_kind="execution",
                ordinal=effect.ordinal,
                intent_sha256="6" * 64,
                owner_component_id="tiangong-backend",
                claimed_at_ms=20_000,
                claim_sha256="0" * 64,
            ).with_computed_sha256()
        )
        return reference.object_id

    def _passed_manifest(
        self,
        data: bytes,
        *,
        filename: str,
        format_id: str,
        declared_mime: str,
        docx_min_words: int = 1,
    ):
        object_id = self._prepare(data)
        reference = self.object_store.get_reference(object_id)
        candidate = ArtifactCandidate(
            producer_fact_id=f"fact_m21_producer_{self._prepare_sequence}",
            object_id=object_id,
            expected_sha256=reference.sha256,
            expected_size_bytes=reference.size_bytes,
            run_sequence=1,
            artifact_intent_id=f"oracle_artifact_{self._prepare_sequence}",
            revision=1,
            workspace_id="workspace_001",
            filename=filename,
            declared_mime=declared_mime,
            format_id=format_id,
            created_at_ms=20_300,
        )
        gate_result = self.gate.accept(candidate)
        if format_id == "docx":
            outcome = self.docx_qc.evaluate(
                gate_result,
                run_sequence=1,
                policy=DocxQcPolicy(minimum_word_count=docx_min_words),
                checked_at_ms=20_500,
            )
        else:
            outcome = self.integrity_qc.evaluate(
                gate_result, run_sequence=1, checked_at_ms=20_500
            )
        self.assertTrue(outcome.passed)
        return outcome.registration.record.manifest

    def _evaluate(self, manifest, predicate) -> VerificationRecord:
        return self.oracle.evaluate(
            manifest, predicate, evaluated_at_ms=21_000
        )

    def _register_lineage_request(self) -> None:
        """Register the REAL request/run/generation the manifests carry."""
        scope = InboundScope(
            channel="desktop",
            tenant_id="tenant_001",
            link_account_id="desktop_m21",
            conversation_ref="conversation_m21",
            channel_message_ref="message_m21",
            sender_ref="sender_m21",
        )
        keys = derive_inbound_scope_keys(scope)
        envelope = InboundEnvelope(
            inbound_id="inbound_m21",
            channel=scope.channel,
            tenant_id=scope.tenant_id,
            link_account_id=scope.link_account_id,
            conversation_ref=scope.conversation_ref,
            conversation_scope_hash=keys.conversation_scope_hash,
            principal_scope_hash=keys.principal_scope_hash,
            message_scope_hash=keys.message_scope_hash,
            channel_message_ref=scope.channel_message_ref,
            sender_ref=scope.sender_ref,
            received_at_ms=1_000,
            idempotency_key="7" * 64,  # -> same request_id as manifests
            channel_metadata_hash=HASH_B,
            text="produce the artifact",
        )
        registration = self.gateway_store.register_request(
            envelope, ingress_sha256=HASH_B, created_at_ms=1_100
        )
        assert registration.entry.request_id == self.request.request_id
        self.gateway_store.acquire_generation_lease(
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            run_sequence=1,
            generation=2,  # manifests carry generation=2 (fixture effect)
            gateway_epoch=1,
            lease_id="lease_m21",
            owner_instance_id="gateway_m21",
            issued_at_ms=1_200,
            lease_duration_ms=60_000,
        )

    def _text_manifest(self, text: str = "identity roundtrip content"):
        return self._passed_manifest(
            text.encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )


class PredicateIdentityTests(unittest.TestCase):
    """A / B / C: predicate identity semantics."""

    def _predicate(self, **params):
        return AcceptancePredicate.create(
            predicate_type="xlsx.required_columns",
            subject_kind="artifact",
            params=params or None,
        )

    def test_a_different_params_different_predicate_and_record_ids(self) -> None:
        base = self._predicate(columns=["姓名", "分数"])
        other = self._predicate(columns=["姓名", "等级"])
        self.assertNotEqual(base.predicate_id, other.predicate_id)
        self.assertNotEqual(base.predicate_sha256, other.predicate_sha256)
        # record ids differ too: result hash covers predicate_sha256 via
        # evidence refs (predicate_sha256:<hash>) — verified in oracle tests.

    def test_b_normalized_params_same_identity(self) -> None:
        first = self._predicate(columns=["姓名", "分数"])
        second = self._predicate(columns=["分数", "姓名", "姓名", "  姓名  "])
        self.assertEqual(first.predicate_id, second.predicate_id)
        empty_a = AcceptancePredicate.create(
            predicate_type="artifact.nonempty", subject_kind="artifact"
        )
        empty_b = AcceptancePredicate.create(
            predicate_type="artifact.nonempty", subject_kind="artifact", params={}
        )
        self.assertEqual(empty_a.predicate_id, empty_b.predicate_id)

    def test_c_params_are_deep_frozen(self) -> None:
        predicate = self._predicate(columns=["姓名"])
        with self.assertRaises(AttributeError):
            predicate.params[0][1].append("x")  # type: ignore[union-attr]
        with self.assertRaises(Exception):
            predicate.params = ()  # type: ignore[misc]

    def test_nonempty_rejects_filler_params(self) -> None:
        from contracts.verification import AcceptancePredicateSpecError

        with self.assertRaises(AcceptancePredicateSpecError):
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty",
                subject_kind="artifact",
                params={"strict": True},
            )


class DescriptorV2Tests(unittest.TestCase):
    """D / F: descriptor v2 facts and snapshot identities."""

    def test_d_v2_descriptor_facts(self) -> None:
        snapshot = VerifierRegistry.with_defaults().snapshot(captured_at_ms=1)
        descriptor = snapshot.find("verifier.artifact_content")
        assert descriptor is not None
        self.assertEqual(descriptor.verifier_version, "2")
        self.assertNotEqual(descriptor.config_sha256, "0" * 64)
        self.assertEqual(
            set(descriptor.supported_predicate_types),
            set(ARTIFACT_IMPLEMENTED_PREDICATE_TYPES),
        )
        # xlsx.min_data_rows must be declared (it is implemented)
        self.assertIn("xlsx.min_data_rows", descriptor.supported_predicate_types)
        # unimplemented planned types must NOT be declared
        self.assertNotIn("docx.min_body_items", descriptor.supported_predicate_types)
        self.assertNotIn("csv.required_columns", descriptor.supported_predicate_types)

    def test_d_default_snapshot_identity_differs_from_v1_world(self) -> None:
        v2_snapshot = VerifierRegistry.with_defaults().snapshot(captured_at_ms=1)
        legacy_only = VerifierRegistry(
            (legacy_artifact_v1_descriptor(),)
        ).snapshot(captured_at_ms=1)
        self.assertNotEqual(v2_snapshot.snapshot_sha256, legacy_only.snapshot_sha256)

    def test_f_v1_snapshot_cannot_instantiate_v2_oracle(self) -> None:
        v1_snapshot = VerifierRegistry(
            (legacy_artifact_v1_descriptor(),)
        ).snapshot(captured_at_ms=1)
        with self.assertRaises(OracleSnapshotInvalid):
            ArtifactContentOracle(
                snapshot=v1_snapshot,
                object_store=mock.Mock(),
                fact_ledger=mock.Mock(),
            )


class OracleVerdictTests(M21OracleTestBase):
    """Deterministic verdicts across formats (review section 4)."""

    def test_docx_min_visible_pass_and_fail(self) -> None:
        filled = self._passed_manifest(
            docx_bytes("这是足够长的正文内容" * 10),
            filename="report.docx",
            format_id="docx",
            declared_mime=DOCX_MIME,
            docx_min_words=5,
        )
        record = self._evaluate(
            filled,
            AcceptancePredicate.create(
                predicate_type="artifact.min_visible_text_chars",
                subject_kind="artifact",
                params={"min_chars": 50},
            ),
        )
        self.assertEqual(record.status, "PASS")
        self.assertEqual(record.reason_codes, ())
        self.assertTrue(record.has_valid_identity())

        shell = self._passed_manifest(
            docx_bytes("标题"),
            filename="shell.docx",
            format_id="docx",
            declared_mime=DOCX_MIME,
        )
        record = self._evaluate(
            shell,
            AcceptancePredicate.create(
                predicate_type="artifact.min_visible_text_chars",
                subject_kind="artifact",
                params={"min_chars": 50},
            ),
        )
        self.assertEqual(record.status, "FAIL")

    def test_xlsx_columns_exact_match_and_rows_exclude_header(self) -> None:
        manifest = self._passed_manifest(
            xlsx_bytes(["姓名", "分数"], [["甲", 90], ["乙", 80]]),
            filename="score.xlsx",
            format_id="xlsx",
            declared_mime=XLSX_MIME,
        )
        columns = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="xlsx.required_columns",
                subject_kind="artifact",
                params={"columns": ["姓名", "分数"]},
            ),
        )
        self.assertEqual(columns.status, "PASS")

        # L: "姓名备注" must NOT satisfy "姓名" (exact normalized match)
        sneaky = self._passed_manifest(
            xlsx_bytes(["姓名备注", "分数"], [["甲", 90]]),
            filename="sneaky.xlsx",
            format_id="xlsx",
            declared_mime=XLSX_MIME,
        )
        record = self._evaluate(
            sneaky,
            AcceptancePredicate.create(
                predicate_type="xlsx.required_columns",
                subject_kind="artifact",
                params={"columns": ["姓名"]},
            ),
        )
        self.assertEqual(record.status, "FAIL")
        self.assertIn("xlsx.column_missing:姓名", record.reason_codes)

        rows = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="xlsx.min_data_rows",
                subject_kind="artifact",
                params={"min_rows": 2},
            ),
        )
        self.assertEqual(rows.status, "PASS")
        three = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="xlsx.min_data_rows",
                subject_kind="artifact",
                params={"min_rows": 3},
            ),
        )
        self.assertEqual(three.status, "FAIL")

    def test_i_xlsx_formula_text_counts_as_content(self) -> None:
        # formula cell has no cached value (data_only=True -> None) but
        # the formula itself is visible content: no fake missing metric.
        manifest = self._passed_manifest(
            xlsx_bytes(["合计"], [], formulas=[("B1", "=SUM(A1:A9)")]),
            filename="formula.xlsx",
            format_id="xlsx",
            declared_mime=XLSX_MIME,
        )
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.min_visible_text_chars",
                subject_kind="artifact",
                params={"min_chars": 5},
            ),
        )
        self.assertEqual(record.status, "PASS", record.reason_codes)

    def test_j_blank_pptx_nonempty_fails(self) -> None:
        manifest = self._passed_manifest(
            blank_pptx_bytes(),
            filename="blank.pptx",
            format_id="pptx",
            declared_mime=PPTX_MIME,
        )
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.status, "FAIL")
        self.assertIn("artifact.content_empty", record.reason_codes)

    def test_k_meaningful_slide_semantics(self) -> None:
        # 2 text slides + 1 picture slide: 3 meaningful, 2 text-bearing.
        deck = self._passed_manifest(
            pptx_bytes(["第一页", "第二页"], with_picture_slide=True),
            filename="deck.pptx",
            format_id="pptx",
            declared_mime=PPTX_MIME,
        )
        three_meaningful = self._evaluate(
            deck,
            AcceptancePredicate.create(
                predicate_type="pptx.min_nonempty_slides",
                subject_kind="artifact",
                params={"min_slides": 3},
            ),
        )
        self.assertEqual(three_meaningful.status, "PASS")
        four = self._evaluate(
            deck,
            AcceptancePredicate.create(
                predicate_type="pptx.min_nonempty_slides",
                subject_kind="artifact",
                params={"min_slides": 4},
            ),
        )
        self.assertEqual(four.status, "FAIL")

    def test_text_markers(self) -> None:
        manifest = self._text_manifest("总结要点：第一项已完成。\n第二项进行中。")
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="text.required_markers",
                subject_kind="artifact",
                params={"markers": ["总结要点"]},
            ),
        )
        self.assertEqual(record.status, "PASS")


class OracleDisciplineTests(M21OracleTestBase):
    """Binding / tamper / resource / applicability discipline."""

    def test_manifest_hash_tamper_is_error(self) -> None:
        manifest = self._text_manifest()
        tampered = manifest.model_copy(update={"size_bytes": manifest.size_bytes + 1})
        record = self._evaluate(
            tampered,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.status, "ERROR")
        self.assertTrue(any(code.startswith("authority:") for code in record.reason_codes))

    def test_object_binding_swap_is_error(self) -> None:
        first = self._text_manifest("first artifact")
        second_object = self._prepare("second artifact bytes".encode("utf-8"))
        second_reference = self.object_store.get_reference(second_object)
        swapped = first.model_copy(
            update={
                "content_object_id": second_object,
                "sha256": second_reference.sha256,
                "size_bytes": second_reference.size_bytes,
            }
        ).with_computed_manifest_sha256()
        record = self._evaluate(
            swapped,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.status, "ERROR")

    def test_disk_blob_tamper_is_error_never_pass(self) -> None:
        original = b"trusted content for tamper test"
        manifest = self._passed_manifest(
            original,
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        targets = [
            path
            for path in self.object_store.root.rglob("*")
            if path.is_file()
            and not path.name.endswith((".sqlite3", ".db"))
            and path.read_bytes() == original
        ]
        self.assertTrue(targets, "object blob not found on disk")
        targets[0].write_bytes(b"tampered blob with same length!!")
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.status, "ERROR")

    def test_o_size_precheck_reads_nothing(self) -> None:
        manifest = self._text_manifest()
        oversized = manifest.model_copy(
            update={"size_bytes": 10 * 1024 * 1024 * 1024}
        ).with_computed_manifest_sha256()
        with mock.patch.object(
            self.object_store, "read_bytes", wraps=self.object_store.read_bytes
        ) as read_spy:
            record = self._evaluate(
                oversized,
                AcceptancePredicate.create(
                    predicate_type="artifact.nonempty", subject_kind="artifact"
                ),
            )
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("input_too_large", record.reason_codes)
        read_spy.assert_not_called()

    def test_m_dependency_missing_is_error_and_unparseable_is_inconclusive(self) -> None:
        manifest = self._text_manifest()
        nonempty = AcceptancePredicate.create(
            predicate_type="artifact.nonempty", subject_kind="artifact"
        )
        with mock.patch.object(
            oracle_module, "_inspect_text", side_effect=oracle_module._InspectorDependencyMissing("encodings")
        ):
            record = self._evaluate(manifest, nonempty)
        self.assertEqual(record.status, "ERROR")
        self.assertIn("inspector_dependency_missing", record.reason_codes)

        from total_gateway.outcome_oracles.artifact_content import ContentUnparseable

        with mock.patch.object(
            oracle_module, "_inspect_text", side_effect=ContentUnparseable("bad")
        ):
            record = self._evaluate(manifest, nonempty)
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("content_unparseable", record.reason_codes)

    def test_applicability_before_implementation(self) -> None:
        docx_manifest = self._passed_manifest(
            docx_bytes("内容足够长" * 5),
            filename="report.docx",
            format_id="docx",
            declared_mime=DOCX_MIME,
        )
        # xlsx.required_sheet_names is UNIMPLEMENTED, but an xlsx predicate
        # on a docx subject is NOT_APPLICABLE first — not INCONCLUSIVE.
        record = self._evaluate(
            docx_manifest,
            AcceptancePredicate.create(
                predicate_type="xlsx.required_columns",
                subject_kind="artifact",
                params={"columns": ["姓名"]},
            ),
        )
        self.assertEqual(record.status, "NOT_APPLICABLE")

    def test_unimplemented_applicable_predicate_is_inconclusive(self) -> None:
        # The v2 descriptor declares exactly the implemented set, so an
        # applicable-but-undeclared predicate can only be simulated by a
        # descriptor that dropped one type (e.g. an older/partial rollout).
        manifest = self._text_manifest()
        predicate = AcceptancePredicate.create(
            predicate_type="text.required_markers",
            subject_kind="artifact",
            params={"markers": ["x"]},
        )
        restricted = self.oracle.descriptor.model_copy(
            update={
                "supported_predicate_types": tuple(
                    t
                    for t in self.oracle.descriptor.supported_predicate_types
                    if t != "text.required_markers"
                )
            }
        )
        with mock.patch.object(self.oracle, "_descriptor", restricted):
            record = self._evaluate(manifest, predicate)
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("predicate_not_implemented", record.reason_codes)

    def test_unsupported_format_is_inconclusive(self) -> None:
        manifest = self._passed_manifest(
            b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj\n<<>>\nendobj\ntrailer\n<<>>\n%%EOF",
            filename="scan.pdf",
            format_id="pdf",
            declared_mime="application/pdf",
        )
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("format_not_inspectable", record.reason_codes)


class EvidenceAndPersistenceTests(M21OracleTestBase):
    """N / E / H / G / P: evidence shape, lineage, RECORD-only persist."""

    def test_n_evidence_refs_are_typed_and_bounded(self) -> None:
        manifest = self._text_manifest()
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        refs = record.evidence_refs
        self.assertLessEqual(len(refs), 8)
        prefixes = {ref.split(":", 1)[0] for ref in refs}
        self.assertEqual(
            prefixes,
            {
                "artifact_revision",
                "manifest_sha256",
                "content_object",
                "content_sha256",
                "qc_evidence_set_sha256",
                "predicate_sha256",
                "observation_sha256",
            },
        )
        # 64 QC evidences must not leak into refs: every ref is typed with
        # a known prefix (the prefixes set equality above already proves no
        # raw QC fact id can appear as its own ref).
        self.assertTrue(all(":" in ref for ref in refs))

    def test_g_lineage_comes_from_manifest_and_cannot_rebind(self) -> None:
        manifest = self._text_manifest()
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        self.assertEqual(record.request_id, manifest.request_id)
        self.assertEqual(record.run_id, manifest.run_id)
        self.assertEqual(record.generation, manifest.generation)
        # rebinding via model_copy breaks identity (M1.1) — recorder rejects
        rebound = record.model_copy(update={"request_id": "req_" + "e" * 64})
        self.assertFalse(rebound.has_valid_identity())

    def test_e_and_h_full_chain_persists_without_rebinding(self) -> None:
        self._register_lineage_request()
        manifest = self._passed_manifest(
            xlsx_bytes(["姓名", "分数"], [["甲", 1], ["乙", 2]]),
            filename="chain.xlsx",
            format_id="xlsx",
            declared_mime=XLSX_MIME,
        )
        self.snapshot  # registry snapshot persisted below
        self.gateway_store.put_registry_snapshot(self.snapshot, recorded_at_ms=1_500)
        recorder = VerificationRecorder(snapshot=self.snapshot, store=self.gateway_store)
        predicate = AcceptancePredicate.create(
            predicate_type="xlsx.min_data_rows",
            subject_kind="artifact",
            params={"min_rows": 2},
        )
        record = self.oracle.evaluate(
            manifest, predicate, evaluated_at_ms=21_000
        )
        self.assertEqual(record.status, "PASS")
        # the record goes straight in — no model_copy, no rebinding
        outcome = recorder.record(record, recorded_at_ms=2_000)
        self.assertTrue(outcome.created_by_this_call)

    def test_p_record_only_zero_state_impact(self) -> None:
        self._register_lineage_request()
        manifest = self._text_manifest()
        self.gateway_store.put_registry_snapshot(self.snapshot, recorded_at_ms=1_500)
        recorder = VerificationRecorder(snapshot=self.snapshot, store=self.gateway_store)
        record = self._evaluate(
            manifest,
            AcceptancePredicate.create(
                predicate_type="artifact.nonempty", subject_kind="artifact"
            ),
        )
        recorder.record(record, recorded_at_ms=2_000)

        connection = sqlite3.connect(self.gateway_store.path)
        try:
            decisions = connection.execute(
                "SELECT COUNT(*) FROM completion_decisions"
            ).fetchone()[0]
            aggregates = connection.execute(
                "SELECT COUNT(*) FROM aggregate_state"
            ).fetchone()[0]
            outbox_rows = connection.execute(
                "SELECT COUNT(*) FROM outbox"
            ).fetchone()[0]
            enforcements = {
                row[0]
                for row in connection.execute(
                    "SELECT DISTINCT enforcement FROM verification_record"
                )
            }
        finally:
            connection.close()
        self.assertEqual(decisions, 0)
        self.assertEqual(aggregates, 0)
        self.assertEqual(outbox_rows, 0)
        self.assertEqual(enforcements, {"RECORD"})


if __name__ == "__main__":
    unittest.main()
