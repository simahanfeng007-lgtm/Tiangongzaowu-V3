"""P19-R2 M2 ArtifactContentOracle tests — review completion gate.

Categories:
1. legal manifest + legal immutable bytes -> normal verdicts;
2. manifest/object/revision/hash binding corruption -> ERROR, never a
   content verdict;
3. object-store bytes tampered -> ERROR, never PASS;
4. deterministic PASS/FAIL for docx/xlsx/pptx/txt/csv predicates;
5. unparseable/unsupported -> INCONCLUSIVE (never FAIL-by-emptiness);
6. predicate/subject mismatch -> NOT_APPLICABLE;
7. oracle exceptions are never silent (ERROR);
8. produced records satisfy M1.1 identity + enforcement=RECORD only;
9. recorder/store roundtrip changes no CompletionDecision/state;
10. unimplemented predicate types -> explicit INCONCLUSIVE.
"""

from __future__ import annotations

import io
import sqlite3
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from contracts import (
    canonical_sha256,
    derive_effect_identity,
    derive_inbound_scope_keys,
    derive_request_identity,
    derive_run_identity,
    InboundEnvelope,
    InboundScope,
)
from contracts.verification import (
    VerificationRecord,
    derive_verification_record_id,
)
from total_gateway.artifact_gate import ArtifactCandidate, ArtifactGate
from total_gateway.artifact_qc import ArtifactIntegrityQcService
from total_gateway.backend_client import BackendClient
from total_gateway.docx_qc import DocxQcPolicy, DocxQcService
from total_gateway.effects import EffectClaim
from total_gateway.fact_ledger import FactLedger
from total_gateway.object_store import ContentAddressedObjectStore
from total_gateway.outcome_oracles.artifact_content import (
    ArtifactContentOracle,
    ArtifactPredicate,
    ArtifactPredicateSpecError,
    _inspect_docx,
)
from total_gateway.store import GatewayStateStore
from total_gateway.verification_registry import VerifierRegistry
from total_gateway.verification_recording import VerificationRecorder
from tests.test_backend_client import FakeBackendTransport, backend_response, signed_ticket
from tests.test_docx_qc import docx_bytes

HASH_B = "b" * 64


def xlsx_bytes(
    header: list[str],
    rows: list[list[object]],
) -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    if header:
        sheet.append(list(header))
    for row in rows:
        sheet.append(list(row))
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def pptx_bytes(slide_texts: list[str]) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for text in slide_texts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = text
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


class ArtifactContentOracleTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        root = Path(self.temporary.name)
        self.gateway_store = GatewayStateStore.open(root / "gateway.sqlite3", now_ms=1_000)
        self.object_store = ContentAddressedObjectStore.open(root / "objects", now_ms=1_000)
        self.fact_ledger = FactLedger.open(root / "facts.sqlite3", self.object_store, now_ms=1_000)
        self.gate = ArtifactGate(self.object_store, self.fact_ledger)
        self.docx_qc = DocxQcService(self.object_store, self.fact_ledger)
        self.integrity_qc = ArtifactIntegrityQcService(self.object_store, self.fact_ledger)
        self.request = derive_request_identity("7" * 64)
        self.run = derive_run_identity(self.request.request_id, 1)
        self.effect = derive_effect_identity(
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            run_sequence=1,
            generation=2,
            effect_kind="execution",
            ordinal=0,
            intent_sha256="6" * 64,
        )
        self.snapshot = VerifierRegistry.with_defaults().snapshot(captured_at_ms=1)
        self._prepare_sequence = 0
        self.oracle = ArtifactContentOracle(
            object_store=self.object_store,
            fact_ledger=self.fact_ledger,
            registry_snapshot_sha256=self.snapshot.snapshot_sha256,
        )

    def tearDown(self) -> None:
        self.fact_ledger.close()
        self.object_store.close()
        self.gateway_store.close()
        self.temporary.cleanup()

    # -- fixture helpers ---------------------------------------------------

    def _prepare(self, data: bytes, *, sequence: int = 1) -> str:
        self._prepare_sequence += 1
        reference = self.object_store.put_bytes(
            data,
            kind="artifact",
            tenant_id="tenant_001",
            link_account_id="wechat_001",
            conversation_scope_hash=HASH_B,
            created_at_ms=20_000,
        ).reference
        # one execution effect per artifact: fact identities are keyed on
        # the effect, so reused effects with different evidence conflict.
        effect = derive_effect_identity(
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            run_sequence=1,
            generation=2,
            effect_kind="execution",
            ordinal=self._prepare_sequence - 1,
            intent_sha256="6" * 64,
        )
        arguments = {"content": f"create artifact for oracle test {self._prepare_sequence}"}
        ticket, capability, trust = signed_ticket(
            arguments,
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            effect_id=effect.effect_id,
        )
        transport = FakeBackendTransport()
        envelope = backend_response(ticket, {"object_id": reference.object_id})
        result = dict(envelope["execution_result"])
        result["result_id"] = f"execution_result_oracle_m2_{self._prepare_sequence}"
        result["fact_ids"] = [f"fact_oracle_producer_{self._prepare_sequence}"]
        result["output_object_refs"] = [reference.object_id]
        envelope["execution_result"] = result
        transport.response = envelope
        response = BackendClient(
            transport,
            self.gateway_store,
            ticket_consumer_instance_id="oracle_m2_test",
        ).execute(
            ticket,
            arguments,
            capability_manifest=capability,
            trust_bundle=trust,
            now_ms=20_000,
            expected_gateway_epoch=3,
            minimum_generation=2,
        )
        self.fact_ledger.record_execution(response, observed_at_ms=20_200)
        self.gateway_store.claim_effect(
            EffectClaim(
                effect_id=effect.effect_id,
                request_id=self.request.request_id,
                run_id=self.run.run_id,
                run_sequence=1,
                generation=2,
                effect_kind="execution",
                ordinal=effect.ordinal,
                intent_sha256="6" * 64,
                owner_component_id="tiangong-backend",
                claimed_at_ms=20_000,
                claim_sha256="0" * 64,
            ).with_computed_sha256()
        )
        return reference.object_id

    def _passed_manifest(
        self,
        data: bytes,
        *,
        filename: str,
        format_id: str,
        declared_mime: str,
        docx_min_words: int = 1,
    ):
        object_id = self._prepare(data)
        reference = self.object_store.get_reference(object_id)
        candidate = ArtifactCandidate(
            producer_fact_id=f"fact_oracle_producer_{self._prepare_sequence}",
            object_id=object_id,
            expected_sha256=reference.sha256,
            expected_size_bytes=reference.size_bytes,
            run_sequence=1,
            artifact_intent_id=f"oracle_artifact_{self._prepare_sequence}",
            revision=1,
            workspace_id="workspace_001",
            filename=filename,
            declared_mime=declared_mime,
            format_id=format_id,
            created_at_ms=20_300,
        )
        gate_result = self.gate.accept(candidate)
        if format_id == "docx":
            outcome = self.docx_qc.evaluate(
                gate_result,
                run_sequence=1,
                policy=DocxQcPolicy(minimum_word_count=docx_min_words),
                checked_at_ms=20_500,
            )
        else:
            outcome = self.integrity_qc.evaluate(
                gate_result, run_sequence=1, checked_at_ms=20_500
            )
        self.assertTrue(outcome.passed)
        return outcome.registration.record.manifest

    def _evaluate(self, manifest, predicate) -> VerificationRecord:
        return self.oracle.evaluate(
            manifest,
            predicate,
            request_id=self.request.request_id,
            run_id=self.run.run_id,
            generation=2,
            evaluated_at_ms=21_000,
        )

    # -- category 1+4: legal chain and deterministic verdicts --------------

    def test_docx_min_visible_text_chars_pass_and_fail(self) -> None:
        filled = self._passed_manifest(
            docx_bytes("这是足够长的正文内容" * 10),
            filename="report.docx",
            format_id="docx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".wordprocessingml.document"
            ),
            docx_min_words=5,
        )
        record = self._evaluate(
            filled,
            ArtifactPredicate(
                predicate_id="vpd_docx_chars",
                predicate_type="artifact.min_visible_text_chars",
                params={"min_chars": 50},
            ),
        )
        self.assertEqual(record.status, "PASS")
        self.assertEqual(record.reason_codes, ())
        self.assertTrue(record.has_valid_identity())

        shell = self._passed_manifest(
            docx_bytes("标题"),
            filename="shell.docx",
            format_id="docx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".wordprocessingml.document"
            ),
        )
        record = self._evaluate(
            shell,
            ArtifactPredicate(
                predicate_id="vpd_docx_chars",
                predicate_type="artifact.min_visible_text_chars",
                params={"min_chars": 50},
            ),
        )
        self.assertEqual(record.status, "FAIL")
        self.assertIn("artifact.visible_text_below_minimum", record.reason_codes)
        self.assertTrue(record.has_valid_identity())

    def test_xlsx_required_columns_and_min_data_rows_exclude_header(self) -> None:
        manifest = self._passed_manifest(
            xlsx_bytes(["姓名", "分数"], [["甲", 90], ["乙", 80]]),
            filename="score.xlsx",
            format_id="xlsx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".spreadsheetml.sheet"
            ),
        )
        columns = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_xlsx_cols",
                predicate_type="xlsx.required_columns",
                params={"columns": ["姓名", "分数"]},
            ),
        )
        self.assertEqual(columns.status, "PASS")

        missing = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_xlsx_cols",
                predicate_type="xlsx.required_columns",
                params={"columns": ["姓名", "等级"]},
            ),
        )
        self.assertEqual(missing.status, "FAIL")
        self.assertIn("xlsx.column_missing:等级", missing.reason_codes)

        # header + 2 data rows: min_data_rows=2 PASS, =3 FAIL (header
        # excluded — the M2-corrected semantics the legacy preflight lacks)
        two = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_xlsx_rows",
                predicate_type="xlsx.min_data_rows",
                params={"min_rows": 2},
            ),
        )
        self.assertEqual(two.status, "PASS")
        three = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_xlsx_rows",
                predicate_type="xlsx.min_data_rows",
                params={"min_rows": 3},
            ),
        )
        self.assertEqual(three.status, "FAIL")

    def test_xlsx_nonempty_and_text_csv_pptx_predicates(self) -> None:
        empty_sheet = self._passed_manifest(
            xlsx_bytes([], []),
            filename="empty.xlsx",
            format_id="xlsx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".spreadsheetml.sheet"
            ),
        )
        record = self._evaluate(
            empty_sheet,
            ArtifactPredicate(
                predicate_id="vpd_nonempty",
                predicate_type="artifact.nonempty",
                params={"strict": True},
            ),
        )
        self.assertEqual(record.status, "FAIL")

        text_manifest = self._passed_manifest(
            "总结要点：第一项已完成。\n第二项进行中。".encode("utf-8"),
            filename="summary.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        markers = self._evaluate(
            text_manifest,
            ArtifactPredicate(
                predicate_id="vpd_markers",
                predicate_type="text.required_markers",
                params={"markers": ["总结要点"]},
            ),
        )
        self.assertEqual(markers.status, "PASS")

        # NOTE: csv.required_columns is deliberately not implemented — the
        # artifact gate has no csv format policy (.csv is not in the text
        # extension allowlist), so no .csv manifest can reach this oracle
        # with authority. See module docstring; capability stays honest.

        deck = self._passed_manifest(
            pptx_bytes(["第一页标题", "第二页标题", "第三页标题"]),
            filename="deck.pptx",
            format_id="pptx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )
        slides = self._evaluate(
            deck,
            ArtifactPredicate(
                predicate_id="vpd_pptx",
                predicate_type="pptx.min_nonempty_slides",
                params={"min_slides": 3},
            ),
        )
        self.assertEqual(slides.status, "PASS")
        four = self._evaluate(
            deck,
            ArtifactPredicate(
                predicate_id="vpd_pptx",
                predicate_type="pptx.min_nonempty_slides",
                params={"min_slides": 4},
            ),
        )
        self.assertEqual(four.status, "FAIL")

    # -- category 2: binding corruption -> ERROR ---------------------------

    def test_manifest_hash_tamper_is_error_not_content_verdict(self) -> None:
        manifest = self._passed_manifest(
            "plain content".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        tampered = manifest.model_copy(update={"size_bytes": manifest.size_bytes + 1})
        self.assertFalse(tampered.has_valid_manifest_sha256())
        record = self._evaluate(
            tampered,
            ArtifactPredicate(
                predicate_id="vpd_markers",
                predicate_type="text.required_markers",
                params={"markers": ["plain"]},
            ),
        )
        self.assertEqual(record.status, "ERROR")
        self.assertTrue(any(code.startswith("authority:") for code in record.reason_codes))

    def test_object_binding_swap_is_error(self) -> None:
        first = self._passed_manifest(
            "first artifact".encode("utf-8"),
            filename="one.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        second_object = self._prepare("second artifact bytes".encode("utf-8"))
        second_reference = self.object_store.get_reference(second_object)
        swapped = first.model_copy(
            update={
                "content_object_id": second_object,
                "sha256": second_reference.sha256,
                "size_bytes": second_reference.size_bytes,
            }
        ).with_computed_manifest_sha256()
        self.assertTrue(swapped.has_valid_manifest_sha256())
        record = self._evaluate(
            swapped,
            ArtifactPredicate(
                predicate_id="vpd_markers",
                predicate_type="text.required_markers",
                params={"markers": ["first"]},
            ),
        )
        # QC facts were registered against the original manifest binding.
        self.assertEqual(record.status, "ERROR")

    # -- category 3: bytes tampered -> ERROR -------------------------------

    def test_object_store_readback_tamper_is_error_never_pass(self) -> None:
        manifest = self._passed_manifest(
            "trusted content".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        original = "trusted content".encode("utf-8")
        blob_files = [
            path
            for path in self.object_store.root.rglob("*")
            if path.is_file() and not path.name.endswith((".sqlite3", ".db"))
        ]
        targets = [path for path in blob_files if path.read_bytes() == original]
        self.assertTrue(targets, "object blob not found on disk")
        targets[0].write_bytes(b"tampered blob with same length!!")
        record = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_markers",
                predicate_type="text.required_markers",
                params={"markers": ["trusted"]},
            ),
        )
        self.assertEqual(record.status, "ERROR")
        self.assertTrue(
            any("authority" in code for code in record.reason_codes),
            record.reason_codes,
        )

    # -- category 5: unparseable/unsupported -> INCONCLUSIVE ---------------

    def test_unsupported_format_is_inconclusive(self) -> None:
        manifest = self._passed_manifest(
            b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj\n<<>>\nendobj\ntrailer\n<<>>\n%%EOF",
            filename="scan.pdf",
            format_id="pdf",
            declared_mime="application/pdf",
        )
        record = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_pdf",
                predicate_type="artifact.nonempty",
                params={"strict": True},
            ),
        )
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("format_not_inspectable", record.reason_codes)

    def test_unparseable_container_is_inconclusive_not_fail(self) -> None:
        manifest = self._passed_manifest(
            "decodable text".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        predicate = ArtifactPredicate(
            predicate_id="vpd_markers",
            predicate_type="text.required_markers",
            params={"markers": ["decodable"]},
        )
        from total_gateway.outcome_oracles.artifact_content import ContentUnparseable

        with mock.patch(
            "total_gateway.outcome_oracles.artifact_content._inspect_text",
            side_effect=ContentUnparseable("text not decodable as utf-8"),
        ):
            record = self._evaluate(manifest, predicate)
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("content_unparseable", record.reason_codes)

    # -- category 6: NOT_APPLICABLE ----------------------------------------

    def test_predicate_subject_mismatch_is_not_applicable(self) -> None:
        manifest = self._passed_manifest(
            docx_bytes("内容足够长" * 5),
            filename="report.docx",
            format_id="docx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".wordprocessingml.document"
            ),
        )
        record = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_mismatch",
                predicate_type="xlsx.required_columns",
                params={"columns": ["姓名"]},
            ),
        )
        self.assertEqual(record.status, "NOT_APPLICABLE")
        self.assertEqual(record.reason_codes, ())

        csv_on_txt = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_mismatch",
                predicate_type="csv.required_columns",
                params={"columns": ["name"]},
            ),
        )
        # gate 无 csv 格式策略：csv 谓词按未实现处理（诚实降级）
        self.assertEqual(csv_on_txt.status, "INCONCLUSIVE")
        self.assertIn("predicate_not_implemented", csv_on_txt.reason_codes)

    # -- category 7: never silent -------------------------------------------

    def test_inspector_runtime_failure_is_error_not_silent(self) -> None:
        manifest = self._passed_manifest(
            "some text".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        predicate = ArtifactPredicate(
            predicate_id="vpd_markers",
            predicate_type="text.required_markers",
            params={"markers": ["some"]},
        )
        real_inspect = _inspect_docx  # prove the module namespace is patchable
        self.assertTrue(callable(real_inspect))
        with mock.patch(
            "total_gateway.outcome_oracles.artifact_content._inspect_text",
            side_effect=RuntimeError("unexpected runtime crash"),
        ):
            # RuntimeError is not ContentUnparseable -> ERROR path
            record = self._evaluate(manifest, predicate)
        self.assertEqual(record.status, "ERROR")
        self.assertIn("inspector_failure", record.reason_codes)

    # -- category 10: honest capability surface -----------------------------

    def test_unimplemented_predicate_is_inconclusive(self) -> None:
        manifest = self._passed_manifest(
            xlsx_bytes(["a"], [["1"]]),
            filename="data.xlsx",
            format_id="xlsx",
            declared_mime=(
                "application/vnd.openxmlformats-officedocument"
                ".spreadsheetml.sheet"
            ),
        )
        record = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_sheets",
                predicate_type="xlsx.required_sheet_names",
                params={"sheets": ["Sheet"]},
            ),
        )
        self.assertEqual(record.status, "INCONCLUSIVE")
        self.assertIn("predicate_not_implemented", record.reason_codes)

    def test_malformed_params_raise_spec_error(self) -> None:
        manifest = self._passed_manifest(
            "text".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        with self.assertRaises(ArtifactPredicateSpecError):
            self._evaluate(
                manifest,
                ArtifactPredicate(
                    predicate_id="vpd_bad",
                    predicate_type="text.required_markers",
                    params={"markers": [""]},
                ),
            )

    # -- categories 8+9: identity + RECORD-only persistence ------------------

    def test_records_carry_identity_and_persist_record_only(self) -> None:
        manifest = self._passed_manifest(
            "identity roundtrip content".encode("utf-8"),
            filename="note.txt",
            format_id="text",
            declared_mime="text/plain",
        )
        record = self._evaluate(
            manifest,
            ArtifactPredicate(
                predicate_id="vpd_markers",
                predicate_type="text.required_markers",
                params={"markers": ["roundtrip"]},
            ),
        )
        self.assertEqual(record.enforcement, "RECORD")
        self.assertTrue(record.has_valid_identity())
        self.assertEqual(record.verifier_id, "verifier.artifact_content")

        # persist through the real recorder/store pair (seeded request)
        scope = InboundScope(
            channel="desktop",
            tenant_id="tenant_001",
            link_account_id="desktop_m2",
            conversation_ref="conversation_m2",
            channel_message_ref="message_m2",
            sender_ref="sender_m2",
        )
        keys = derive_inbound_scope_keys(scope)
        envelope = InboundEnvelope(
            inbound_id="inbound_m2",
            channel=scope.channel,
            tenant_id=scope.tenant_id,
            link_account_id=scope.link_account_id,
            conversation_ref=scope.conversation_ref,
            conversation_scope_hash=keys.conversation_scope_hash,
            principal_scope_hash=keys.principal_scope_hash,
            message_scope_hash=keys.message_scope_hash,
            channel_message_ref=scope.channel_message_ref,
            sender_ref=scope.sender_ref,
            received_at_ms=1_000,
            idempotency_key=keys.idempotency_key,
            channel_metadata_hash=HASH_B,
            text="produce the note",
        )
        registration = self.gateway_store.register_request(
            envelope, ingress_sha256=HASH_B, created_at_ms=1_100
        )
        request_id = registration.entry.request_id
        run_id = derive_run_identity(request_id, 1).run_id
        self.gateway_store.acquire_generation_lease(
            request_id=request_id,
            run_id=run_id,
            run_sequence=1,
            generation=1,
            gateway_epoch=1,
            lease_id="lease_m2",
            owner_instance_id="gateway_m2",
            issued_at_ms=1_200,
            lease_duration_ms=60_000,
        )
        self.gateway_store.put_registry_snapshot(self.snapshot, recorded_at_ms=1_500)
        recorder = VerificationRecorder(snapshot=self.snapshot, store=self.gateway_store)
        bound = record.model_copy(
            update={
                "request_id": request_id,
                "run_id": run_id,
                "generation": 1,
            }
        ).with_computed_sha256()
        bound = bound.model_copy(
            update={
                "verification_record_id": derive_verification_record_id(
                    result_sha256=bound.result_sha256
                )
            }
        )
        outcome = recorder.record(bound, recorded_at_ms=2_000)
        self.assertTrue(outcome.created_by_this_call)

        connection = sqlite3.connect(self.gateway_store.path)
        try:
            decisions = connection.execute(
                "SELECT COUNT(*) FROM completion_decisions"
            ).fetchone()[0]
            aggregates = connection.execute(
                "SELECT COUNT(*) FROM aggregate_state"
            ).fetchone()[0]
            enforcements = {
                row[0]
                for row in connection.execute(
                    "SELECT DISTINCT enforcement FROM verification_record"
                )
            }
        finally:
            connection.close()
        self.assertEqual(decisions, 0)  # zero CompletionDecision changes
        self.assertEqual(aggregates, 0)  # zero request state transitions
        self.assertEqual(enforcements, {"RECORD"})  # RECORD only


if __name__ == "__main__":
    unittest.main()
