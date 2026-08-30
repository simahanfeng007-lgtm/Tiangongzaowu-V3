"""P19-R2 M3 Artifact oracle golden cases — image/chart meaningful slides.

Pre-BLOCK calibration requested by the M2.2 review: lock the behaviour
of ``artifact.nonempty`` / ``pptx.min_nonempty_slides`` on picture-only
and chart-bearing decks. No verifier version bump — this only locks
already-shipped v3 semantics.
"""

from __future__ import annotations

import io
import unittest

from contracts.verification import AcceptancePredicate
from total_gateway.verification_registry import VerifierRegistry
from total_gateway.outcome_oracles.artifact_content import ArtifactContentOracle

from tests.test_p19_m2_1_artifact_oracle import (
    M21OracleTestBase,
    PPTX_MIME,
    pptx_bytes,
)


def chart_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("series", (1, 2, 3))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, 300, 200, chart_data
    )
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


class ArtifactImageChartGoldenTests(M21OracleTestBase):
    def _nonempty(self):
        return AcceptancePredicate.create(
            predicate_type="artifact.nonempty", subject_kind="artifact"
        )

    def _min_slides(self, minimum: int):
        return AcceptancePredicate.create(
            predicate_type="pptx.min_nonempty_slides",
            subject_kind="artifact",
            params={"min_slides": minimum},
        )

    def test_picture_only_deck_is_nonempty(self) -> None:
        manifest = self._passed_manifest(
            pptx_bytes([], with_picture_slide=True),
            filename="pictures.pptx",
            format_id="pptx",
            declared_mime=PPTX_MIME,
        )
        record = self.oracle.evaluate(
            manifest, self._nonempty(), evaluated_at_ms=21_000
        )
        self.assertEqual(record.status, "PASS", record.reason_codes)

    def test_picture_only_deck_counts_toward_min_slides(self) -> None:
        manifest = self._passed_manifest(
            pptx_bytes(["标题页"], with_picture_slide=True),
            filename="mixed.pptx",
            format_id="pptx",
            declared_mime=PPTX_MIME,
        )
        record = self.oracle.evaluate(
            manifest, self._min_slides(2), evaluated_at_ms=21_000
        )
        self.assertEqual(record.status, "PASS", record.reason_codes)
        three = self.oracle.evaluate(
            manifest, self._min_slides(3), evaluated_at_ms=21_000
        )
        self.assertEqual(three.status, "FAIL")

    def test_chart_only_deck_is_nonempty(self) -> None:
        manifest = self._passed_manifest(
            chart_pptx_bytes(),
            filename="chart.pptx",
            format_id="pptx",
            declared_mime=PPTX_MIME,
        )
        record = self.oracle.evaluate(
            manifest, self._nonempty(), evaluated_at_ms=21_000
        )
        self.assertEqual(record.status, "PASS", record.reason_codes)
        slides = self.oracle.evaluate(
            manifest, self._min_slides(1), evaluated_at_ms=21_000
        )
        self.assertEqual(slides.status, "PASS", slides.reason_codes)

    def test_artifact_descriptor_stays_v3(self) -> None:
        snapshot = VerifierRegistry.with_defaults().snapshot(captured_at_ms=1)
        descriptor = snapshot.find("verifier.artifact_content")
        assert descriptor is not None
        self.assertEqual(descriptor.verifier_version, "3")
        # golden cases must not widen the declared capability set
        self.assertNotIn("pptx.required_slide_titles", descriptor.supported_predicate_types)


if __name__ == "__main__":
    unittest.main()
