from __future__ import annotations

import ast
import importlib.util
import json
from pathlib import Path
import tempfile
import typing
import unittest


ROOT = Path(__file__).resolve().parents[1]
RUNTIME_SOURCE = ROOT / "app" / "backend" / "tiangong-backend" / "_internal" / "omni_body_skill" / "tools" / "omni_body_tool.py"
READABLE_RUNTIME_SOURCE = ROOT / "readable-python-source" / "omni_body_skill" / "tools" / "omni_body_tool.py"
PPT_DESIGN_SOURCE = ROOT / "app" / "backend" / "tiangong-backend" / "_internal" / "omni_body_skill" / "tools" / "ppt_design.py"
READABLE_PPT_DESIGN_SOURCE = ROOT / "readable-python-source" / "omni_body_skill" / "tools" / "ppt_design.py"
DELIVERY_SOURCE = ROOT / "app" / "backend" / "tiangong-backend" / "_internal" / "omni_body_skill" / "tools" / "delivery_kernel.py"
READABLE_DELIVERY_SOURCE = ROOT / "readable-python-source" / "omni_body_skill" / "tools" / "delivery_kernel.py"
DESIGN_TEMPLATE = ROOT / "app" / "backend" / "tiangong-backend" / "_internal" / "omni_body_skill" / "templates" / "executive_ppt.design.json"


def _load_file_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _load_registry_finalizer():
    tree = ast.parse(RUNTIME_SOURCE.read_text(encoding="utf-8"), filename=str(RUNTIME_SOURCE))
    required_node = next(
        node
        for node in tree.body
        if isinstance(node, ast.AnnAssign)
        and isinstance(node.target, ast.Name)
        and node.target.id == "_REQUIRED_SINGLE_TOOL_ACTIONS"
    )
    finalizer = next(
        node
        for node in tree.body
        if isinstance(node, ast.FunctionDef) and node.name == "_finalize_delivery_registry"
    )
    required = ast.literal_eval(required_node.value)
    namespace = {
        "Any": typing.Any,
        "Dict": typing.Dict,
        "Tuple": typing.Tuple,
        "_REQUIRED_SINGLE_TOOL_ACTIONS": required,
    }
    exec(compile(ast.Module(body=[finalizer], type_ignores=[]), str(RUNTIME_SOURCE), "exec"), namespace)
    return required, namespace["_finalize_delivery_registry"]


class _WorkspaceRuntime:
    def __init__(self, workspace: Path):
        self.workspace = workspace.resolve()

    def _resolve(self, value, must_exist=False):
        path = Path(value)
        if not path.is_absolute():
            path = self.workspace / path
        path = path.resolve()
        path.relative_to(self.workspace)
        if must_exist and not path.exists():
            raise FileNotFoundError(path)
        return path

    def _rel(self, path):
        return Path(path).resolve().relative_to(self.workspace).as_posix()


class PptDesignPipelineTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.ppt_design = _load_file_module("test_tiangong_ppt_design", PPT_DESIGN_SOURCE)
        cls.delivery = _load_file_module("test_tiangong_delivery_kernel", DELIVERY_SOURCE)
        # A direct file import has no package for the relative import used by
        # delivery_kernel.  Inject the same production inspector explicitly.
        cls.delivery._ppt_inspection = cls.ppt_design.inspect_presentation

    def test_source_and_readable_mirrors_are_identical(self) -> None:
        self.assertEqual(RUNTIME_SOURCE.read_bytes(), READABLE_RUNTIME_SOURCE.read_bytes())
        self.assertEqual(PPT_DESIGN_SOURCE.read_bytes(), READABLE_PPT_DESIGN_SOURCE.read_bytes())
        self.assertEqual(DELIVERY_SOURCE.read_bytes(), READABLE_DELIVERY_SOURCE.read_bytes())

    def test_required_actions_degrade_to_explicit_unavailable_not_unknown(self) -> None:
        required, finalizer = _load_registry_finalizer()
        registry, errors = finalizer({}, "ImportError: simulated delivery failure")
        self.assertEqual(set(registry), set(required))
        self.assertEqual(set(errors), set(required))
        for name in required:
            self.assertFalse(registry[name]["implemented"])
            self.assertEqual(registry[name]["capability_state"], "unavailable")
            self.assertIn("simulated delivery failure", registry[name]["unavailable_reason"])

    def test_template_apply_emits_compact_machine_readable_design_sidecar(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            runtime = _WorkspaceRuntime(Path(temporary))
            result = self.delivery._template_apply(
                runtime,
                "story.md",
                {"template_id": "executive_ppt", "variables": {"title": "治理汇报", "audience": "决策委员会"}},
            )
            story = runtime._resolve("story.md", must_exist=True)
            sidecar = runtime._resolve(result["design_spec"]["path"], must_exist=True)
            design = json.loads(sidecar.read_text(encoding="utf-8"))
            self.assertTrue(story.is_file())
            self.assertEqual(sidecar.name, "story.design.json")
            self.assertEqual(design["schema"], "tiangong.v3.ppt_design.v1")
            self.assertEqual(result["next_action_args"]["design_spec"], "story.design.json")

    def test_design_generator_is_widescreen_placeholder_free_and_qc_ready(self) -> None:
        content = """# 潮汐算法让排水决策可解释、可追溯
面向城市治理决策委员会的执行汇报

---

## 核心结论：统一事件账本可消除跨部门口径冲突
- 每次调度都记录输入、规则版本、责任主体和输出证据
- 审计人员可以从结果反向追溯到原始依据

---

## 建议：风险阈值必须先于自动调度进入生产
- 红线风险由人工确认，系统不得绕过
- 中低风险通过策略版本和证据链自动执行

---

## 实施路径将按三个阶段逐步降低切换风险
- 第一阶段建立只读影子账本并核对差异
- 第二阶段开放低风险建议但保留人工执行
- 第三阶段只对通过验收的场景授权自动动作

---

## 下一步行动：批准四周试点并指定唯一责任人
- 本周确认试点区域、验收指标与回滚条件
- 每周复盘误差、异常和人工接管记录
"""
        with tempfile.TemporaryDirectory() as temporary:
            workspace = Path(temporary).resolve()
            source = workspace / "story.md"
            source.write_text(content, encoding="utf-8")
            source.with_suffix(".design.json").write_bytes(DESIGN_TEMPLATE.read_bytes())
            resolve = _WorkspaceRuntime(workspace)._resolve
            prs, metadata = self.ppt_design.build_presentation(
                resolve,
                {"source": str(source), "template_id": "executive_ppt"},
                content=content,
                source_path=source,
            )
            output = workspace / "designed.pptx"
            prs.save(output)
            inspection = self.ppt_design.inspect_presentation(output)
            self.assertTrue(inspection["is_widescreen"])
            self.assertEqual(inspection["placeholder_count"], 0)
            self.assertGreaterEqual(inspection["designed_visual_count"], inspection["slide_count"])
            self.assertIn("Microsoft YaHei", inspection["font_names"])
            self.assertTrue(metadata["placeholder_free"])

            qc = self.delivery._qc_ppt(_WorkspaceRuntime(workspace), output, {})["result"]
            self.assertTrue(qc["hard_gate_passed"], qc["issues"])
            self.assertTrue(qc["acceptance"], qc)
            self.assertEqual(qc["visual_coverage"], 1.0)

    def test_legacy_default_placeholder_deck_is_a_hard_failure(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as temporary:
            workspace = Path(temporary).resolve()
            output = workspace / "legacy.pptx"
            prs = Presentation()
            for index in range(5):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"第{index + 1}页资料"
                slide.placeholders[1].text = "只有默认占位符和文字，没有图表、图片、表格或设计组件。"
            prs.save(output)

            qc = self.delivery._qc_ppt(_WorkspaceRuntime(workspace), output, {})["result"]
            codes = {item["code"] for item in qc["issues"]}
            self.assertFalse(qc["acceptance"])
            self.assertFalse(qc["hard_gate_passed"])
            self.assertIn("no_meaningful_visuals", codes)
            self.assertIn("default_placeholder_layout", codes)
            self.assertIn("legacy_aspect_ratio", codes)
            self.assertLessEqual(qc["score"], 59)

    def test_structured_chart_and_table_are_native_evidence_not_fake_numbers(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            workspace = Path(temporary).resolve()
            args = {
                "title": "治理试点证据支持分阶段决策",
                "subtitle": "所有数值均由调用方显式提供",
                "slides": [
                    {
                        "title": "结论：人工接管率连续三周下降",
                        "chart": {
                            "categories": ["第1周", "第2周", "第3周"],
                            "series": [{"name": "人工接管率", "values": [18, 12, 7]}],
                        },
                    },
                    {
                        "title": "下一步行动：仅开放通过验收的场景",
                        "table": {
                            "headers": ["场景", "状态", "责任人"],
                            "rows": [["低风险预警", "通过", "值班主管"], ["红线调度", "人工确认", "总指挥"]],
                        },
                    },
                ],
            }
            prs, _metadata = self.ppt_design.build_presentation(_WorkspaceRuntime(workspace)._resolve, args)
            output = workspace / "evidence.pptx"
            prs.save(output)
            inspection = self.ppt_design.inspect_presentation(output)
            self.assertEqual(inspection["slide_count"], 3)
            self.assertGreaterEqual(inspection["native_visual_count"], 2)
            qc = self.delivery._qc_ppt(_WorkspaceRuntime(workspace), output, {"min_slides": 3})["result"]
            self.assertTrue(qc["acceptance"], qc)
            self.assertEqual(qc["native_visual_count"], 2)


if __name__ == "__main__":
    unittest.main()
